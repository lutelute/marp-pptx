"""Deck audit: does it find the defects, and does it stay quiet otherwise?

Fixtures are built with python-pptx directly so each test controls the exact
geometry that should (or should not) trip a check. The last group runs the
audit over decks this project builds — the gate that keeps our own output
clean.
"""
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from marp_pptx.audit import audit_pptx, contrast_ratio, format_findings
from marp_pptx.pkgcheck import check_package


def _deck(tmp_path, build) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build(slide)
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    return out


def _text(slide, left, top, w, h, text, size=18, color=RGBColor(0, 0, 0),
          grow=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    # python-pptx writes <a:spAutoFit/> into every new textbox; a deck built by
    # this library pins the height instead, and that is the case under test.
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT if grow else MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = "Arial"
    run.font.color.rgb = color
    return tb


def _kinds(findings, kind, severity=None):
    return [f for f in findings
            if f.kind == kind and (severity is None or f.severity == severity)]


# ── overflow ────────────────────────────────────────────────────────────────
def test_overflow_is_reported_with_the_overshoot(tmp_path):
    long_line = "This sentence is far too long for the small box it was put in, "
    deck = _deck(tmp_path, lambda s: _text(s, 1, 1, 2.0, 0.4, long_line * 3))
    found = _kinds(audit_pptx(deck), "overflow", "error")
    assert found, "a 6-line paragraph in a 0.4in box must be reported"
    assert found[0].detail["needed_pt"] > found[0].detail["box_pt"]
    assert "pt" in found[0].fix


def test_a_box_that_grows_is_a_warning_not_a_clip(tmp_path):
    """spAutoFit doesn't clip — it shoves whatever is underneath out of place,
    which is a different (and quieter) problem than cut-off text."""
    long_line = "This sentence is far too long for the small box it was put in, "
    deck = _deck(tmp_path, lambda s: _text(s, 1, 1, 2.0, 0.4, long_line * 3,
                                           grow=True))
    found = _kinds(audit_pptx(deck), "overflow")
    assert found and found[0].severity == "warn"


def test_text_that_fits_is_not_reported(tmp_path):
    deck = _deck(tmp_path, lambda s: _text(s, 1, 1, 6.0, 1.0, "Short title"))
    assert not _kinds(audit_pptx(deck), "overflow")


def test_japanese_is_measured_with_the_east_asian_typeface(tmp_path):
    """Measured with a Latin face, 40 full-width characters look ~40% narrower
    and the overflow disappears — the check has to read <a:ea>."""
    jp = "日本語の本文がここに入ります。" * 4

    def build(slide):
        tb = _text(slide, 1, 1, 3.0, 0.5, jp)
        rpr = tb.text_frame.paragraphs[0].runs[0]._r.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr")
        ea = rpr.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}ea", {})
        ea.set("typeface", "Hiragino Sans")
        rpr.append(ea)

    assert _kinds(audit_pptx(_deck(tmp_path, build)), "overflow", "error")


# ── geometry ────────────────────────────────────────────────────────────────
def test_shape_past_the_slide_edge_is_an_error(tmp_path):
    deck = _deck(tmp_path, lambda s: _text(s, 12.0, 1, 4.0, 1.0, "off the edge"))
    assert _kinds(audit_pptx(deck), "offslide", "error")


def test_text_crowding_the_edge_is_a_warning(tmp_path):
    deck = _deck(tmp_path, lambda s: _text(s, 0.05, 3, 4.0, 1.0, "hugging the edge"))
    assert _kinds(audit_pptx(deck), "offslide", "warn")


def test_overlapping_text_blocks_are_reported(tmp_path):
    def build(slide):
        _text(slide, 2, 2, 4.0, 1.0, "first block of text")
        _text(slide, 2.3, 2.2, 4.0, 1.0, "second block on top")

    assert _kinds(audit_pptx(_deck(tmp_path, build)), "overlap", "error")


def test_side_by_side_columns_are_not_overlaps(tmp_path):
    def build(slide):
        _text(slide, 1, 2, 5.0, 1.0, "left column")
        _text(slide, 7, 2, 5.0, 1.0, "right column")

    assert not _kinds(audit_pptx(_deck(tmp_path, build)), "overlap")


# ── tables ──────────────────────────────────────────────────────────────────
def _table(slide, rows, cols, left, top, w, h, fill):
    gf = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                Inches(w), Inches(h))
    for ri in range(rows):
        for ci in range(cols):
            cell = gf.table.cell(ri, ci)
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = fill(ri, ci)
            run.font.size = Pt(14)
            run.font.name = "Arial"
    return gf


def test_a_row_too_short_for_its_cell_is_reported(tmp_path):
    """A table grows rather than clips, so a cramped row shoves everything
    below it down — invisible to every text-frame check, because a table is a
    graphic frame."""
    long = "この列には折り返しが必要な長い説明文が入っていて一行では収まりません"
    deck = _deck(tmp_path, lambda s: _table(
        s, 2, 2, 1, 1, 4.0, 0.6,
        lambda r, c: long if (r == 1 and c == 1) else "x"))
    found = _kinds(audit_pptx(deck), "overflow")
    assert found, "a wrapped cell in a 0.3in row must be reported"
    assert "表" in found[0].shape


def test_a_table_running_off_the_slide_is_an_error(tmp_path):
    long = "折り返しが必要な長い説明文" * 3
    deck = _deck(tmp_path, lambda s: _table(
        s, 6, 2, 1, 5.0, 4.0, 2.0, lambda r, c: long if c else "x"))
    assert _kinds(audit_pptx(deck), "overflow", "error")


def test_a_table_that_fits_is_quiet(tmp_path):
    deck = _deck(tmp_path, lambda s: _table(
        s, 3, 3, 1, 1, 9.0, 1.8, lambda r, c: f"r{r}c{c}"))
    assert not _kinds(audit_pptx(deck), "overflow")


# ── colour ──────────────────────────────────────────────────────────────────
def test_contrast_ratio_matches_the_wcag_reference():
    assert contrast_ratio((0, 0, 0), (255, 255, 255)) == pytest.approx(21.0, abs=0.1)
    assert contrast_ratio((255, 255, 255), (255, 255, 255)) == pytest.approx(1.0)


def test_low_contrast_body_text_is_flagged(tmp_path):
    deck = _deck(tmp_path, lambda s: _text(
        s, 1, 3, 6.0, 1.0, "pale grey on white", size=12,
        color=RGBColor(0xCC, 0xCC, 0xCC)))
    assert _kinds(audit_pptx(deck), "contrast", "warn")


def test_black_on_white_passes(tmp_path):
    deck = _deck(tmp_path, lambda s: _text(s, 1, 3, 6.0, 1.0, "readable", size=12))
    assert not _kinds(audit_pptx(deck), "contrast")


# ── package integrity ───────────────────────────────────────────────────────
def test_a_freshly_written_deck_is_structurally_clean(tmp_path):
    deck = _deck(tmp_path, lambda s: _text(s, 1, 1, 6.0, 1.0, "hello"))
    assert [f for f in check_package(deck) if f.severity == "error"] == []


def test_a_dangling_relationship_is_caught(tmp_path):
    deck = _deck(tmp_path, lambda s: _text(s, 1, 1, 6.0, 1.0, "hello"))
    broken = tmp_path / "broken.pptx"
    with zipfile.ZipFile(deck) as src, zipfile.ZipFile(broken, "w") as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/_rels/presentation.xml.rels":
                data = data.replace(b'Target="slides/slide1.xml"',
                                    b'Target="slides/slide99.xml"')
            dst.writestr(item, data)
    errors = [f for f in check_package(broken) if f.severity == "error"]
    assert errors and any("slide99" in f.message for f in errors)


def test_not_a_pptx_at_all(tmp_path):
    junk = tmp_path / "junk.pptx"
    junk.write_bytes(b"this is not a zip")
    assert check_package(junk)[0].severity == "error"


# ── end-to-end over our own output ──────────────────────────────────────────
def _build(src: Path, out_dir: Path, palette: str = "claude") -> Path:
    import contextlib
    import io

    from marp_pptx.builder import PptxBuilder
    from marp_pptx.parser import parse_marp
    from marp_pptx.theme import ThemeConfig, get_default_theme_path, get_palette_path

    tc = ThemeConfig.from_css(get_default_theme_path())
    with contextlib.redirect_stderr(io.StringIO()):
        pal = get_palette_path(palette)
        if pal:
            tc.apply_palette(pal)
        builder = PptxBuilder(base_path=src.parent, theme=tc)
        builder.build_all(parse_marp(str(src)))
        out = out_dir / f"{src.stem}.pptx"
        builder.save(str(out))
    return out


def _template_ids(project_root):
    return sorted(p.stem for p in (project_root / "templates").glob("*.md"))


def test_every_bundled_template_builds_without_visible_defects(tmp_path, project_root):
    """The gate this whole module exists for: nothing we ship clips, overlaps,
    leaves the slide, or arrives structurally broken. Every type, every time."""
    bad = []
    for name in _template_ids(project_root):
        out = _build(project_root / "templates" / f"{name}.md", tmp_path)
        errors = [f for f in check_package(out) + audit_pptx(out)
                  if f.severity == "error"]
        if errors:
            bad.append(f"{name}: " + format_findings(errors))
    assert not bad, "\n".join(bad)


@pytest.mark.parametrize("preset", ["academic-talk", "academic-dense", "lecture",
                                    "product", "minimal"])
def test_presets_build_without_visible_defects(tmp_path, project_root, preset):
    src = project_root / "src" / "marp_pptx" / "data" / "presets" / f"{preset}.md"
    if not src.exists():
        pytest.skip(f"{preset} preset not bundled")
    out = _build(src, tmp_path)
    errors = [f for f in check_package(out) + audit_pptx(out)
              if f.severity == "error"]
    assert not errors, format_findings(errors)


@pytest.mark.parametrize("scale,density", [(1.22, "keynote"), (0.8, "academic"),
                                           (1.3, "academic")])
def test_templates_survive_the_shipped_scale_options(tmp_path, project_root,
                                                     scale, density):
    """`--density keynote` and `--font-scale` are shipped options, so the
    reserved boxes have to follow the type. They didn't: fixed-inch heights and
    a paragraph gap that scaled while the builder wrote a literal Pt(4) put 40+
    errors into these settings while the default read clean."""
    import contextlib
    import io

    from marp_pptx.builder import PptxBuilder
    from marp_pptx.parser import parse_marp
    from marp_pptx.theme import ThemeConfig, get_default_theme_path, get_palette_path

    bad = []
    for name in _template_ids(project_root):
        src = project_root / "templates" / f"{name}.md"
        tc = ThemeConfig.from_css(get_default_theme_path())
        with contextlib.redirect_stderr(io.StringIO()):
            pal = get_palette_path("claude")
            if pal:
                tc.apply_palette(pal)
            tc.font_scale = scale
            tc.margin_scale = 1.12 if density == "keynote" else 1.0
            tc.density = density
            builder = PptxBuilder(base_path=src.parent, theme=tc)
            builder.build_all(parse_marp(str(src)))
            out = tmp_path / f"{name}.pptx"
            builder.save(str(out))
        errors = [f for f in audit_pptx(out) if f.severity == "error"]
        if errors:
            bad.append(f"{name}: " + format_findings(errors))
    assert not bad, "\n".join(bad)


@pytest.mark.parametrize("palette", ["claude", "minimal", "beamer", "tmu-cs"])
def test_a_dense_deck_survives_every_theme(tmp_path, project_root, palette):
    """Themes move the geometry — a band title, a footer bar, wider margins —
    so a deck that fits in one can clip in another."""
    src = project_root / "src" / "marp_pptx" / "data" / "presets" / "academic-dense.md"
    if not src.exists():
        pytest.skip("academic-dense preset not bundled")
    out = _build(src, tmp_path, palette)
    errors = [f for f in check_package(out) + audit_pptx(out)
              if f.severity == "error"]
    assert not errors, format_findings(errors)


def _every_text(slide):
    """Every string on a slide, tables included.

    A table is a graphic frame with no text frame of its own, so walking
    `shape.text_frame` alone reads past it — which is how `$\\sqrt{n}$` sat in
    an appendix table while the markup check reported the deck clean."""
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame.text
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text


def _decks(project_root):
    yield from sorted((project_root / "templates").glob("*.md"))
    yield from sorted((project_root / "src" / "marp_pptx" / "data" / "presets")
                      .glob("*.md"))
    yield from sorted(p for p in (project_root / "demo").glob("*.md")
                      if p.stem != "EVALUATION")


def test_no_markdown_or_latex_reaches_the_slide(tmp_path, project_root):
    """`**bold**` and `$x^2$` are instructions, not content. Any path that
    writes a paragraph's text directly instead of going through the inline
    renderer prints them verbatim — which is how theorem titles, footnotes,
    numbered steps, definition notes and table cells used to ship."""
    import re

    from pptx import Presentation

    leak = re.compile(r"\$[^$\n]{1,60}\$|\*\*\S")
    bad = []
    for src in _decks(project_root):
        out = _build(src, tmp_path)
        for i, slide in enumerate(Presentation(str(out)).slides, 1):
            for text in _every_text(slide):
                if leak.search(text):
                    bad.append(f"{src.stem} slide {i}: {text.strip()[:60]}")
    assert not bad, "\n".join(bad)


def test_format_findings_says_so_when_clean():
    assert "clean" in format_findings([])
