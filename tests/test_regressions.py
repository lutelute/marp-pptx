"""Regression tests for the 2026-06-08 quality-verification findings.

Each test pins one of the four real-content bugs found while authoring
demo/hosting-capacity.md:

1. title: inline $math$ / HTML entities leaked verbatim into author lines
2. timeline-h: <br> inside tl-h-detail was collapsed to a single line
3. takeaway: a wrapped ta-main overlapped ta-points (height under-estimate)
4. table-slide: the rendered table grew into the box below it (no headroom)
"""
import tempfile
from pathlib import Path

from pptx.util import Inches, Pt

from marp_pptx.parser import parse_slide, strip_html, text_with_breaks
from marp_pptx.theme import ThemeConfig
from marp_pptx.builder import PptxBuilder


def _make_builder(tmp_path=None):
    if tmp_path is None:
        tmp_path = Path(tempfile.mkdtemp())
    return PptxBuilder(base_path=tmp_path, theme=ThemeConfig())


def _all_text(slide):
    return "\n".join(s.text_frame.text for s in slide.shapes
                     if getattr(s, "has_text_frame", False))


# --- 1. entities + inline math on the title slide ---------------------------

def test_strip_html_unescapes_entities():
    assert strip_html("A &ensp;|&ensp; B") == "A \u2002|\u2002 B"
    assert strip_html("R&amp;D") == "R&D"


def test_title_author_line_renders_inline_markup():
    md = ("<!-- _class: title -->\n"
          "# Test Title\n"
          "## Subtitle\n\n"
          "Alice $^{1}$, Bob $^{2}$\n\n"
          "$^{1}$ Univ A &ensp;|&ensp; 2026")
    b = _make_builder()
    b.build_title(parse_slide(0, md))
    joined = _all_text(b.prs.slides[0])
    assert "Alice" in joined
    assert "$" not in joined          # raw $...$ must not leak as plain text
    assert "&ensp;" not in joined     # entities must be unescaped


# --- 2. <br> inside timeline details ----------------------------------------

def test_text_with_breaks_preserves_br():
    assert text_with_breaks("foo<br>bar") == "foo\nbar"
    assert text_with_breaks("foo <br/> bar") == "foo\nbar"
    assert text_with_breaks("a  <BR >  b") == "a\nb"


TL_H = """<!-- _class: timeline-h -->
# History
<div class="tl-h-container">
<div class="tl-h-item">
  <div class="tl-h-block">
    <span class="tl-h-year">2012</span>
    <span class="tl-h-text">CNN</span>
    <div class="tl-h-detail">line1<br>line2</div>
  </div>
</div>
</div>"""


def test_timeline_h_detail_keeps_br_as_newline():
    sd = parse_slide(0, TL_H)
    assert sd.timeline_items[0]["detail"] == "line1\nline2"


# --- 3. takeaway main wrapping must push the points down --------------------

def test_estimate_height_accounts_for_wrapping():
    b = _make_builder()
    long_jp = "感" * 60     # 60 em of full-width text
    h_nowrap = b._estimate_text_height([long_jp], Pt(28))
    h_wrap = b._estimate_text_height([long_jp], Pt(28), width=Inches(10))
    # 60 em at 28pt in a 720pt-wide box needs >= 3 lines
    assert h_wrap.pt >= h_nowrap.pt * 2


def test_takeaway_points_follow_wrapped_main():
    long_main = "とても長いキーメッセージで折り返しが必ず発生する" * 4   # ~96 em
    md = ("<!-- _class: takeaway -->\n# Key\n"
          f'<div class="ta-main">{long_main}</div>\n'
          '<div class="ta-points">\n<li>p1</li>\n<li>p2</li>\n</div>')
    b = _make_builder()
    b.build_takeaway(parse_slide(0, md))
    slide = b.prs.slides[0]
    boxes = [s for s in slide.shapes
             if getattr(s, "has_text_frame", False) and s.text_frame.text]
    main = next(s for s in boxes if "とても長い" in s.text_frame.text)
    pts = next(s for s in boxes if "p1" in s.text_frame.text)
    # the main box must be tall enough for >= 3 wrapped lines ...
    assert main.height >= int(Pt(28 * 1.25 * 3))
    # ... and the points box must start below it
    assert pts.top >= main.top + main.height


# --- 4. table block keeps headroom above whatever sits below it -------------

def test_table_block_reserves_headroom():
    rows = "| A | B |\n|---|---|\n" + "| a | b |\n" * 5   # header + 5 body rows
    md = ("<!-- _class: table-slide -->\n# Compare\n\n" + rows +
          '\n<div class="box-accent">\n\nNote below the table\n\n</div>')
    b = _make_builder()
    b.build_table(parse_slide(0, md))
    slide = b.prs.slides[0]
    tbl = next(s for s in slide.shapes if getattr(s, "has_table", False))
    # The frame is exactly its rows: a renderer has no reason to grow one, so
    # the last rule stays inside the frame and the block below keeps its gap.
    assert tbl.height == sum(r.height for r in tbl.table.rows)
    note = next(s for s in slide.shapes
                if getattr(s, "has_text_frame", False) and "Note below" in s.text_frame.text)
    assert note.top >= tbl.top + tbl.height


# --- 5. appendix with a table must stay ONE slide ----------------------------

def test_appendix_table_stays_on_one_slide():
    md = ("<!-- _class: appendix -->\n# Params\n\n"
          '<span class="appendix-label">Appendix A</span>\n\n'
          "| k | v |\n|---|---|\n| lr | 3e-4 |\n| bs | 64 |\n")
    b = _make_builder()
    b.build_appendix(parse_slide(0, md))
    assert len(b.prs.slides) == 1          # was split into 2 slides before
    slide = b.prs.slides[0]
    assert any(getattr(s, "has_table", False) for s in slide.shapes)


# --- 6. unclosed <div> must terminate (used to spin forever) ----------------

def test_unclosed_div_terminates_and_salvages():
    from marp_pptx.parser import extract_child_divs
    out = extract_child_divs('<div class="a">x<div class="b">content')
    assert out and "content" in out[0]      # salvaged, and returned at all
    md = ('<!-- _class: kpi -->\n# K\n<div class="kpi-container">\n'
          '<div class="kpi-item">\n<span class="kpi-value">42</span>\n'
          '<span class="kpi-label">answer</span>\n')
    sd = parse_slide(0, md)                  # must not hang
    assert sd.slide_class == "kpi"


# --- 7. non-fatal authoring warnings are surfaced, not silent ---------------

def test_unknown_slide_type_warns():
    b = _make_builder()
    b.build_all([parse_slide(0, "<!-- _class: no-such-type -->\n# H\n\nbody")])
    assert any("unknown type" in w and "no-such-type" in w for w in b.warnings)
    assert len(b.prs.slides) == 1            # still renders (as plain slide)


def test_missing_image_warns():
    b = _make_builder()
    assert b._resolve_image("does-not-exist.png") is None
    assert any("image not found" in w for w in b.warnings)


def test_warnings_are_deduplicated():
    b = _make_builder()
    b._resolve_image("nope.png")
    b._resolve_image("nope.png")
    assert sum("nope.png" in w for w in b.warnings) == 1


def test_known_type_does_not_warn():
    b = _make_builder()
    b.build_all([parse_slide(0, "<!-- _class: title -->\n# T\n## sub")])
    assert b.warnings == []


# --- 8. visual_lint must ignore the page-number footer chrome ---------------

def _synth_png(rects, size=(1280, 720), bg=(250, 249, 245)):
    """White slide with black rects [(x0,y0,x1,y1), …]. Returns a temp path."""
    from PIL import Image
    im = Image.new("RGB", size, bg)
    px = im.load()
    for x0, y0, x1, y1 in rects:
        for y in range(y0, y1):
            for x in range(x0, x1):
                px[x, y] = (20, 20, 20)
    p = tempfile.mktemp(suffix=".png")
    im.save(p)
    return p


def test_visual_lint_ignores_page_number_corner():
    from marp_pptx.visuallint import visual_lint
    # A page-number-like blob in the bottom-right corner + normal centered
    # content must NOT be flagged as edge overflow (regression: it was).
    png = _synth_png([(400, 250, 880, 470),      # centered content, well inside
                      (1180, 690, 1260, 712)])    # "N / M" footer chrome
    warns = visual_lint(png, "kpi")
    assert not any("reaches the slide edge" in w for w in warns)


def test_visual_lint_still_flags_real_bottom_overflow():
    from marp_pptx.visuallint import visual_lint
    # Content spanning most of the width and running to the bottom edge is a
    # genuine overflow and must still be flagged.
    png = _synth_png([(100, 400, 900, 718)])
    warns = visual_lint(png, "kpi")
    assert any("reaches the slide edge" in w for w in warns)


# --- 9. timeline: two-class spans + highlight on the item div ---------------

_TL_H = """<!-- _class: timeline-h -->
# H
<div class="tl-h-container">
<div class="tl-h-item">
  <div class="tl-h-block">
    <span class="tl-h-year">2012</span>
    <span class="tl-h-text">AlexNet</span>
  </div>
</div>
<div class="tl-h-item highlight">
  <div class="tl-h-block">
    <span class="tl-h-year">2024</span>
    <span class="tl-h-text bold">本研究</span>
  </div>
</div>
</div>"""

_TL_V = """<!-- _class: timeline -->
# V
<div class="tl-container">
<div class="tl-item">
  <span class="tl-year">2012</span>
  <span class="tl-text">AlexNet</span>
</div>
<div class="tl-item highlight">
  <span class="tl-year">2024</span>
  <span class="tl-text bold">本研究</span>
</div>
</div>"""


def test_timeline_h_two_class_text_and_highlight():
    sd = parse_slide(0, _TL_H)
    hi = [it for it in sd.timeline_items if it["highlight"]]
    assert len(hi) == 1
    assert hi[0]["text"] == "本研究"      # was dropped (two-class regex)
    assert hi[0]["year"] == "2024"


def test_timeline_v_two_class_text_and_highlight():
    sd = parse_slide(0, _TL_V)
    hi = [it for it in sd.timeline_items if it["highlight"]]
    assert len(hi) == 1
    assert hi[0]["text"] == "本研究"
    # the non-highlighted item must NOT be marked highlighted
    assert sum(it["highlight"] for it in sd.timeline_items) == 1


# --- 10. zone-matrix must render BOTH axis labels ---------------------------

def test_zone_matrix_renders_axis_labels():
    md = ("<!-- _class: zone-matrix -->\n# M\n"
          '<div class="zm-container">\n'
          '<div class="zm-ylabel">YAXIS_PRECISION</div>\n'
          '<div class="zm-xlabel">XAXIS_COST</div>\n'
          '<div class="zm-cell zm-tl"><span class="zm-label">A</span>'
          '<span class="zm-body">a</span></div>\n'
          '</div>')
    b = _make_builder()
    b.build_zone_matrix(parse_slide(0, md))
    txt = "".join(s.text_frame.text for sl in b.prs.slides
                  for s in sl.shapes if getattr(s, "has_text_frame", False))
    assert "XAXIS_COST" in txt      # x-axis label was dropped entirely
    assert "YAXIS_PRECISION" in txt


# --- 11. overview/result must render the figure caption ---------------------

def test_overview_renders_figure_caption():
    from PIL import Image
    d = Path(tempfile.mkdtemp())
    Image.new("RGB", (400, 260), (230, 235, 242)).save(d / "f.png")
    md = ("<!-- _class: overview -->\n# O\n\n"
          '<div class="ov-lead">lead</div>\n\n'
          "![w:600](f.png)\n\n"
          '<div class="caption">Fig. 1. CAPTION_TOKEN here</div>\n\n'
          '<div class="ov-points">\n<li>p1</li>\n</div>')
    b = PptxBuilder(base_path=d, theme=ThemeConfig())
    b.build_all([parse_slide(0, md)])
    txt = "".join(s.text_frame.text for sl in b.prs.slides
                  for s in sl.shapes if getattr(s, "has_text_frame", False))
    assert "CAPTION_TOKEN" in txt      # caption was dropped by _build_image_points


# --- 12. gallery-img must render each image's caption -----------------------

def test_gallery_img_renders_captions():
    from PIL import Image
    d = Path(tempfile.mkdtemp())
    Image.new("RGB", (300, 200), (235, 235, 240)).save(d / "g.png")
    md = ("<!-- _class: gallery-img -->\n# G\n"
          '<div class="gi-container">\n'
          '<div class="gi-item">\n\n![w:300](g.png)\n\n'
          '<div class="gi-caption">CAP_ALPHA</div>\n</div>\n'
          '<div class="gi-item">\n\n![w:300](g.png)\n\n'
          '<div class="gi-caption">CAP_BETA</div>\n</div>\n'
          '</div>')
    b = PptxBuilder(base_path=d, theme=ThemeConfig())
    b.build_all([parse_slide(0, md)])
    txt = "".join(s.text_frame.text for sl in b.prs.slides
                  for s in sl.shapes if getattr(s, "has_text_frame", False))
    assert "CAP_ALPHA" in txt and "CAP_BETA" in txt   # captions were dropped


# --- 13. rq: a long question must not overflow the card onto rq_sub ----------

def test_rq_long_main_does_not_overlap_sub():
    long_q = "長系列データにおいて推論精度を維持したまま計算量を大幅に削減し" * 3 + "できるか"
    md = ("<!-- _class: rq -->\n# RQ\n"
          f'<div class="rq-main">\n{long_q}\n</div>\n'
          '<div class="rq-sub">\n— サブ説明テキスト\n</div>')
    b = _make_builder()
    b.build_rq(parse_slide(0, md))
    shapes = [s for s in b.prs.slides[0].shapes
              if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()]
    main = next(s for s in shapes if "長系列" in s.text_frame.text)
    sub = next(s for s in shapes if "サブ説明" in s.text_frame.text)
    # the long question wraps to several lines, so the card must be sized for
    # at least 2 lines (SZ_H2 = 24pt). The width-unaware estimate sized it for
    # ~1 line + padding and the text spilled out the bottom onto rq_sub.
    assert main.height >= int(Pt(24 * 1.25 * 2)) + int(Inches(0.4))
    assert sub.top >= main.top + main.height - int(Inches(0.05))


# --- 14. quote / definition: wrapped body must not collide with what's below

def test_quote_long_text_no_overlap_with_source():
    long_q = "科学とは単なる知識の集積ではなくむしろ世界を理解するための体系的な思考の方法であり" * 2
    md = (f'<!-- _class: quote -->\n# Q\n<div class="qt-text">{long_q}</div>\n'
          '<div class="qt-source">Carl Sagan (1995)</div>')
    b = _make_builder()
    b.build_quote(parse_slide(0, md))
    sh = [s for s in b.prs.slides[0].shapes
          if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()]
    quote = next(s for s in sh if "科学とは" in s.text_frame.text)
    src = next(s for s in sh if "Carl Sagan" in s.text_frame.text)
    assert src.top >= quote.top + quote.height - int(Inches(0.05))


def test_definition_long_body_no_overlap_with_note():
    long_b = "入力系列の全要素ペア間ではなく選択的に注意重みを計算することで計算量を大幅に削減する手法" * 2
    md = ('<!-- _class: definition -->\n# D\n'
          '<div class="df-term">用語</div>\n'
          f'<div class="df-body">{long_b}</div>\n'
          '<div class="df-note">関連: NOTE_TOKEN</div>')
    b = _make_builder()
    b.build_definition(parse_slide(0, md))
    sh = [s for s in b.prs.slides[0].shapes
          if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()]
    body = next(s for s in sh if "入力系列" in s.text_frame.text)
    note = next(s for s in sh if "NOTE_TOKEN" in s.text_frame.text)
    assert note.top >= body.top + body.height - int(Inches(0.05))


# --- 15. highlight band must grow for a long (wrapping) message --------------

def test_highlight_band_fits_long_text():
    long_h = "動的スパースマスクの導入により計算量を大幅に削減しながら推論精度を維持できる" * 2
    md = f'<!-- _class: highlight -->\n# H\n<div class="hl-text">{long_h}</div>'
    b = _make_builder()
    b.build_highlight(parse_slide(0, md))
    box = next(s for s in b.prs.slides[0].shapes
               if getattr(s, "has_text_frame", False) and "動的スパース" in s.text_frame.text)
    # the band/textbox must be tall enough for >= 3 wrapped lines at 30pt
    # (was sized for ~1 line, so text spilled out and the rule hit the text)
    assert box.height >= int(Pt(30 * 1.25 * 3))


# --- 16. geometric text-collision detector ----------------------------------

def _mk_textbox(slide, l, t, w, h, text, pt):
    from pptx.util import Inches as _In, Pt as _Pt
    tb = slide.shapes.add_textbox(_In(l), _In(t), _In(w), _In(h))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.size = _Pt(pt)
    return tb


def test_collision_detector_flags_real_overflow():
    from pptx import Presentation
    from marp_pptx.visuallint import detect_text_collisions
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _mk_textbox(s, 1, 1, 8, 0.5, "とても長い文章で必ず複数行に折り返してしまう" * 3, 24)  # short box
    _mk_textbox(s, 1, 1.7, 8, 0.5, "下に置かれた別の要素", 18)
    hits = detect_text_collisions(prs)
    assert hits and hits[0]["slide"] == 1


def test_collision_detector_clean_when_box_fits():
    from pptx import Presentation
    from marp_pptx.visuallint import detect_text_collisions
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _mk_textbox(s, 1, 1, 8, 2.0, "短い文", 24)        # box has room
    _mk_textbox(s, 1, 3.2, 8, 0.5, "下の要素", 18)
    assert detect_text_collisions(prs) == []


def test_collision_detector_ignores_page_number_chrome():
    from pptx import Presentation
    from marp_pptx.visuallint import detect_text_collisions
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # a full-width footnote whose box spans the page number must NOT collide
    _mk_textbox(s, 0.5, 6.8, 9, 0.4, "脚注テキストがここに入る", 12)
    _mk_textbox(s, 8.5, 6.9, 1, 0.3, "7 / 14", 11)   # page-number chrome
    assert detect_text_collisions(prs) == []


def test_builder_self_check_clean_on_all_skeletons():
    import glob
    from marp_pptx.parser import parse_marp
    base = Path(__file__).resolve().parents[1] / "src/marp_pptx/data/templates"
    for f in sorted(glob.glob(str(base / "*.md"))):
        b = PptxBuilder(base_path=base, theme=ThemeConfig())
        b.build_all(parse_marp(f))
        bad = [w for w in b.warnings if "overflows onto" in w]
        assert not bad, f"{Path(f).name}: {bad}"


# --- 17. profile bio is width-aware; detector catches off-slide overflow -----

def test_profile_long_bio_is_width_aware():
    # A bio of long bullets wraps in the right column; the block height (hence
    # vertical centering) must reflect the wrapped height, not 1 line each.
    bio = "\n".join(f'<li>項目{k}: ' + "長い説明文がここに入って必ず折り返す" * 2 + "</li>"
                    for k in range(4))
    md = ("<!-- _class: profile -->\n# P\n"
          '<div class="pf-container">\n'
          '<div class="pf-name">山田 太郎</div>\n'
          '<div class="pf-affiliation">所属</div>\n'
          f'<div class="pf-bio">\n{bio}\n</div>\n</div>')
    b = _make_builder()
    b.build_profile(parse_slide(0, md))
    bio_box = next(s for s in b.prs.slides[0].shapes
                   if getattr(s, "has_text_frame", False) and "項目0" in s.text_frame.text)
    # 4 bullets each wrapping to ~2 lines → box must be much taller than 4 lines
    from marp_pptx.visuallint import _needed_height_emu
    assert bio_box.height >= _needed_height_emu(bio_box) - int(Inches(0.05))


def test_collision_detector_flags_off_slide_overflow():
    from pptx import Presentation
    from marp_pptx.visuallint import detect_text_collisions
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # box low on the slide with content that runs well past the bottom edge
    _mk_textbox(s, 1, 6.5, 8, 0.4, "下端を確実に超える長い文章" * 6, 24)
    hits = detect_text_collisions(prs)
    assert any("off the slide bottom" in h["onto"] for h in hits)


def test_collision_detector_ignores_bottom_footnote():
    from pptx import Presentation
    from marp_pptx.visuallint import detect_text_collisions
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # a one-line footnote anchored near the bottom must NOT trip the off-slide
    # check (0.2in margin absorbs the estimate's slight over-count)
    _mk_textbox(s, 0.5, 7.0, 9, 0.3, "脚注テキストが一行で入る程度の長さ", 11)
    assert detect_text_collisions(prs) == []


# --- 17. diagram/chart captions must wrap instead of running off-slide ------
# python-pptx textboxes default to wrap="none"; a multi-sentence CJK caption
# rendered as one line past both slide edges.

def _find_shape_with_text(slide, token):
    for s in slide.shapes:
        if getattr(s, "has_text_frame", False) and token in s.text_frame.text:
            return s
    return None


def test_diagram_long_caption_wraps_and_grows():
    long_cap = "PV の逆潮流で末端ほど電圧が上昇し、上限に触れた時点の合計連系量が受入限界を決める。" * 3
    md = ('<!-- _class: diagram -->\n'
          '# T\n\n'
          f'<div class="caption">{long_cap}</div>')
    b = _make_builder()
    b.build_diagram(parse_slide(0, md))
    cap = _find_shape_with_text(b.prs.slides[0], "逆潮流")
    assert cap is not None
    assert cap.text_frame.word_wrap
    assert cap.height > int(Inches(0.5))


def test_diagram_short_caption_keeps_legacy_height():
    md = ('<!-- _class: diagram -->\n'
          '# T\n\n'
          '<div class="caption">短い一行キャプション</div>')
    b = _make_builder()
    b.build_diagram(parse_slide(0, md))
    cap = _find_shape_with_text(b.prs.slides[0], "キャプション")
    assert cap is not None
    assert cap.height == int(Inches(0.5))


def test_chart_long_caption_wraps_and_grows():
    long_cap = "同一計算機・単スレッドでの実測。提案手法はノード数に対しほぼ線形にスケールする。" * 3
    md = ('<!-- _class: chart -->\n'
          '<!-- _chart: column -->\n'
          '# T\n\n'
          '| x | a | b |\n|---|---|---|\n| 1 | 2 | 3 |\n\n'
          f'<div class="chart-caption">{long_cap}</div>')
    b = _make_builder()
    b.build_chart(parse_slide(0, md))
    cap = _find_shape_with_text(b.prs.slides[0], "同一計算機")
    assert cap is not None
    assert cap.text_frame.word_wrap
    assert cap.height > int(Inches(0.4))


# --- 18. cols-2 column heads must start at one shared top --------------------
# Per-column vertical centering put short and tall columns at different
# heights, so the two column headings never lined up.

def test_columns_share_one_top():
    md = ('<!-- _class: cols-2 -->\n'
          '# T\n\n'
          '<div class="columns">\n<div>\n\n'
          '## 短い列\n\n- 一行だけ\n\n'
          '</div>\n<div>\n\n'
          '## 長い列\n\n' + "\n".join(f"- 説明テキスト{i}が長く続く" for i in range(8)) + '\n\n'
          '</div>\n</div>')
    b = _make_builder()
    b.build_columns(parse_slide(0, md))
    short = _find_shape_with_text(b.prs.slides[0], "短い列")
    long_ = _find_shape_with_text(b.prs.slides[0], "長い列")
    assert short is not None and long_ is not None
    assert short.top == long_.top


# --- 19. png-mode inline-math fallback must unwrap \mathrm{...} --------------
# \mathrm{pv} degraded to "mathrmpv" (the brace strip kept the command name).

def test_math_text_fallback_unwraps_wrappers():
    b = _make_builder()
    assert b._math_text_fallback(r"P_{\mathrm{pv}}(t)") == "Ppv(t)"
    assert b._math_text_fallback(r"\text{SOC}_{\min}") == "SOCmin"
    assert b._math_text_fallback(r"\mathbf{x}^\top") == "x⊤"  # \top maps to the glyph


# --- 20. title_align: left must left-align the whole hero stack --------------
# ThemeLayout.title_align existed but build_title hardcoded CENTER.

def test_math_annotate_directives_are_parsed():
    md = ('<!-- _class: equation -->\n'
          '# DFT\n\n'
          '<div class="eq-main">\n\n'
          '$$\n'
          'X_k % [!math-annotate label="出力" note="k 番目の周波数成分"]\n'
          '= \\sum_{n=0}^{N-1} % [!math-annotate note="全サンプルの総和"]\n'
          'x_n e^{-i 2\\pi kn/N}\n'
          '$$\n\n'
          '</div>')
    sd = parse_slide(0, md)
    assert len(sd.eq_annotations) == 3
    assert sd.eq_annotations[0] == ("X_k", "出力", "k 番目の周波数成分", None)
    assert sd.eq_annotations[1][2] == "全サンプルの総和"
    assert sd.eq_annotations[2][2] is None          # plain line keeps no note
    assert "[!math-annotate" not in sd.eq_main      # directives stripped


def test_math_annotate_builds_cards_and_connectors():
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    md = ('<!-- _class: equation -->\n'
          '# DFT\n\n'
          '<div class="eq-main">\n\n'
          '$$\n'
          'X_k % [!math-annotate label="出力" note="k 番目の周波数成分"]\n'
          '= \\sum_{n=0}^{N-1} x_n % [!math-annotate note="全サンプルの総和" color="#0d558e"]\n'
          '$$\n\n'
          '</div>')
    b = _make_builder()
    b.build_equation(parse_slide(0, md))
    slide = b.prs.slides[0]
    assert _find_shape_with_text(slide, "k 番目の周波数成分") is not None
    assert _find_shape_with_text(slide, "全サンプルの総和") is not None
    lines = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(lines) == 2                          # one connector per note


# --- 22. tmu-cs directive: code [!step N action[:M]] -------------------------

def test_code_steps_parse_and_strip():
    code = ('for i in range(10):   # ループ [!step 1 warning]\n'
            '    print(i)          # [!step 2 info]\n'
            'done = True')
    md = ('<!-- _class: code -->\n# T\n\n<div class="cd-code">\n\n'
          '```python\n' + code + '\n```\n\n</div>')
    sd = parse_slide(0, md)
    assert sd.code_steps == [(0, 1, "warning", 1), (1, 2, "info", 1)]
    lines = sd.code_text.split("\n")
    assert lines[0].endswith("# ループ")            # visible comment kept
    assert lines[1] == "    print(i)"               # directive-only comment removed
    assert "[!step" not in sd.code_text


def test_step_slides_expand_per_step(tmp_path):
    md = ("---\nmarp: true\n---\n\n"
          '<!-- _class: code -->\n# T\n\n<div class="cd-code">\n\n'
          '```python\na = 1  # [!step 1 highlight]\nb = 2  # [!step 2 focus]\n```\n\n</div>\n')
    f = tmp_path / "t.md"
    f.write_text(md, encoding="utf-8")
    from marp_pptx.parser import parse_marp
    slides = parse_marp(f)
    assert len(slides) == 2                          # one slide per step
    assert [s.active_step for s in slides] == [1, 2]
    assert [s.index for s in slides] == [0, 1]


def test_code_step_band_renders_behind_text():
    md = ('<!-- _class: code -->\n# T\n\n<div class="cd-code">\n\n'
          '```python\na = 1  # [!step 1 warning]\nb = 2\n```\n\n</div>')
    sd = parse_slide(0, md)
    sd.active_step = 1
    b = _make_builder()
    b.build_code(sd)
    slide = b.prs.slides[0]
    fills = []
    for s in slide.shapes:
        try:
            fills.append(s.fill.fore_color.rgb)
        except (TypeError, AttributeError):
            pass
    from pptx.dml.color import RGBColor as _C
    assert _C(0x45, 0x3E, 0x28) in fills            # warning band
    assert _C(0xF9, 0xE2, 0xAF) in fills            # warning tick bar
    assert _find_shape_with_text(slide, "STEP 1 / 1") is not None


def test_blocks_parse_variants_and_content():
    md = ('<!-- _class: blocks -->\n# T\n\n'
          '<div class="bk-container">\n\n'
          '<div class="bk theorem">\n'
          '  <span class="bk-title">定理 1</span>\n'
          '  <span class="bk-body">本文 A</span>\n'
          '</div>\n\n'
          '<div class="bk alert">\n'
          '  <span class="bk-title">注意</span>\n'
          '  <span class="bk-body">本文 B</span>\n'
          '</div>\n\n'
          '</div>')
    sd = parse_slide(0, md)
    assert sd.block_items == [("theorem", "定理 1", "本文 A"),
                              ("alert", "注意", "本文 B")]


def test_blocks_render_with_variant_colors():
    from pptx.dml.color import RGBColor as _C
    md = ('<!-- _class: blocks -->\n# T\n\n'
          '<div class="bk-container">\n\n'
          '<div class="bk theorem">\n'
          '  <span class="bk-title">定理 1</span>\n'
          '  <span class="bk-body">本文 A</span>\n'
          '</div>\n\n'
          '<div class="bk alert">\n'
          '  <span class="bk-title">注意</span>\n'
          '  <span class="bk-body">本文 B</span>\n'
          '</div>\n\n'
          '</div>')
    b = _make_builder()
    b.build_blocks(parse_slide(0, md))
    slide = b.prs.slides[0]
    assert _find_shape_with_text(slide, "定理 1") is not None
    assert _find_shape_with_text(slide, "本文 B") is not None
    fills = []
    for s in slide.shapes:
        try:
            fills.append(s.fill.fore_color.rgb)
        except (TypeError, AttributeError):
            pass
    assert _C(0x9F, 0x1D, 0x1D) in fills        # alert bar keeps beamer red


def test_footer_bar_renders_title_section_page():
    from marp_pptx.theme import ThemeLayout
    title_md = '<!-- _class: title -->\n# デッキ題名\n## サブ\n著者 / 2026'
    div_md = '<!-- _class: divider -->\n# 1. 理論\n## 収束解析'
    body_md = '<!-- _class: statement -->\n\n本文スライド。'
    b = _make_builder()
    b.theme.layout = ThemeLayout(footer_bar=True)
    b.build_all([parse_slide(0, title_md), parse_slide(1, div_md),
                 parse_slide(2, body_md)])
    last = b.prs.slides[2]
    assert _find_shape_with_text(last, "3 / 3") is not None
    assert _find_shape_with_text(last, "理論") is not None    # section, no "1."
    assert _find_shape_with_text(last, "デッキ題名") is not None
    # the title slide keeps its hero layout (no bar text)
    assert _find_shape_with_text(b.prs.slides[0], "1 / 3") is None


def test_title_align_left_is_honored():
    from pptx.enum.text import PP_ALIGN
    from marp_pptx.theme import ThemeLayout
    md = ('<!-- _class: title -->\n'
          '# 左揃えタイトル\n'
          '## キッカー行\n\n'
          '著者名 / 2026')
    b = _make_builder()
    b.theme.layout = ThemeLayout(title_align="left")
    b.build_title(parse_slide(0, md))
    t = _find_shape_with_text(b.prs.slides[0], "左揃えタイトル")
    assert t is not None
    assert t.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT

    b2 = _make_builder()  # default stays centered
    b2.theme.layout = ThemeLayout(title_align="center")
    b2.build_title(parse_slide(0, md))
    t2 = _find_shape_with_text(b2.prs.slides[0], "左揃えタイトル")
    assert t2.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER


# --- feature-gallery must keep reproducing (math-annotate v2 + steps) -------

def test_feature_gallery_reproduces_with_math(tmp_path):
    """demo/feature-gallery.md is the living demo of the tmu-cs feature port:
    %[!math-annotate] equations (PNG segments + note cards + connectors) and
    #[!step] slide expansion. It must keep building with the math intact and
    stay clean under audit — including the annotation-label / dimmed-code
    contrast that used to sit below AA."""
    from marp_pptx.parser import parse_marp
    from marp_pptx.audit import audit_pptx

    root = Path(__file__).resolve().parents[1]
    md = root / "demo" / "feature-gallery.md"
    slides = parse_marp(md)          # step expansion happens inside
    tc = ThemeConfig()
    tc.math_mode = "png"
    b = PptxBuilder(base_path=md.parent, theme=tc)
    b.build_all(slides)
    assert len(b.prs.slides) == 9          # 7 sources -> the step slide x3
    out = tmp_path / "feature-gallery.pptx"
    b.save(str(out))
    findings = audit_pptx(out)
    assert not [f for f in findings if f.severity == "error"], findings
    assert not [f for f in findings if f.kind == "contrast"], (
        "annotation labels / dimmed code must stay AA-readable")
