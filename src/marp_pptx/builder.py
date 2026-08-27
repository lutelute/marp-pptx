"""PptxBuilder: converts parsed SlideData into PowerPoint slides.

All theme-dependent values come from self.theme (ThemeConfig instance),
eliminating the global state of the original convert_v2.py.
"""
from __future__ import annotations

import re
import sys
import tempfile
import hashlib
from dataclasses import replace as dc_replace
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

from marp_pptx import metrics
from marp_pptx.parser import SlideData, strip_html
from marp_pptx.theme import ThemeConfig
from marp_pptx.layout import (
    SW, SH, MARGIN_L, MARGIN_R, MARGIN_T, MARGIN_B, CONTENT_W,
    TITLE_H, TITLE_TOP, TITLE_GAP, KICKER_H, BODY_TOP, BODY_H, FOOTER_RESERVE,
    SZ_DISPLAY, SZ_TITLE, SZ_KICKER, SZ_H2, SZ_H3, SZ_BODY, SZ_COL, SZ_SMALL,
    SZ_FOOT, SZ_METRIC, SZ_EQ, SZ_EQ_VAR, SZ_ZONE_L, SZ_ZONE_B,
    LINE_BODY, LINE_TITLE, PARA_GAP, BLOCK_GAP, SECTION_GAP,
    CARD_GAP, CARD_PAD, CARD_RADIUS, CARD_ACCENT_H, HAIRLINE_W, ACCENT_RULE_W,
)
from marp_pptx.math.omml import latex_to_omml_element, set_omml_size, omml_glyph_count, OmmlError
from marp_pptx.math.renderer import render_latex_png

try:
    import cairosvg
    HAS_CAIROSVG = True
except ImportError:
    HAS_CAIROSVG = False

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


class PptxBuilder:
    def __init__(self, base_path: Path, theme: ThemeConfig):
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self.base_path = base_path
        self.theme = theme
        self._img_cache: dict = {}
        self._divider_no = 0
        # Non-fatal authoring problems surfaced during build (unknown slide
        # type, missing image, …). Printed once each to stderr and collected
        # here so callers (CLI / web / tests) can report them.
        self.warnings: list[str] = []
        self._warned: set = set()

    def _warn(self, msg: str):
        """Record a non-fatal authoring warning (deduplicated) and echo it to
        stderr. The build still succeeds — this just stops the tool from
        silently producing something other than what the author wrote."""
        if msg in self._warned:
            return
        self._warned.add(msg)
        self.warnings.append(msg)
        print(f"[warn] {msg}", file=sys.stderr)

    def _fs(self, pt_val):
        """Scale a Pt value by theme.font_scale. Min 8pt.

        Accepts either a Pt object or a raw int (points).
        """
        from pptx.util import Pt as _Pt
        scale = getattr(self.theme, "font_scale", 1.0)
        try:
            base = pt_val.pt
        except AttributeError:
            base = float(pt_val)
        return _Pt(max(8, base * scale))

    def _unscale(self, pt_val):
        """Undo `_fs`, for a size already computed in rendered points that is
        about to be handed to a helper which scales it again."""
        from pptx.util import Pt as _Pt
        scale = getattr(self.theme, "font_scale", 1.0) or 1.0
        base = pt_val.pt if hasattr(pt_val, "pt") else float(pt_val)
        return _Pt(base / scale)

    def _ms(self, inch_val):
        """Scale an Inches value by theme.margin_scale. Min 0.1 inch.

        Accepts either an Inches object (EMU) or a raw float (inches).
        Used for margins, padding, gaps, and box positions/sizes.
        """
        from pptx.util import Inches as _Inches
        scale = getattr(self.theme, "margin_scale", 1.0)
        try:
            base = inch_val.inches
        except AttributeError:
            base = float(inch_val)
        return _Inches(max(0.1, base * scale))

    # ══════════════════════════════════════════════
    # Design-system helpers (refined-minimal, center-balanced)
    # ══════════════════════════════════════════════
    def _tint(self, rgb, t):
        """Blend rgb toward white by t (0=rgb, 1=white). For hairlines/ghosts."""
        return RGBColor(
            int(rgb[0] + (255 - rgb[0]) * t),
            int(rgb[1] + (255 - rgb[1]) * t),
            int(rgb[2] + (255 - rgb[2]) * t),
        )

    def _text_safe(self, color: RGBColor, bg: RGBColor | None = None,
                   min_ratio: float = 4.5) -> RGBColor:
        """Darken `color` toward black until it reads on `bg` (measured).

        For small text drawn in a decorative color (annotation-card labels,
        user-supplied hexes) — the hue survives, the contrast is guaranteed.
        """
        from marp_pptx.audit import contrast_ratio
        if bg is None:
            bg = self.SURFACE
        c = color
        for _ in range(14):
            if contrast_ratio(tuple(c), tuple(bg)) >= min_ratio:
                return c
            c = RGBColor(int(c[0] * 0.82), int(c[1] * 0.82), int(c[2] * 0.82))
        return c

    def _hero_fill_color(self) -> RGBColor | None:
        """ThemeLayout.hero_fill parsed, or None. Dark/light is decided by
        measured luminance where it's used, not by naming convention."""
        hx = (self.LAYOUT.hero_fill or "").lstrip("#")
        if len(hx) != 6:
            return None
        try:
            return RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        except ValueError:
            return None

    def _hero_accent(self, min_ratio: float = 2.5) -> RGBColor:
        """Accent for rules/labels on a PRIMARY-dark hero slide.

        The theme accent when it actually reads against PRIMARY (measured
        WCAG ratio, not guessed), white otherwise. min_ratio: 2.5 for
        graphics (rules), 4.5 for small text (kickers).
        """
        from marp_pptx.audit import contrast_ratio
        if contrast_ratio(tuple(self.ACCENT), tuple(self.PRIMARY)) >= min_ratio:
            return self.ACCENT
        return self.WHITE

    @property
    def HAIRLINE(self):
        return getattr(self.theme, "hairline", None) or self._tint(self.MUTED, 0.62)

    @property
    def SURFACE(self):
        # Slightly off-white panel fill that reads against pure-white bg.
        return self.LIGHT

    def _no_shadow(self, shape):
        """Kill python-pptx's heavy default preset shadow (mandatory for minimal)."""
        try:
            shape.shadow.inherit = False
        except Exception:
            pass
        return shape

    def _soft_shadow(self, shape):
        """Subtle downward drop shadow — 15% black, wide blur, short offset.

        Cards on a near-white bg with only a hairline border read as flat
        outlines; a soft shadow is the elevation cue that doesn't add a line.
        Opt-in per theme via ThemeLayout.card_shadow.
        """
        sp_pr = shape.fill._xPr  # <p:spPr> — python-pptx has no shadow API
        for old in sp_pr.findall(qn("a:effectLst")):
            sp_pr.remove(old)
        effect = etree.SubElement(sp_pr, qn("a:effectLst"))
        shdw = etree.SubElement(effect, qn("a:outerShdw"))
        shdw.set("blurRad", "90000")     # ~7pt blur
        shdw.set("dist", "19050")        # 1.5pt down
        shdw.set("dir", "5400000")       # 90° (downward)
        shdw.set("rotWithShape", "0")
        clr = etree.SubElement(shdw, qn("a:srgbClr"))
        clr.set("val", "141413")
        alpha = etree.SubElement(clr, qn("a:alpha"))
        alpha.set("val", "15000")
        return shape

    def _content_region(self, has_title: bool = True, *, full: bool = False,
                        lead: bool = False):
        """(left, top, width, height) EMU for the body canvas.

        has_title: start below the title band; else from the top margin.
        full: don't reserve the footer band (figures / full-bleed images).
        lead: a colored lead line sits under the title — push the body down.
        Single source of truth — every build_* should start from this.
        """
        left = int(MARGIN_L)
        width = int(CONTENT_W)
        footer = 0 if full else int(FOOTER_RESERVE)
        if self.LAYOUT.footer_bar:
            footer += int(Inches(0.30))   # the beamer bar claims the bottom band
        if has_title:
            if self.LAYOUT.h1_deco == "band":
                # thin frametitle band → the body starts higher (denser canvas)
                top = int(Inches(0.08) + TITLE_H + TITLE_GAP)
            else:
                top = int(MARGIN_T + TITLE_H + TITLE_GAP)
        else:
            top = int(MARGIN_T)
        if lead:
            # One lead line at SZ_H3 — tracks font_scale so a scaled deck
            # doesn't clip the line into the body below it.
            top += int(Pt(26 * getattr(self.theme, "font_scale", 1.0)))
        height = int(SH - top - MARGIN_B - footer)
        return (left, top, width, height)

    def _add_lead(self, slide, text):
        """Hearing-style colored lead under the title: 「トピック｜結論」.

        The single most recognizable device of the reference hearing deck —
        every content slide answers "which thread, and what's the claim"
        before the body starts.
        """
        band = self.LAYOUT.h1_deco == "band"
        title_top = int(Inches(0.08)) if band else int(TITLE_TOP)
        top = title_top + int(TITLE_H) + int(Inches(0.08))
        lead_h = int(self._fs(Pt(SZ_H3.pt * metrics.DEFAULT_LINE_FACTOR)))
        tb = self._add_textbox(slide, int(MARGIN_L), top,
                               int(CONTENT_W), lead_h)
        p = tb.text_frame.paragraphs[0]
        self._set_rich_text(p, text, SZ_H3, self.SECONDARY)
        for r in p.runs:
            r.font.bold = True
        return tb

    def _content_region_with_lead(self, slide, sd, *, full: bool = False):
        """_content_region + the h2 lead line for types that don't otherwise
        consume sd.h2. With no h2 the geometry is exactly _content_region."""
        lead = (sd.h2 or "").strip() if (sd.h1 and getattr(sd, "h2", "")) else ""
        region = self._content_region(has_title=bool(sd.h1), full=full,
                                      lead=bool(lead))
        if lead:
            self._add_lead(slide, lead)
        return region

    def _stack_tops(self, heights, region, mode="center", gap=None):
        """Return the top (EMU) for each block stacked vertically in region.

        mode: "center" (block group centered), "justify" (free space spread
        between blocks), "top" (legacy top-pile), "fill" (each block expanded
        to share the region equally — caller uses the returned step as height).
        This is the core fix for the "empty bottom 40-60%" defect.
        """
        left, rtop, width, rheight = region
        if gap is None:
            gap = int(BLOCK_GAP)
        n = len(heights)
        if n == 0:
            return []
        content_h = sum(int(h) for h in heights)
        if mode == "justify" and n > 1:
            free = max(0, rheight - content_h)
            g = free // (n - 1)
            tops, y = [], rtop
            for h in heights:
                tops.append(int(y)); y += int(h) + g
            return tops
        if mode == "fill":
            free = max(0, rheight - gap * (n - 1))
            each = free // n
            tops, y = [], rtop
            for _ in heights:
                tops.append(int(y)); y += each + gap
            return tops
        total = content_h + gap * (n - 1)
        if mode == "center":
            y = rtop + max(0, (rheight - total) // 2)
        else:  # "top"
            y = rtop
        tops = []
        for h in heights:
            tops.append(int(y)); y += int(h) + gap
        return tops

    def _vcenter_shapes(self, shapes, region):
        """Shift already-placed shapes so their bounding box is vertically
        centered in region. Use ONLY for shapes with known height (tables,
        pictures, explicit-height textboxes) — never autofit textboxes."""
        shapes = [s for s in shapes if s is not None]
        if not shapes:
            return
        _, rtop, _, rheight = region
        tops = [s.top for s in shapes]
        bottoms = [s.top + s.height for s in shapes]
        bbox_top, bbox_bot = min(tops), max(bottoms)
        offset = rtop + (rheight - (bbox_bot - bbox_top)) // 2 - bbox_top
        for s in shapes:
            s.top = int(s.top + offset)

    def _hairline(self, slide, left, top, width, thickness=None, color=None,
                  vertical=False):
        """Thin accent/divider rule. Horizontal by default."""
        if thickness is None:
            thickness = HAIRLINE_W
        if color is None:
            color = self.HAIRLINE
        w, h = (int(thickness), int(width)) if vertical else (int(width), int(thickness))
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top), w, h)
        ln.fill.solid()
        ln.fill.fore_color.rgb = color
        ln.line.fill.background()
        self._no_shadow(ln)
        return ln

    _GREEK = {
        "alpha": "α", "beta": "β", "gamma": "γ", "delta": "δ", "epsilon": "ε",
        "zeta": "ζ", "eta": "η", "theta": "θ", "iota": "ι", "kappa": "κ",
        "lambda": "λ", "mu": "μ", "nu": "ν", "xi": "ξ", "pi": "π", "rho": "ρ",
        "sigma": "σ", "tau": "τ", "phi": "φ", "chi": "χ", "psi": "ψ", "omega": "ω",
        "Gamma": "Γ", "Delta": "Δ", "Theta": "Θ", "Lambda": "Λ", "Pi": "Π",
        "Sigma": "Σ", "Phi": "Φ", "Psi": "Ψ", "Omega": "Ω",
        "times": "×", "cdot": "·", "leq": "≤", "geq": "≥", "le": "≤", "ge": "≥",
        "in": "∈", "forall": "∀", "infty": "∞", "nabla": "∇", "partial": "∂",
        "sum": "Σ", "approx": "≈", "neq": "≠", "to": "→", "top": "⊤",
        "sqrt": "√", "odot": "⊙", "otimes": "⊗", "oplus": "⊕", "pm": "±",
        "cup": "∪", "cap": "∩", "subset": "⊂", "supset": "⊃", "propto": "∝",
        "sim": "∼", "equiv": "≡", "langle": "⟨", "rangle": "⟩", "prod": "∏",
        "gtrsim": "≳", "lesssim": "≲", "ll": "≪", "gg": "≫",
        "int": "∫", "exists": "∃", "leftarrow": "←", "rightarrow": "→",
    }

    def _math_text_fallback(self, latex: str) -> str:
        """Approximate a LaTeX snippet as plain text/unicode (png-mode fallback
        when OMML is disabled). Maps common greek/operators, strips the rest."""
        s = latex
        # Unwrap argument-carrying commands BEFORE the generic strip below —
        # otherwise \mathrm{pv} degrades to "mathrmpv" instead of "pv".
        wrapper = re.compile(
            r"\\(?:mathrm|mathbf|mathit|mathcal|mathbb|mathsf|mathtt|text|textrm"
            r"|textbf|textit|bm|boldsymbol|operatorname)\s*\{([^{}]*)\}")
        for _ in range(3):  # a few rounds for nested wrappers
            s2 = wrapper.sub(r"\1", s)
            if s2 == s:
                break
            s = s2
        # \frac{a}{b} -> a/b (otherwise the brace strip leaves "fracab")
        s = re.sub(r"\\[td]?frac\s*\{([^{}]*)\}\s*\{([^{}]*)\}", r"\1/\2", s)
        for name, ch in self._GREEK.items():
            s = re.sub(r"\\" + name + r"(?![a-zA-Z])", ch, s)
        s = s.replace("\\,", " ").replace("\\;", " ").replace("\\ ", " ")
        s = re.sub(r"[\\{}$]", "", s)
        s = s.replace("_", "").replace("^", "")
        return s.strip()

    def _set_tracking(self, para, spc=160):
        """Wide letter-spacing (1/100 pt units) on a paragraph's runs."""
        for run in para.runs:
            run._r.get_or_add_rPr().set("spc", str(int(spc)))

    def _kicker_para(self, para, text, size=None, color=None, align=PP_ALIGN.LEFT,
                     tracking=180):
        """Small-caps-style label: UPPERCASE (latin only) + wide tracking.

        Tracking is a Latin small-caps device. Japanese is already spaced by
        its own em box, so widening it there only makes the label longer than
        the card it sits in — a KPI caption that measured as one line came out
        as two. Latin keeps the full tracking; anything else gets a token
        amount.
        """
        if size is None:
            size = SZ_KICKER
        if color is None:
            color = self.ACCENT_TEXT
        t = (text or "")
        latin = t.isascii()
        para.text = t.upper() if latin else t
        para.alignment = align
        para.font.name = self.FONT_HEAD
        para.font.size = self._fs(size)
        para.font.bold = True
        para.font.color.rgb = color
        self._set_tracking(para, tracking if latin else min(tracking, 40))
        return para

    def _eyebrow(self, slide, text, left, top, width, align=PP_ALIGN.LEFT,
                 color=None, size=None):
        """Standalone small-caps eyebrow label textbox."""
        tb = self._add_textbox(slide, int(left), int(top), int(width),
                               int(KICKER_H))
        self._kicker_para(tb.text_frame.paragraphs[0], text, size=size,
                          color=color or self.ACCENT_TEXT, align=align)
        return tb

    def _add_card(self, slide, region, *, label="", body="", body_lines=None,
                  value=None, accent=None, accent_bar=False, fill=None,
                  anchor="top", label_size=None, body_size=None,
                  value_size=None, label_color=None):
        """Refined-minimal card: surface fill + hairline border, no drop shadow.

        Elevation cue is the optional top accent bar (accent_bar=True), used
        sparingly (KPI, highlighted panels). Returns (bg_shape, text_shape).
        """
        left, top, width, height = (int(v) for v in region)
        if accent is None:
            accent = self.ACCENT
        if fill is None:
            fill = self.SURFACE
        pad = int(CARD_PAD)
        acc_h = int(CARD_ACCENT_H) if accent_bar else 0

        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        bg.adjustments[0] = CARD_RADIUS
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = self.HAIRLINE
        bg.line.width = HAIRLINE_W
        if self.LAYOUT.card_shadow:
            self._soft_shadow(bg)
        else:
            self._no_shadow(bg)

        if accent_bar:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, acc_h)
            bar.fill.solid(); bar.fill.fore_color.rgb = accent
            bar.line.fill.background()
            self._no_shadow(bar)

        tx = left + pad
        ty = top + acc_h + pad
        tw = width - pad * 2
        th = height - acc_h - pad * 2
        tb = self._add_textbox(slide, tx, ty, max(int(Inches(0.3)), tw),
                               max(int(Inches(0.2)), th))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE if anchor == "middle" else MSO_ANCHOR.TOP

        first = True
        if value is not None:
            p = tf.paragraphs[0]; first = False
            run = p.add_run(); run.text = value
            run.font.name = self.FONT_HEAD
            run.font.size = value_size or self._fs(SZ_METRIC)
            run.font.bold = True
            run.font.color.rgb = accent
            p.alignment = PP_ALIGN.CENTER
        if label:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            if not first:
                p.space_before = Pt(6)
            first = False
            align = PP_ALIGN.CENTER if value is not None else PP_ALIGN.LEFT
            lsize = label_size or SZ_ZONE_L
            if value is not None:
                lsize = self._fs(lsize)        # fit at the size it renders
                # A metric caption belongs on one line: wrapping it orphans a
                # closing bracket and pushes the values out of alignment with
                # the card next to it. The label is set bold and tracked, so
                # both come off the width before fitting — measuring the plain
                # face against the full width left it a point too big.
                track = int(Pt(0.4 * len(label)))
                lsize = self._fit_size(label, lsize,
                                       max(int(Inches(0.5)), int((tw - track) * 0.97)),
                                       int(Pt(lsize.pt * metrics.DEFAULT_LINE_FACTOR)),
                                       bold=True, min_ratio=0.75)
                lsize = self._unscale(lsize)  # _kicker_para re-applies the scale
            self._kicker_para(p, label, size=lsize,
                              color=label_color or self.SECONDARY, align=align,
                              tracking=120)
        if body or body_lines:
            if first:
                p2 = tf.paragraphs[0]
            else:
                p2 = tf.add_paragraph(); p2.space_before = Pt(6)
            # A card is a fixed cell in a grid: its copy cannot push the layout
            # around, so when it doesn't fit, the type gives way. Only inside
            # cards — prose that overruns should be edited, not shrunk.
            bsize = self._fit_card_body(body, body_lines, tw, th,
                                        body_size or SZ_ZONE_B,
                                        value, value_size, label, label_size)
            if body_lines:
                self._fill_multiline_box(tf, "\n".join(body_lines), bsize, self.FG)
            else:
                self._set_text_with_inline_math(p2, body, bsize, self.FG)
        return bg, tb

    def _fit_card_body(self, body, body_lines, tw, th, size,
                       value=None, value_size=None, label="", label_size=None):
        """Body size that fits what is left of the card after value + label."""
        text = "\n".join(body_lines) if body_lines else (body or "")
        if not text or tw <= 0 or th <= 0:
            return size
        used = 0.0
        if value is not None:
            used += (value_size or self._fs(SZ_METRIC)).pt * metrics.DEFAULT_LINE_FACTOR
        if label:
            used += (self._fs(label_size or SZ_ZONE_L).pt
                     * metrics.DEFAULT_LINE_FACTOR + 6)
        avail = th - int(Pt(used))
        if avail <= int(Pt(8)):
            return size
        # Fit at the size it renders at, then hand back an unscaled value —
        # the callers re-apply font_scale.
        fitted = self._fit_size(text, self._fs(size), tw, avail, min_ratio=0.72)
        return self._unscale(fitted)

    def save(self, path: str):
        self._ensure_ea_font()
        self.prs.save(path)

    # ── Theme shortcuts ──
    @property
    def PRIMARY(self): return self.theme.primary
    @property
    def SECONDARY(self): return self.theme.secondary
    @property
    def ACCENT(self): return self.theme.accent
    @property
    def ACCENT_TEXT(self):
        # Darker accent for small text. A palette may hand-pick it; without
        # one, darken the accent until it measures 4.5:1 on the slide bg —
        # a bare palette used to inherit whatever ratio the accent happened
        # to have (default theme: 3.8:1).
        picked = getattr(self.theme, "accent_text", None)
        if picked is not None:
            return picked
        return self._text_safe(self.theme.accent, self.theme.bg)
    @property
    def FG(self): return self.theme.fg
    @property
    def MUTED(self): return self.theme.muted
    @property
    def LIGHT(self): return self.theme.light
    @property
    def WHITE(self): return self.theme.white
    @property
    def FONT(self): return self.theme.font
    @property
    def FONT_HEAD(self): return self.theme.font_head
    @property
    def FONT_EA(self): return self.theme.font_ea
    @property
    def FONT_MONO(self): return self.theme.font_mono
    @property
    def LAYOUT(self): return self.theme.layout

    # ── EA font injection ──
    def _ensure_ea_font(self):
        """Patch every run-property element with an East-Asian typeface so
        CJK text uses the EA font in PowerPoint. Covers slides AND the slide
        master/layout defaults (so inherited runs are handled too)."""
        rpr_tags = (f"{{{NS_A}}}rPr", f"{{{NS_A}}}defRPr", f"{{{NS_A}}}endParaRPr")
        roots = [s._element for s in self.prs.slides]
        for layout in self.prs.slide_layouts:
            roots.append(layout._element)
            roots.append(layout.slide_master._element)
        seen = set()
        for root in roots:
            if id(root) in seen:
                continue
            seen.add(id(root))
            for tag in rpr_tags:
                for rpr in root.iter(tag):
                    self._patch_rpr(rpr)

    def _patch_rpr(self, rpr):
        """Ensure <a:ea> and <a:cs> typefaces exist, in correct schema order
        (latin, ea, cs). Works even when <a:latin> is absent."""
        A = f"{{{NS_A}}}"

        def _insert_ordered(child, after_tags):
            for t in after_tags:
                prev = rpr.find(f"{A}{t}")
                if prev is not None:
                    prev.addnext(child)
                    return
            rpr.insert(0, child)

        ea = rpr.find(f"{A}ea")
        if ea is None:
            ea = etree.Element(f"{A}ea")
            _insert_ordered(ea, ["latin"])
        ea.set("typeface", self.FONT_EA)
        cs = rpr.find(f"{A}cs")
        if cs is None:
            cs = etree.Element(f"{A}cs")
            _insert_ordered(cs, ["ea", "latin"])
        cs.set("typeface", self.FONT_EA)

    # ── Math helpers ──
    def _omml_element(self, latex: str, display: bool, pt_size: int | None = None):
        # Allow the theme to force PNG fallback (preview mode for viewers
        # with poor OMML support like LibreOffice).
        if getattr(self.theme, "math_mode", "omml") == "png":
            return None
        try:
            el = latex_to_omml_element(latex, display=display)
        except OmmlError as e:
            print(f"  OMML failed: {latex[:40]}... ({e}) — falling back to PNG", file=sys.stderr)
            return None
        if pt_size is not None:
            # Pin run size in half-points so every renderer draws it identically.
            eff = max(8.0, float(pt_size) * getattr(self.theme, "font_scale", 1.0))
            set_omml_size(el, int(round(eff * 2)))
        return el

    def _add_math_omml_display(self, slide, latex, left, top, width, pt_size=30):
        el = self._omml_element(latex, display=True, pt_size=pt_size)
        if el is None:
            return None
        from marp_pptx.math.omml import NS_M
        has_frac = el.find(f".//{{{NS_M}}}f") is not None
        line_factor = 2.4 if has_frac else 1.6
        eff_pt = max(8.0, pt_size * getattr(self.theme, "font_scale", 1.0))
        height = Emu(int(eff_pt * line_factor * 12700))   # 1pt = 12700 EMU
        tb = self._add_textbox(slide, left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = ""
        run.font.name = self.FONT
        run.font.size = self._fs(Pt(pt_size))
        run.font.color.rgb = self.FG
        p._p.append(el)
        return tb

    def _append_math_omml_inline(self, para, latex, size, color):
        pt = size.pt if hasattr(size, "pt") else float(size)
        el = self._omml_element(latex, display=False, pt_size=int(round(pt)))
        if el is None:
            return False
        run = para.add_run()
        run.text = ""
        run.font.name = self.FONT
        run.font.size = size
        if color is not None:
            run.font.color.rgb = color
        para._p.append(el)
        return True

    _MATH_DPI = 220   # crisp on retina / when projected

    def _render_math(self, latex: str, display: bool = False, fontsize: int = 28) -> str | None:
        """Render LaTeX to PNG via matplotlib. Returns path to PNG."""
        color_hex = f"#{self.FG}"
        return render_latex_png(latex, fontsize=fontsize, display=display,
                                color=color_hex, dpi=self._MATH_DPI)

    def _add_math_image(self, slide, latex, left, top, max_width, display=True, fontsize=28):
        png = self._render_math(latex, display=display, fontsize=fontsize)
        if not png:
            return None
        from PIL import Image
        with Image.open(png) as im:
            iw, ih = im.size
        dpi = self._MATH_DPI
        pw = int(iw * 914400 / dpi)
        ph = int(ih * 914400 / dpi)
        if pw > max_width:
            scale = max_width / pw
            pw = int(pw * scale)
            ph = int(ph * scale)
        img_left = left + (max_width - pw) // 2
        slide.shapes.add_picture(png, img_left, top, pw, ph)
        return (pw, ph)

    # ── Basic shape helpers ──
    def _blank_slide(self):
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        # Apply the theme background (e.g. warm cream) to every slide. Hero
        # builders may override this with their own bg.
        bg = self.theme.bg
        if (bg[0], bg[1], bg[2]) != (0xff, 0xff, 0xff):
            self._set_bg(slide, bg)
        return slide

    def _add_textbox(self, slide, left, top, width, height):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        return tb

    def _set_bg(self, slide, color):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = color

    def _set_gradient_bg(self, slide, c1, c2):
        bg = slide.background
        bg.fill.gradient()
        bg.fill.gradient_stops[0].color.rgb = c1
        bg.fill.gradient_stops[0].position = 0.0
        bg.fill.gradient_stops[1].color.rgb = c2
        bg.fill.gradient_stops[1].position = 1.0

    def _add_title(self, slide, text, top=None, color=None, kicker=None,
                   width=None):
        """Slide title: a small-caps kicker (provides the color cue) + the H1,
        separated from the body by whitespace.

        We deliberately AVOID a decorative accent line under the title — a
        repeated underline rule is a hallmark of AI-generated decks; hierarchy
        comes from whitespace + the kicker + size contrast instead. Optional
        under-title rules are only drawn when a palette explicitly opts in via
        ThemeLayout (h1_deco='bottom-line' for a neutral hairline, or
        accent_rule='short-left' for a colored tick)."""
        if color is None:
            color = self.PRIMARY
        band = self.LAYOUT.h1_deco == "band"
        if top is None:
            # The frametitle band hugs the top edge — pull the title up so the
            # band stays thin (real beamer) and the body area grows below it.
            top = int(Inches(0.08)) if band else TITLE_TOP
        text_left = int(MARGIN_L)
        text_w = int(width if width is not None else CONTENT_W)
        if self.LAYOUT.h1_deco == "left-bar":
            text_w -= int(Pt(6) + Pt(12))

        # Fit the headline to the band instead of letting a long one spill into
        # the body. One line is the design; a second is allowed (with the band
        # grown to match) before the type shrinks further.
        title_pt, title_h, top = self._fit_title(text, text_w, top)

        # Full-width headline band (beamer frametitle): the title sits in a
        # structure-colored band flush with the top edge.
        if band:
            bb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(SW),
                                        int(top + title_h + Inches(0.06)))
            bb.fill.solid(); bb.fill.fore_color.rgb = self.PRIMARY
            bb.line.fill.background(); self._no_shadow(bb)
            color = self.WHITE

        # Small-caps eyebrow — the title's color accent without a line.
        if kicker:
            kb = self._add_textbox(slide, text_left, int(top - KICKER_H + Inches(0.02)),
                                   text_w, int(KICKER_H))
            self._kicker_para(kb.text_frame.paragraphs[0], kicker, color=self.ACCENT_TEXT)

        left_bar = self.LAYOUT.h1_deco == "left-bar"
        if left_bar:
            deco_w = Pt(6)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(MARGIN_L), int(top), int(deco_w), int(title_h))
            bar.fill.solid(); bar.fill.fore_color.rgb = self.ACCENT
            bar.line.fill.background(); self._no_shadow(bar)
            text_left = int(MARGIN_L + deco_w + Pt(12))

        tb = self._add_textbox(slide, text_left, int(top), text_w, int(title_h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.FONT_HEAD
        p.font.size = Pt(title_pt)
        p.font.bold = True
        p.font.color.rgb = color
        p.line_spacing = LINE_TITLE

        # Opt-in under-title rules only (off by default to avoid the AI-deck look).
        if not left_bar:
            rule_y = int(top + title_h + Inches(0.04))
            if self.LAYOUT.h1_deco == "bottom-line":
                self._hairline(slide, MARGIN_L, rule_y, CONTENT_W, color=self.HAIRLINE)
            if self.LAYOUT.accent_rule == "short-left":
                self._hairline(slide, MARGIN_L, rule_y, Inches(0.7),
                               thickness=ACCENT_RULE_W, color=self.ACCENT)
        return tb

    def _fit_title(self, text: str, text_w: int, top: int) -> tuple[float, int, int]:
        """(font size in pt, band height in EMU, adjusted top) for one H1.

        A headline longer than the band used to keep its 30pt and spill over
        the body — invisible at author time, obvious in the room. The order of
        preference is: keep the design (one line at full size) → shrink a
        little to keep it on one line → let it take two lines and grow the band
        around it → shrink until two lines fit.
        """
        # title_scale adds emphasis, but never past the band's two-line
        # envelope (~1.32 × SZ_TITLE): beyond that a wrapped title pushes the
        # band into the body no matter how the whitespace is redistributed.
        # With font_scale 1.3 the clamp lands where the pre-title_scale
        # design already sat, so large-scale decks keep their old geometry.
        base = min(self._fs(SZ_TITLE).pt * getattr(self.LAYOUT, "title_scale", 1.0),
                   SZ_TITLE.pt * 1.32)
        if not text:
            return base, int(TITLE_H), int(top)
        width_pt = text_w / 12700.0
        plain = self._plain(text)
        pitch = LINE_TITLE * metrics.DEFAULT_LINE_FACTOR

        def lines_at(size: float) -> int:
            return metrics.line_count(plain, self.FONT_HEAD, size, width_pt,
                                      ea_font=self.FONT_EA, bold=True)

        def one_line(size: float):
            # Even a single line can outgrow the band once font_scale and
            # title_scale stack (30 × 1.3 × 1.13 = 44pt needs 58pt); grow the
            # band out of the whitespace above instead of letting audit see a
            # clipped title.
            needed = Pt(size * pitch)
            if needed <= TITLE_H:
                return size, int(TITLE_H), int(top)
            grow = int(needed) - int(TITLE_H)
            new_top = max(int(Inches(0.08)), int(top) - grow // 2)
            return size, int(needed), new_top

        if lines_at(base) <= 1:
            return one_line(base)

        # A modest shrink often pulls a title back onto one line; anything
        # deeper reads as a different design, so stop and go two-line instead.
        size = base
        while size > base * 0.82:
            size -= 0.5
            if lines_at(size) <= 1:
                return one_line(size)

        size = base
        while size > base * 0.60 and lines_at(size) > 2:
            size -= 0.5
        needed = Pt(size * pitch * min(2, lines_at(size)))
        if needed <= TITLE_H:
            return size, int(TITLE_H), int(top)
        # Grow the band symmetrically, but never above the top margin: the
        # extra room comes out of the whitespace above and the title→body gap.
        grow = int(needed) - int(TITLE_H)
        new_top = max(int(Inches(0.08)), int(top) - grow // 2)
        return size, int(needed), new_top

    def _add_para(self, tf, text, size=None, color=None, bold=False, italic=False, space_before=Pt(4)):
        if size is None:
            size = SZ_BODY
        if color is None:
            color = self.FG
        p = tf.add_paragraph()
        p.text = text
        p.font.name = self.FONT
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_before = space_before
        return p

    def _visual_em_width(self, s: str) -> float:
        """Display width of a string in em units, from the real font.

        Was a two-bucket guess (1.0 em for CJK, 0.55 for everything else),
        which is off by 3x between "llll" and "WWWW" and misjudged every
        mixed-script line. `metrics` reads the advance widths out of the
        typeface this deck actually names.
        """
        return metrics.measure_em(s, self.FONT, ea_font=self.FONT_EA)

    @staticmethod
    def _plain(s: str) -> str:
        """Strip inline markup so measurement sees the rendered characters."""
        return (s.replace("**", "").replace("`", "").replace("$", "")
                 .replace("==", "").lstrip("#").lstrip())

    def _wrapped_lines(self, text: str, size_pt: float, width_emu, *,
                       bold: bool = False) -> int:
        """How many lines `text` takes in a `width_emu`-wide box."""
        if not width_emu:
            return 1
        return metrics.line_count(self._plain(text), self.FONT, size_pt,
                                  width_emu / 12700.0, ea_font=self.FONT_EA,
                                  bold=bold)

    def _estimate_text_height(self, lines, size, width=None, gap=None,
                              line_spacing=1.0):
        """Tight height for a block of markdown lines at given font size.

        Slightly over-estimates so the shape hugs content but never clips.
        When `width` (EMU) is given, wrapping is measured with real font
        metrics — including 禁則処理 and Latin word boundaries — so callers
        stacking blocks below this one don't overlap when a line wraps.

        `gap` is the paragraph spacing the caller will actually set (points, or
        a Pt). It defaults to the 4pt of `_add_para`; pass `PARA_GAP` when the
        block uses the wider list rhythm, or the box comes out short.

        `line_spacing` must match what the caller sets on the paragraph
        (`LINE_BODY`, `LINE_PROSE`, …) — a block written at 1.30 and reserved
        at 1.0 is 30% short, which is enough to clip the last line.
        """
        from pptx.util import Pt as _Pt
        base = size.pt if hasattr(size, "pt") else float(size)
        scale = getattr(self.theme, "font_scale", 1.0)
        eff = base * scale
        line_h = eff * metrics.DEFAULT_LINE_FACTOR * line_spacing
        gap_pt = 4.0 if gap is None else (gap.pt if hasattr(gap, "pt") else float(gap))
        total = 0.0
        first = True
        for line in lines:
            s = line.strip()
            if not s:
                continue
            mult = 1.35 if s.startswith(("## ", "### ")) else 1.0
            wraps = self._wrapped_lines(s, eff * mult, width) if width else 1
            total += line_h * mult * wraps
            if not first:
                total += gap_pt      # written as a literal Pt(); never scaled
            first = False
        # +6pt tail breathing room to prevent clipping when autofit is off
        return _Pt(max(18, total + 6))

    def _fit_size(self, text: str, size, width_emu, height_emu, *,
                  bold: bool = False, min_ratio: float = 0.72,
                  font: str | None = None) -> "Pt":
        """Largest size ≤ `size` at which `text` fits the given box.

        Titles, captions and card labels are written by an author who cannot
        see the box; this keeps a long one inside it instead of clipping. The
        floor is `min_ratio` of the requested size — below that the content is
        genuinely too long and shrinking further would just make it unreadable.
        """
        from pptx.util import Pt as _Pt
        base = size.pt if hasattr(size, "pt") else float(size)
        if not text or not width_emu or not height_emu:
            return _Pt(base)
        fitted = metrics.fit_size(
            self._plain(text), font or self.FONT, width_emu / 12700.0,
            height_emu / 12700.0, max_size=base, min_size=base * min_ratio,
            ea_font=self.FONT_EA, bold=bold)
        return _Pt(fitted)

    def _add_body_text(self, slide, lines, left=None, top=None, width=None, height=None, size=None):
        # Safety net: \x00 box sentinels (parser.column_lines) are consumed by
        # _add_column_content; any that reach a plain text path must not print.
        lines = [l for l in lines if not l.strip().startswith("\x00")]
        if size is None:
            size = SZ_BODY
        if left is None: left = MARGIN_L
        if top is None: top = BODY_TOP
        if width is None: width = CONTENT_W
        explicit_height = height is not None
        if height is None:
            estimated = self._estimate_text_height(lines, size, width=width)
            height = min(BODY_H, int(estimated))

        tb = self._add_textbox(slide, left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        # When the caller specified a height (zone reservation), lock the size
        # so viewers that ignore spAutoFit don't grow the textbox and break
        # downstream placement (conclusion boxes, footnotes, etc.).
        if explicit_height:
            tf.auto_size = MSO_AUTO_SIZE.NONE

        first = True
        for line in lines:
            s = line.strip()
            if not s:
                continue
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            is_h2 = s.startswith("## ")
            is_h3 = s.startswith("### ")
            is_bullet = s.startswith("- ") or s.startswith("* ")
            is_numbered = re.match(r"^\d+\.\s", s)

            if is_h2:
                p.text = strip_html(s[3:])
                p.font.name = self.FONT_HEAD
                p.font.size = self._fs(SZ_H2)
                p.font.bold = True
                p.font.color.rgb = self.SECONDARY
                p.space_before = Pt(10)
            elif is_h3:
                p.text = strip_html(s[4:])
                p.font.name = self.FONT_HEAD
                p.font.size = self._fs(SZ_H3)
                p.font.bold = True
                p.font.color.rgb = self.MUTED
                p.space_before = Pt(6)
            elif is_bullet:
                # Apply rich text (bold markdown **...**) to bullet content
                self._set_rich_text(p, s[2:], size, self.FG)
                p.level = 0
                p.space_before = Pt(4)
                pPr = p._p.get_or_add_pPr()
                buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
                for existing in pPr.findall(qn("a:buChar")):
                    pPr.remove(existing)
                for existing in pPr.findall(qn("a:buNone")):
                    pPr.remove(existing)
                pPr.append(buChar)
            elif is_numbered:
                # Numbered items get the same inline treatment as every other
                # line. They used to be written verbatim, so a step written as
                # "1. $f$ の **平滑性**" reached the slide with its markup and
                # its LaTeX showing.
                self._set_rich_text(p, s, size, self.FG)
                p.space_before = Pt(4)
            else:
                self._set_rich_text(p, s, size, self.FG)
                p.space_before = Pt(4)

        # Only enable autofit when caller left sizing to us
        if not explicit_height:
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        return tb

    def _rich_line(self, para, text, size=None, color=None, *, font=None,
                   bold=False):
        """Author prose into one paragraph, with the inline markup honoured.

        Captions, notes, definitions and quotes are prose: they carry
        `**bold**`, `` `code` `` and `$math$` like any other line, and setting
        `para.text` directly puts the markup on the slide verbatim. Styles are
        applied to the runs afterwards because `_set_rich_text` writes explicit
        run properties that paragraph-level defaults no longer reach.
        """
        self._set_rich_text(para, text, size, color)
        if font or bold:
            for r in para.runs:
                if font:
                    r.font.name = font
                if bold:
                    r.font.bold = True
        return para

    def _add_plain_run(self, para, text, size, color, bold=False, mono=False):
        """Append a single styled run to para. No-op if text is empty."""
        if not text:
            return
        run = para.add_run()
        run.text = text
        run.font.name = self.FONT_MONO if mono else self.FONT
        run.font.size = self._fs(size) if hasattr(size, "pt") else size
        if mono:
            # Make inline code visually distinct from prose.
            run.font.color.rgb = self.SECONDARY
        elif color is not None:
            run.font.color.rgb = color
        if bold:
            run.font.bold = True

    # Combined inline markup: **bold**, `code`, $math$, ==marker==
    _RICH_PATTERN = re.compile(
        r"(\*\*[^\*\n]+?\*\*)"
        r"|(`[^`\n]+?`)"
        r"|(\$[^\$\n]+?\$)"
        r"|(==[^=\n]+?==)"
    )

    def _run_highlight(self, run, rgb):
        """Marker-pen highlight behind a run (a:highlight — no python-pptx API).

        Schema order: highlight must precede a:latin in rPr, so insert
        before it when font.name has already written one.
        """
        rPr = run._r.get_or_add_rPr()
        hl = rPr.makeelement(qn("a:highlight"), {})
        clr = etree.SubElement(hl, qn("a:srgbClr"))
        clr.set("val", str(rgb))
        latin = rPr.find(qn("a:latin"))
        if latin is not None:
            latin.addprevious(hl)
        else:
            rPr.append(hl)

    @property
    def MARKER(self):
        """Highlight fill for ==marker== — the deck's own accent, tinted pale
        so black text stays AA on top (a foreign yellow reads as a sticker)."""
        return getattr(self.theme, "marker", None) or self._tint(self.ACCENT, 0.72)

    def _set_rich_text(self, para, text, size=None, color=None):
        """Render inline markup in a SINGLE paragraph / SINGLE textbox.

        Handles **bold**, `code` (monospace), and $math$ (OMML) without
        breaking the containing textbox. Runs co-exist with Japanese+Latin
        mixed text via the ea-font patch applied at save time.
        """
        if size is None:
            size = SZ_BODY
        if color is None:
            color = self.FG
        para.clear()
        if not text:
            return

        pos = 0
        for m in self._RICH_PATTERN.finditer(text):
            if m.start() > pos:
                self._add_plain_run(para, text[pos:m.start()], size, color)
            if m.group(1):  # **bold**
                self._add_plain_run(para, m.group(1)[2:-2], size, color, bold=True)
            elif m.group(2):  # `code`
                self._add_plain_run(para, m.group(2)[1:-1], size, color, mono=True)
            elif m.group(3):  # $math$
                latex = m.group(3)[1:-1]
                if not self._append_math_omml_inline(para, latex, size, color):
                    # OMML disabled (png preview mode): show the symbol without
                    # the $ delimiters, italicized, as a readable approximation.
                    run = para.add_run()
                    run.text = self._math_text_fallback(latex)
                    run.font.name = self.FONT
                    run.font.size = self._fs(size) if hasattr(size, "pt") else size
                    run.font.italic = True
                    if color is not None:
                        run.font.color.rgb = color
            elif m.group(4):  # ==marker== highlight
                self._add_plain_run(para, m.group(4)[2:-2], size, color)
                self._run_highlight(para.runs[-1], self.MARKER)
            pos = m.end()
        if pos < len(text):
            self._add_plain_run(para, text[pos:], size, color)

    # Backward-compatible alias for callers that used the math-only name
    def _set_text_with_inline_math(self, para, text, size, color):
        return self._set_rich_text(para, text, size=size, color=color)

    def _image_or_placeholder(self, img_path: str) -> str | None:
        """Resolved image path, or a generated placeholder PNG when missing.

        A missing figure used to leave the slide near-empty (warn + skip).
        Every image slot now stays visible — a themed frame with an image
        glyph and the missing filename — so a draft deck still reads as a
        deck and the author can see exactly which asset to fix.
        Returns None only when no image was asked for at all.
        """
        if not img_path:
            return None
        resolved = self._resolve_image(img_path)
        if resolved:
            return resolved
        try:
            return self._placeholder_png(Path(img_path).name)
        except Exception:
            return None   # placeholder is best-effort; the warn already fired

    def _placeholder_png(self, label: str) -> str:
        """Render (and cache) a 16:10 placeholder PNG in theme colors."""
        cache = getattr(self, "_ph_cache", None)
        if cache is None:
            import tempfile as _tf
            self._ph_dir = _tf.TemporaryDirectory(prefix="marp_ph_")
            cache = self._ph_cache = {}
        if label in cache:
            return cache[label]
        from PIL import Image as PImage, ImageDraw, ImageFont
        W, H = 1440, 900
        light, muted, hair, accent = (tuple(self.LIGHT), tuple(self.MUTED),
                                      tuple(self.HAIRLINE), tuple(self.ACCENT))
        im = PImage.new("RGB", (W, H), light)
        d = ImageDraw.Draw(im)
        d.rectangle([1, 1, W - 2, H - 2], outline=hair, width=3)
        # image glyph: frame + sun + mountains
        gx, gy, gw, gh = (W - 340) // 2, 190, 340, 240
        d.rounded_rectangle([gx, gy, gx + gw, gy + gh], radius=18,
                            outline=muted, width=6)
        d.ellipse([gx + 58, gy + 48, gx + 106, gy + 96], fill=accent)
        d.polygon([(gx + 24, gy + gh - 24), (gx + 132, gy + 82),
                   (gx + 208, gy + gh - 24)], fill=muted)
        d.polygon([(gx + 150, gy + gh - 24), (gx + 238, gy + 118),
                   (gx + gw - 24, gy + gh - 24)], fill=muted)
        try:
            fp = metrics.font_file(self.FONT_EA) or metrics.font_file(self.FONT)
            f_big = ImageFont.truetype(fp, 46)
            f_small = ImageFont.truetype(fp, 32)
        except Exception:
            f_big = f_small = ImageFont.load_default()
        d.text((W // 2, gy + gh + 96), label, fill=muted, font=f_big, anchor="mm")
        d.text((W // 2, gy + gh + 162), "image not found", fill=muted,
               font=f_small, anchor="mm")
        out = str(Path(self._ph_dir.name) / f"ph_{len(cache)}.png")
        im.save(out)
        cache[label] = out
        return out

    def _resolve_image(self, img_path: str) -> str | None:
        p = self.base_path / img_path
        if not p.exists():
            self._warn(f"image not found, placeholder drawn: {img_path}")
            return None
        if p.suffix.lower() == ".svg":
            png_path = p.with_suffix(".png")
            if not png_path.exists() or png_path.stat().st_mtime < p.stat().st_mtime:
                if HAS_CAIROSVG:
                    cairosvg.svg2png(url=str(p), write_to=str(png_path), output_width=1400, dpi=300)
                else:
                    self._warn(f"SVG needs cairosvg to embed, skipping: {img_path} "
                               "(install marp-pptx[ingest] or pre-render to PNG)")
                    return None
            return str(png_path)
        return str(p)

    def _fill_multiline_box(self, tf, text, size, color):
        """Fill a textbox with text that may contain bullets and continuation lines.

        Handles:
        - `- item` and `* item` as bullets
        - Continuation lines indented MORE than their parent bullet
        - **bold** markdown via _set_rich_text
        """
        # Strip common leading whitespace (dedent) so HTML-source indentation
        # doesn't get mistaken for content continuation.
        raw = [line for line in text.split("\n") if line.strip()]
        if not raw:
            return
        indents = [len(l) - len(l.lstrip()) for l in raw]
        base_indent = min(indents) if indents else 0
        dedented = [l[base_indent:] if len(l) >= base_indent else l for l in raw]

        # Merge continuation: a line is a continuation only if the previous
        # line was a bullet AND this line is further indented.
        merged: list[str] = []
        last_was_bullet = False
        for line in dedented:
            stripped = line.lstrip()
            line_indent = len(line) - len(stripped)
            is_bullet = stripped.startswith("- ") or stripped.startswith("* ")
            if last_was_bullet and line_indent > 0 and not is_bullet and merged:
                merged[-1] = merged[-1] + " " + stripped
            else:
                merged.append(stripped)
                last_was_bullet = is_bullet

        first = True
        for line in merged:
            s = line.strip()
            if not s:
                continue
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            is_bullet = s.startswith("- ") or s.startswith("* ")
            if is_bullet:
                self._set_rich_text(p, s[2:], size, color)
                p.space_before = Pt(4)
                pPr = p._p.get_or_add_pPr()
                buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
                for existing in pPr.findall(qn("a:buChar")):
                    pPr.remove(existing)
                pPr.append(buChar)
            else:
                self._set_rich_text(p, s, size, color)
                p.space_before = Pt(4)

    def _add_accent_box(self, slide, text, left, top, width, height, border_color=None):
        if border_color is None:
            border_color = self.ACCENT
        left, top, width, height = int(left), int(top), int(width), int(height)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        bg.adjustments[0] = CARD_RADIUS
        bg.fill.solid(); bg.fill.fore_color.rgb = self.SURFACE
        bg.line.color.rgb = self.HAIRLINE; bg.line.width = HAIRLINE_W
        self._no_shadow(bg)
        bdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, int(Pt(4)), height)
        bdr.fill.solid(); bdr.fill.fore_color.rgb = border_color
        bdr.line.fill.background(); self._no_shadow(bdr)
        tb = self._add_textbox(slide, left + int(Pt(18)), top + int(Pt(8)),
                               width - int(Pt(34)), height - int(Pt(16)))
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._fill_multiline_box(tf, text, SZ_ZONE_B, self.FG)
        return tb

    def _add_conclusion_box(self, slide, text, left, top, width, height):
        left, top, width, height = int(left), int(top), int(width), int(height)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        bg.adjustments[0] = CARD_RADIUS
        bg.fill.solid(); bg.fill.fore_color.rgb = self.SURFACE
        bg.line.color.rgb = self.HAIRLINE; bg.line.width = HAIRLINE_W
        self._no_shadow(bg)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, int(CARD_ACCENT_H))
        bar.fill.solid(); bar.fill.fore_color.rgb = self.ACCENT
        bar.line.fill.background(); self._no_shadow(bar)
        tb = self._add_textbox(slide, left + int(Pt(18)), top + int(Pt(10)),
                               width - int(Pt(36)), height - int(Pt(20)))
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._fill_multiline_box(tf, text, SZ_ZONE_B, self.FG)
        return tb

    def _add_footnote(self, slide, text):
        left = int(MARGIN_L)
        top = int(SH - Inches(0.62))
        if self.LAYOUT.footer_bar:
            top -= int(Inches(0.30))      # sit above the beamer bar
        # Stop short of the page number in the same band — a long source line
        # used to run underneath it.
        width = int(CONTENT_W - Inches(1.2))
        self._hairline(slide, left, top, Inches(2.8), color=self.HAIRLINE)
        tb = self._add_textbox(slide, left, top + int(Pt(5)), width, int(Inches(0.4)))
        tf = tb.text_frame
        tf.word_wrap = True
        # Footnotes cite formulas — "加速版は手順 4 の重み付けを $t_k^2$ に" —
        # and used to print the dollar signs.
        self._set_rich_text(tf.paragraphs[0], text, SZ_FOOT, self.MUTED)
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    def _add_zone_box(self, slide, left, top, width, height,
                      label="", body="", fill_color=None, fill=None,
                      label_size=None, body_size=None,
                      accent_bar=False, accent=None, anchor="middle"):
        """Backward-compatible wrapper over the refined-minimal card.

        Every zone-based builder (zone-flow/compare/matrix/process, steps,
        stack, card-grid, split-text, before-after, multi-result) routes
        through here, so they all gain the visible surface + hairline at once.
        anchor defaults to "middle": a short label+body pinned to the top of a
        tall card leaves the bottom two-thirds empty, which is the single most
        common "looks unfinished" complaint in real decks.
        """
        return self._add_card(
            slide, (left, top, width, height),
            label=label, body=body, fill=fill if fill is not None else fill_color,
            label_size=label_size, body_size=body_size,
            accent_bar=accent_bar, accent=accent, anchor=anchor,
        )

    def _set_cell_border(self, cell, edges=("bottom",), color=None, pt=0.75):
        """Inject cell borders via lxml (python-pptx has no public API).

        `pt=None` writes an explicit *no* border. That is not the same as
        leaving the edge alone: a table with no style id still inherits a full
        grid from the renderer's default, which is why every table used to come
        out looking like a spreadsheet.
        """
        if color is None:
            color = self.HAIRLINE
        hexcol = str(color)
        tcPr = cell._tc.get_or_add_tcPr()
        tag = {"left": "a:lnL", "right": "a:lnR", "top": "a:lnT", "bottom": "a:lnB"}
        for edge in edges:
            qname = qn(tag[edge])
            for old in tcPr.findall(qname):
                tcPr.remove(old)
            if pt is None:
                ln = tcPr.makeelement(qname, {})
                ln.append(ln.makeelement(qn("a:noFill"), {}))
                tcPr.insert(0, ln)
                continue
            ln = tcPr.makeelement(qname, {"w": str(int(Pt(pt))), "cap": "flat",
                                          "cmpd": "sng", "algn": "ctr"})
            fill = ln.makeelement(qn("a:solidFill"), {})
            clr = fill.makeelement(qn("a:srgbClr"), {"val": hexcol})
            fill.append(clr); ln.append(fill)
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "solid"}))
            tcPr.insert(0, ln)

    # Table geometry — column widths and row heights come from the content,
    # measured with the same font metrics as everything else.
    TBL_PAD_X = Pt(10)
    TBL_PAD_Y = Pt(6)

    _NUM_RE = re.compile(r"^[\s$¥€£]*[-+]?[\d,]+(\.\d+)?\s*[%x×kKMB]?\s*$")

    def _table_numeric_cols(self, rows_data) -> set:
        """Columns whose body is mostly numbers — right-aligned, and never
        the column that absorbs slack width."""
        cols = max((len(r) for r in rows_data), default=1)
        out = set()
        for ci in range(cols):
            body = [self._plain(strip_html(r[ci])) for r in rows_data[1:]
                    if ci < len(r) and r[ci].strip()]
            if body and sum(bool(self._NUM_RE.match(b)) for b in body) >= max(
                    1, len(body) * 0.6):
                out.add(ci)
        return out

    def _table_cols(self, rows_data, width, size=None):
        """Column widths (EMU) sized to what each column actually holds.

        Equal columns waste half the slide on a two-character "12" and squeeze
        the column carrying prose. Slack goes to the text columns rather than
        being spread over every column, which would strand a right-aligned "12"
        an inch away from its label.
        """
        cols = max((len(r) for r in rows_data), default=1)
        if cols <= 0:
            return []
        size_pt = (size or self._fs(SZ_SMALL)).pt
        pad = int(self.TBL_PAD_X) * 2
        floor = int(Inches(0.7))
        ceil = int(width * 0.5)
        natural = []
        for ci in range(cols):
            w = 0.0
            for r in rows_data:
                cell = self._plain(strip_html(r[ci])) if ci < len(r) else ""
                w = max(w, metrics.measure_pt(cell, self.FONT, size_pt,
                                              ea_font=self.FONT_EA))
            natural.append(min(max(int(Pt(w)) + pad, floor), ceil))
        total = sum(natural)
        # A table narrower than ~3/4 of the content width reads as a stray
        # fragment; a wider one has to be squeezed back in.
        target = min(int(width), max(total, int(width * 0.78)))
        widths = list(natural)
        if total > target:
            widths = [max(floor, int(n * target / total)) for n in natural]
        elif total < target:
            grow = sorted(set(range(cols)) - self._table_numeric_cols(rows_data)) \
                or list(range(cols))
            share = (target - total) // len(grow)
            for ci in grow:
                widths[ci] += share
        drift = target - sum(widths)
        widths[widths.index(max(widths))] += drift
        return widths

    def _table_rows(self, rows_data, widths, size=None):
        """Row heights (EMU) from the wrapped content of each row."""
        size_pt = (size or self._fs(SZ_SMALL)).pt
        line_h = size_pt * metrics.DEFAULT_LINE_FACTOR
        pad = int(self.TBL_PAD_Y) * 2
        out = []
        for r in rows_data:
            lines = 1
            for ci, w in enumerate(widths):
                cell = self._plain(strip_html(r[ci])) if ci < len(r) else ""
                if not cell:
                    continue
                inner = (w - int(self.TBL_PAD_X) * 2) / 12700.0
                lines = max(lines, metrics.line_count(cell, self.FONT, size_pt,
                                                      max(20.0, inner),
                                                      ea_font=self.FONT_EA))
            out.append(int(Pt(lines * line_h)) + pad)
        return out

    def _table_height(self, rows_data, width, size=None) -> int:
        """Height the table will occupy — what callers must reserve for it."""
        if not rows_data:
            return 0
        widths = self._table_cols(rows_data, width, size)
        return sum(self._table_rows(rows_data, widths, size))

    def _styled_table(self, slide, rows_data, left, top, width, height=None):
        """Themed table: measured columns and rows, horizontal rules only.

        `height` is ignored except as a floor — a table's height is the sum of
        its rows, and forcing a different one just means the renderer grows the
        rows past the frame and clips the last rule.
        Header style follows theme.layout.table_header_style (fill | rule).
        """
        rows = len(rows_data)
        cols = max(len(r) for r in rows_data) if rows_data else 1
        col_w = self._table_cols(rows_data, width)
        row_h = self._table_rows(rows_data, col_w)
        # Centred in the block it was given, since it may be narrower than it.
        x = int(left) + max(0, (int(width) - sum(col_w)) // 2)
        gf = slide.shapes.add_table(rows, cols, x, int(top),
                                    sum(col_w), sum(row_h))
        table = gf.table
        for ci, w in enumerate(col_w):
            table.columns[ci].width = Emu(int(w))
        for ri, h in enumerate(row_h):
            table.rows[ri].height = Emu(int(h))
        try:
            table.first_row = False
            table.horz_banding = False
            tblPr = table._tbl.find(qn("a:tblPr"))
            if tblPr is not None:
                for sid in tblPr.findall(qn("a:tableStyleId")):
                    tblPr.remove(sid)
        except Exception:
            pass
        rule_header = getattr(self.LAYOUT, "table_header_style", "fill") == "rule"
        # Numbers read better right-aligned (CSE table guidance).
        numeric_cols = self._table_numeric_cols(rows_data)
        last = rows - 1
        for ri, row in enumerate(rows_data):
            for ci in range(cols):
                cell = table.cell(ri, ci)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = cell.margin_right = self.TBL_PAD_X
                cell.margin_top = cell.margin_bottom = self.TBL_PAD_Y
                raw = row[ci] if ci < len(row) else ""
                is_head = (ri == 0)
                emph = ("**" in raw)   # markdown-bolded cell (e.g. **Ours**)
                if is_head:
                    color = self.FG if rule_header else self.WHITE
                elif emph:
                    color = self.ACCENT_TEXT
                elif ci == 0:
                    color = self.SECONDARY
                else:
                    color = self.FG
                bold = is_head or emph or (ci == 0 and not is_head)
                # Cells are author text: `**Ours**` and `$\sqrt{n}$` have to be
                # rendered, not printed.
                p = cell.text_frame.paragraphs[0]
                self._rich_line(p, strip_html(raw), SZ_SMALL, color,
                                font=self.FONT_HEAD if (is_head or emph) else None,
                                bold=bold)
                p.alignment = PP_ALIGN.RIGHT if ci in numeric_cols else PP_ALIGN.LEFT

                # Rules run horizontally only: a full grid reads as a
                # spreadsheet, and the vertical lines carry no information the
                # column alignment doesn't already give.
                self._set_cell_border(cell, ("left", "right"), pt=None)
                if is_head and not rule_header:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.PRIMARY
                else:
                    cell.fill.background()
                if is_head:
                    self._set_cell_border(cell, ("top",), color=self.PRIMARY, pt=1.25)
                    self._set_cell_border(cell, ("bottom",), color=self.PRIMARY, pt=1.0)
                elif ri == last:
                    self._set_cell_border(cell, ("top",), pt=None)
                    self._set_cell_border(cell, ("bottom",), color=self.PRIMARY, pt=1.25)
                else:
                    self._set_cell_border(cell, ("top",), pt=None)
                    self._set_cell_border(cell, ("bottom",), color=self.HAIRLINE, pt=0.5)
        return gf

    def _add_column_content(self, slide, lines, left, top, width, height, size=None):
        if size is None:
            size = SZ_COL
        text_lines = []
        images = []
        for line in lines:
            img_m = re.match(r"!\[(?:w:\d+)?\]\(([^)]+)\)", line.strip())
            if img_m:
                images.append(img_m.group(1))
            else:
                text_lines.append(line)
        cur_top = top
        for img_path in images:
            img_file = self._image_or_placeholder(img_path)
            if img_file:
                from PIL import Image
                with Image.open(img_file) as im:
                    iw, ih = im.size
                max_w = int(width * 0.95)
                max_h = int(Inches(2.5))
                scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
                pw = int(iw * scale * 914400 / 96)
                ph = int(ih * scale * 914400 / 96)
                img_left = left + (width - pw) // 2
                slide.shapes.add_picture(img_file, img_left, cur_top, pw, ph)
                cur_top += ph + Inches(0.1)
        # Render text and box segments in author order. Boxes arrive as
        # \x00-sentinel runs from parser.column_lines — before this, a
        # <div class="box-accent"> in a column silently rendered as plain
        # text (the styling the skeletons advertise never appeared).
        for kind, seg in self._column_segments(text_lines):
            seg_lines = [l for l in seg if l.strip()]
            if not seg_lines:
                continue
            remaining_h = top + height - cur_top
            if remaining_h <= 0:
                break
            if kind == "text":
                # Hug content height (capped at remaining), don't force full
                # column. Measure at the column's own width — without it this
                # re-estimate threw away the wrap-aware height the caller had
                # just computed and gave every wrapped bullet a single line.
                estimated = self._estimate_text_height(seg_lines, size, width=width)
                use_h = min(int(remaining_h), int(estimated))
                tb = self._add_body_text(slide, seg_lines, left=left, top=int(cur_top),
                                         width=int(width), height=use_h, size=size)
                # Lock the column textbox size — some viewers (Keynote,
                # LibreOffice) ignore spAutoFit and render the stored height.
                if tb is not None:
                    tb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                cur_top += use_h + int(Inches(0.12))
            else:
                inner_w = int(width - Pt(34))
                est = self._estimate_text_height(seg_lines, SZ_ZONE_B, width=inner_w)
                box_h = min(int(remaining_h), int(est + Pt(16)))
                border = {"accent": self.ACCENT, "primary": self.SECONDARY}.get(kind)
                if border is not None:
                    self._add_accent_box(slide, "\n".join(seg_lines), left, cur_top,
                                         width, box_h, border_color=border)
                else:
                    self._add_card(slide, (left, int(cur_top), int(width), box_h),
                                   body_lines=seg_lines, anchor="middle")
                cur_top += box_h + int(Inches(0.12))

    @staticmethod
    def _column_segments(lines):
        """Split column lines into ("text" | "plain" | "accent" | "primary", lines)."""
        segs, kind, cur = [], "text", []
        for line in lines:
            s = line.strip()
            if s.startswith("\x00BOX"):
                if cur:
                    segs.append((kind, cur))
                parts = s.split()
                kind, cur = (parts[1] if len(parts) > 1 else "plain"), []
            elif s == "\x00END":
                segs.append((kind, cur))
                kind, cur = "text", []
            else:
                cur.append(line)
        if cur:
            segs.append((kind, cur))
        return segs

    def _column_height(self, lines, size, width):
        """Height a column needs — box chrome included (mirror of the renderer)."""
        total, first = 0, True
        for kind, seg in self._column_segments(lines):
            seg_lines = [l for l in seg if l.strip()]
            if not seg_lines:
                continue
            if not first:
                total += int(Inches(0.12))
            first = False
            if kind == "text":
                total += int(self._estimate_text_height(seg_lines, size, width=width))
            else:
                total += int(self._estimate_text_height(
                    seg_lines, SZ_ZONE_B, width=int(width - Pt(34))) + Pt(16))
        return total

    # ══════════════════════════════════════════════
    # Slide type builders
    # ══════════════════════════════════════════════

    def _hero_meta(self, sd):
        """Split a title slide's body into (subtitle_lines, meta_tokens).
        A line with a 4-digit year or an author/date marker becomes meta."""
        subs, meta = [], []
        past_h1 = False
        for line in sd.raw.split("\n"):
            s = line.strip()
            if s.startswith("# ") and not past_h1:
                past_h1 = True
                continue
            if s.startswith("## "):
                continue   # H2 is used as the kicker
            if past_h1 and s and not s.startswith(("<!--", "<div", "</div", ">")):
                t = strip_html(s)
                if not t:
                    continue
                if re.search(r"\b(19|20)\d{2}\b", t) or t.startswith(("by ", "@", "—", "·")):
                    meta.append(t.lstrip("—@·· ").strip())
                else:
                    subs.append(t)
        return subs, meta

    def build_title(self, sd: SlideData):
        slide = self._blank_slide()
        hero = self._hero_fill_color()
        if self.LAYOUT.title_bg == "gradient":
            self._set_gradient_bg(slide, self.PRIMARY, self.SECONDARY)
        elif self.LAYOUT.title_bg == "dark":
            self._set_bg(slide, self.PRIMARY)
        elif self.LAYOUT.title_bg == "light":
            self._set_bg(slide, self.LIGHT)
        elif hero is not None:
            self._set_bg(slide, hero)
        is_dark = self.LAYOUT.title_bg in ("gradient", "dark")
        if hero is not None and self.LAYOUT.title_bg == "white":
            from marp_pptx.audit import contrast_ratio
            is_dark = contrast_ratio(tuple(hero), (255, 255, 255)) > 4.0
        is_box = self.LAYOUT.title_bg == "box"   # Madrid: navy hero box on white
        h_color = self.WHITE if (is_dark or is_box) else self.PRIMARY
        sub_color = self._tint(self.PRIMARY, 0.85) if is_dark else self.MUTED
        accent = self._hero_accent(2.5) if is_dark else self.ACCENT       # rule (graphic)
        kicker_color = self._hero_accent(4.5) if is_dark else self.ACCENT_TEXT  # small text
        if hero is not None and not is_dark:
            # A light hero fill eats a little contrast — re-measure the small text.
            sub_color = self._text_safe(sub_color, hero)
            kicker_color = self._text_safe(kicker_color, hero)
        cx = int(SW // 2)
        is_left = self.LAYOUT.title_align == "left"
        align = PP_ALIGN.LEFT if is_left else PP_ALIGN.CENTER
        tx = int(MARGIN_L) if is_left else None  # shared left edge for the stack

        subs, meta = self._hero_meta(sd)
        kicker = sd.h2 if (sd.h2 and len(sd.h2.split()) <= 6) else None

        # Vertically center the kicker→title→rule→subtitle→meta stack.
        # (1) kicker
        if kicker:
            kb = self._add_textbox(slide, tx or int(Inches(1.0)), int(Inches(2.55)),
                                   int(SW - (tx or int(Inches(1.0))) * 2), int(KICKER_H))
            self._kicker_para(kb.text_frame.paragraphs[0], kicker, color=kicker_color,
                              align=align, tracking=200)
        # (2) accent hairline above the title (skipped for the hero box)
        if not is_box:
            self._hairline(slide, tx if is_left else int(cx - Inches(0.55)), int(Inches(3.16)),
                           Inches(1.1), thickness=Pt(1.6), color=accent)
        # (3) title
        # Two lines at the size and spacing the title is actually written
        # with — reserving at 1.0 while rendering at LINE_TITLE clipped the
        # second line of any hero title that wrapped.
        hero_pt = SZ_DISPLAY.pt * getattr(self.LAYOUT, "title_scale", 1.0)
        title_h = int(self._fs(Pt(hero_pt * metrics.DEFAULT_LINE_FACTOR
                                  * LINE_TITLE * 2)))
        if is_box:
            # Madrid title page: rounded structure-colored box, white title.
            # Sized for two lines regardless — line-count estimates miss by a
            # hair on borderline CJK titles, and a roomy box centers fine.
            scale = getattr(self.theme, "font_scale", 1.0)
            box_h = int(Pt(hero_pt * metrics.DEFAULT_LINE_FACTOR
                           * LINE_TITLE * 2 * scale) + Inches(0.34))
            box_y = int(Inches(3.30) + (title_h - box_h) // 2)
            bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        int(Inches(0.95)), box_y,
                                        int(SW - Inches(1.9)), box_h)
            bx.adjustments[0] = 0.16
            bx.fill.solid(); bx.fill.fore_color.rgb = self.PRIMARY
            bx.line.fill.background(); self._no_shadow(bx)
            tb = self._add_textbox(slide, int(Inches(1.15)), box_y,
                                   int(SW - Inches(2.3)), box_h)
            tf = tb.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            tb = self._add_textbox(slide, tx or int(Inches(0.8)), int(Inches(3.34)),
                                   int(SW - (tx or int(Inches(0.8))) * 2), title_h)
            tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]; p.text = sd.h1
        p.font.name = self.FONT_HEAD; p.font.size = self._fs(Pt(hero_pt))
        p.font.bold = True; p.font.color.rgb = h_color; p.alignment = align
        p.line_spacing = LINE_TITLE
        # (4) subtitle — rich text so $math$ (e.g. author superscripts) renders
        if subs:
            # Follow the title block instead of a fixed y: at font_scale 1.3
            # a two-line hero title reached down into the subtitle.
            title_bottom = (box_y + box_h) if is_box else (int(Inches(3.34)) + title_h)
            sub_y = max(int(Inches(5.18) if is_box else Inches(5.05)),
                        title_bottom + int(Inches(0.25)))
            sb = self._add_textbox(slide, tx or int(Inches(1.2)), int(sub_y),
                                   int(SW - (tx or int(Inches(1.2))) * 2), int(Inches(0.95)))
            sb.text_frame.word_wrap = True
            for i, line in enumerate(subs[:3]):
                sp = sb.text_frame.paragraphs[0] if i == 0 else sb.text_frame.add_paragraph()
                self._set_rich_text(sp, line, SZ_BODY, sub_color)
                sp.alignment = align
                if i:
                    sp.space_before = Pt(4)
        # (5) author / affiliation / date
        if meta:
            meta_y = int(Inches(6.35))
            if subs:
                meta_y = max(meta_y, int(sub_y) + int(Inches(1.0)))
            mb = self._add_textbox(slide, tx or int(Inches(1.0)), meta_y,
                                   int(SW - (tx or int(Inches(1.0))) * 2), int(Inches(0.45)))
            mp = mb.text_frame.paragraphs[0]
            self._set_rich_text(mp, "   ·   ".join(meta), SZ_SMALL, sub_color)
            mp.alignment = align

    def build_divider(self, sd: SlideData):
        slide = self._blank_slide()
        hero = self._hero_fill_color()
        if hero is not None:
            self._set_bg(slide, hero)
        is_center = self.LAYOUT.divider_align != "left"
        align = PP_ALIGN.CENTER if is_center else PP_ALIGN.LEFT
        cx = int(SW // 2)
        # ghost section number
        m = re.match(r"^\s*0?(\d{1,2})\b", sd.h2 or "")
        ghost = (m.group(1) if m else str(getattr(self, "_divider_no", 1))).zfill(2)
        if self.LAYOUT.divider_number:
            gb = self._add_textbox(
                slide, int(Inches(0.85)), int(Inches(0.85)), int(Inches(5.0)),
                int(Pt(self._fs(Pt(150)).pt * metrics.DEFAULT_LINE_FACTOR)))
            gp = gb.text_frame.paragraphs[0]; gp.text = ghost
            gp.font.name = self.FONT_HEAD; gp.font.size = self._fs(Pt(150))
            gp.font.bold = True; gp.font.color.rgb = self._tint(self.PRIMARY, 0.88)
            gp.alignment = PP_ALIGN.LEFT
        x = int(Inches(1.0)) if is_center else int(Inches(1.4))
        w = int(SW - Inches(2.0)) if is_center else int(SW - Inches(2.8))
        # kicker
        kb = self._add_textbox(slide, x, int(Inches(3.5)), w, int(KICKER_H))
        self._kicker_para(kb.text_frame.paragraphs[0],
                          ("Section " + ghost) if self.LAYOUT.divider_number else "Section",
                          color=self.ACCENT_TEXT, align=align, tracking=200)
        # title — the rule and subtitle below it follow the title's measured
        # height, not fixed inches, or a larger font_scale walks the title into
        # the subtitle.
        title_pt = self._fs(Pt(34))
        title_y = int(Inches(3.92))
        title_lines = min(2, metrics.line_count(
            self._plain(sd.h1 or ""), self.FONT_HEAD, title_pt.pt, w / 12700.0,
            ea_font=self.FONT_EA, bold=True))
        title_h = int(Pt(title_pt.pt * LINE_TITLE * metrics.DEFAULT_LINE_FACTOR
                         * max(1, title_lines)))
        tb = self._add_textbox(slide, x, title_y, w, title_h)
        tf = tb.text_frame; tf.word_wrap = True
        tp = tf.paragraphs[0]; tp.text = sd.h1
        tp.font.name = self.FONT_HEAD; tp.font.size = title_pt
        tp.font.bold = True; tp.font.color.rgb = self.PRIMARY; tp.alignment = align
        tp.line_spacing = LINE_TITLE
        # accent rule
        rule_y = title_y + title_h + int(Inches(0.18))
        rx = cx if is_center else int(Inches(1.4) + Inches(0.35))
        self._hairline(slide, int(rx - Inches(0.35)), rule_y, Inches(0.7),
                       thickness=ACCENT_RULE_W, color=self.ACCENT)
        # subtitle (only when h2 wasn't just the number)
        if sd.h2 and not m:
            sb = self._add_textbox(slide, x, rule_y + int(Inches(0.23)), w,
                                   int(self._estimate_text_height([sd.h2], SZ_BODY,
                                                                  width=w)))
            sb.text_frame.word_wrap = True
            sp = sb.text_frame.paragraphs[0]; sp.text = sd.h2
            sp.font.name = self.FONT; sp.font.size = self._fs(SZ_BODY)
            sp.font.color.rgb = self.MUTED; sp.alignment = align

    def build_default(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        region = self._content_region_with_lead(slide, sd)
        left, _, width, _ = region

        # Measure each block, then center/justify the stack in the body region.
        blocks = []  # (kind, height)
        if sd.body_lines:
            bh = int(self._estimate_text_height(sd.body_lines, SZ_BODY, width=width))
            blocks.append(("body", min(bh, region[3])))
        if sd.table_rows:
            blocks.append(("table", self._table_height(sd.table_rows, width)))
        if sd.bottom_text:
            blocks.append(("accent", int(Inches(1.0))))

        heights = [h for _, h in blocks]
        mode = "center" if len(heights) <= 2 else "justify"
        tops = self._stack_tops(heights, region, mode=mode)
        for (kind, h), t in zip(blocks, tops):
            if kind == "body":
                self._add_body_text(slide, sd.body_lines, left=left, top=t,
                                    width=width, height=h)
            elif kind == "table":
                self._styled_table(slide, sd.table_rows, left, t, width, h)
            elif kind == "accent":
                self._add_accent_box(slide, sd.bottom_text, left, t, width, h)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_equation(self, sd: SlideData):
        if any(r[2] for r in sd.eq_annotations):
            return self.build_equation_annotated_lines(sd)
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        region = self._content_region(has_title=bool(sd.h1))
        _, rtop, _, rheight = region
        # A short identity (E=mc², a definition) at a fixed 34pt floats small
        # in 12in of canvas — let brevity buy size. The glyph proxy strips
        # commands/braces so \frac{a}{b} counts its visible glyphs only.
        glyphs = len(re.sub(r"\\[a-zA-Z]+|[{}^_\s]", "", sd.eq_main or ""))
        EQ_PT = 44 if glyphs <= 10 else 40 if glyphs <= 18 else 36 if glyphs <= 30 else 34
        # Estimate the equation block + variable legend, center the stack.
        eq_h_est = int(Emu(int(EQ_PT * 2.4 * 12700)))
        var_n = len(sd.eq_vars)
        row_h = int(Inches(0.58))
        vars_h = row_h * var_n
        gap = int(Inches(0.35))
        total = eq_h_est + (gap + vars_h if var_n else 0)
        eq_top = rtop + max(0, (rheight - total) // 2)
        omml_box = self._add_math_omml_display(slide, sd.eq_main, MARGIN_L, eq_top, CONTENT_W, pt_size=EQ_PT)
        if omml_box is not None:
            var_top = eq_top + omml_box.height + Inches(0.25)
        else:
            result = self._add_math_image(slide, sd.eq_main, MARGIN_L, eq_top, CONTENT_W, display=True, fontsize=EQ_PT)
            if result:
                _, eq_h = result
                var_top = eq_top + eq_h + Inches(0.3)
            else:
                tb = self._add_textbox(slide, MARGIN_L, eq_top, CONTENT_W, Inches(1.2))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sd.eq_main
                p.font.name = self.FONT
                p.font.size = self._fs(Pt(28))
                p.font.color.rgb = self.FG
                p.alignment = PP_ALIGN.CENTER
                var_top = eq_top + Inches(1.6)
        if sd.eq_vars:
            self._add_var_legend(slide, sd.eq_vars, var_top, sym_pt=20, desc_pt=17)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_equation_annotated_lines(self, sd: SlideData):
        """tmu-cs [!math-annotate]: the equation stays ONE line — segments are
        measured and composed on a shared baseline, and note cards hang below
        with a connector pointing at their own term (HTML sample #12 parity).

        Pointing accuracy needs build-time term coordinates, so this mode
        renders the math as measured PNG segments rather than OMML."""
        from marp_pptx.math.renderer import render_latex_png_measured
        from PIL import Image
        rows = sd.eq_annotations
        EQ_PT = 30
        segs = []
        for (tex, label, note, color) in rows:
            r = render_latex_png_measured(tex, fontsize=EQ_PT,
                                          color=f"#{self.FG}", dpi=self._MATH_DPI)
            if r is None:
                segs = None
                break
            png, depth = r
            with Image.open(png) as im:
                iw, ih = im.size
            segs.append({"png": png, "w": iw, "h": ih, "depth": depth,
                         "label": label, "note": note, "color": color})
        if segs is None:
            # mathtext failed on a segment — degrade to the plain equation
            # with the notes as a variable legend so no content is dropped.
            fb = dc_replace(sd, eq_annotations=[],
                            eq_vars=sd.eq_vars or [(r[0], r[2]) for r in rows if r[2]])
            self._warn(f"slide {sd.index + 1}: math-annotate segment failed; "
                       "rendered as plain equation + legend")
            return self.build_equation(fb)

        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))

        emu_per_px = 914400 / self._MATH_DPI
        gap_px = self._MATH_DPI * 0.05
        total_px = sum(s["w"] for s in segs) + gap_px * (len(segs) - 1)
        fit = min(1.0, (rwidth * 0.96) / (total_px * emu_per_px))
        asc = max(s["h"] - s["depth"] for s in segs)   # px above the baseline
        dep = max(s["depth"] for s in segs)            # px below the baseline
        line_h = int((asc + dep) * emu_per_px * fit)

        annotated = [s for s in segs if s["note"]]
        note_pt = 12
        card_gap = int(Inches(0.25))
        n = len(annotated)
        # Predict each segment's center (same arithmetic as placement below)
        # so the card width can shrink to the pitch between annotated terms —
        # wide cards shove each other sideways and the connectors degrade
        # into long diagonals that no longer read as "pointing".
        run = 0.0
        for s in segs:
            s["cx_pred"] = run + s["w"] / 2.0
            run += s["w"] + gap_px
        cap = int(Inches(3.4))
        if n >= 2:
            cs = sorted(s["cx_pred"] * emu_per_px * fit for s in annotated)
            min_pitch = min(b - a for a, b in zip(cs, cs[1:]))
            cap = min(cap, max(int(Inches(1.5)), int(min_pitch) - card_gap))
        card_w = int(min(cap, (rwidth - card_gap * (n - 1)) / max(n, 1)))
        # Card height from the note measured at the card's real inner width,
        # plus the card's own padding. Narrowing the cards to the term pitch
        # (above) makes the notes wrap more, so this has to be measured, not
        # assumed: a fixed 0.55in card clipped every note over two lines.
        inner_w = card_w - int(CARD_PAD) * 2
        card_h = 0
        for s in annotated:
            body_h = int(self._estimate_text_height([s["note"]], Pt(note_pt),
                                                    width=inner_w))
            label_h = int(Pt(10 * metrics.DEFAULT_LINE_FACTOR)) if s["label"] else 0
            s_h = body_h + label_h + int(CARD_PAD) * 2
            card_h = max(card_h, max(int(Inches(0.55)), s_h))
        drop = int(Inches(0.55))                       # connector run below the row
        block_h = line_h + ((drop + card_h) if n else 0)
        eq_top = rtop + max(0, (rheight - block_h) // 2)

        # Equation row, baseline-aligned.
        x = rleft + max(0, (rwidth - int(total_px * emu_per_px * fit)) // 2)
        for s in segs:
            y = eq_top + int((asc - (s["h"] - s["depth"])) * emu_per_px * fit)
            pw = int(s["w"] * emu_per_px * fit)
            ph = int(s["h"] * emu_per_px * fit)
            slide.shapes.add_picture(s["png"], int(x), int(y), pw, ph)
            s["cx"] = int(x + pw / 2)
            s["bottom"] = int(y + ph)
            x += pw + int(gap_px * emu_per_px * fit)

        # Note cards: centered under their own term where possible, nudged
        # apart when neighbours collide, even spread as the last resort.
        xs = []
        cur = rleft
        for s in annotated:
            xi = max(int(s["cx"] - card_w // 2), int(cur))
            xs.append(xi)
            cur = xi + card_w + card_gap
        if xs:
            over = xs[-1] + card_w - (rleft + rwidth)
            if over > 0:
                xs = [xi - over for xi in xs]
            if xs[0] < rleft:
                spread = card_w * n + card_gap * (n - 1)
                x0 = rleft + max(0, (rwidth - spread) // 2)
                xs = [x0 + i * (card_w + card_gap) for i in range(n)]
        card_top = int(eq_top + line_h + drop)
        for s, card_x in zip(annotated, xs):
            acc = self.ACCENT
            if s["color"]:
                hx = s["color"].lstrip("#")
                try:
                    acc = RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
                except (ValueError, IndexError):
                    pass
            _, tb = self._add_card(slide, (card_x, card_top, card_w, card_h),
                                   anchor="middle")
            tf = tb.text_frame
            p = tf.paragraphs[0]
            if s["label"]:
                self._kicker_para(p, s["label"], size=Pt(10),
                                  color=self._text_safe(acc), tracking=120)
                p2 = tf.add_paragraph()
            else:
                p2 = p
            self._rich_line(p2, s["note"], Pt(note_pt), self.FG)
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                s["cx"], int(s["bottom"] + Inches(0.03)),
                int(card_x + card_w // 2), int(card_top - Inches(0.02)))
            conn.line.color.rgb = acc
            conn.line.width = Pt(1.1)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def _add_var_legend(self, slide, eq_vars, top, sym_pt=18, desc_pt=15,
                        row_h=None):
        """Refined "where: <sym> | <description>" legend, centered as a group:
        right-aligned accent symbol, thin vertical hairline, left description."""
        if row_h is None:
            row_h = int(Inches(0.5))
        sym_w = int(Inches(1.7))
        desc_w = int(Inches(6.2))
        group_w = sym_w + int(Inches(0.5)) + desc_w
        gx = int(MARGIN_L + (CONTENT_W - group_w) // 2)
        # "where" eyebrow
        self._eyebrow(slide, "where", gx, int(top - KICKER_H + Inches(0.04)),
                      sym_w, align=PP_ALIGN.RIGHT, color=self.MUTED)
        n = len(eq_vars)
        # vertical hairline divider spanning the rows
        div_x = gx + sym_w + int(Inches(0.25))
        self._hairline(slide, div_x, int(top + Inches(0.04)),
                       int(row_h * n - Inches(0.08)), thickness=HAIRLINE_W,
                       color=self.HAIRLINE, vertical=True)
        for vi, (sym, desc) in enumerate(eq_vars):
            row_top = int(top + row_h * vi)
            sym_latex = sym.strip().strip("$")
            stb = self._add_textbox(slide, gx, row_top, sym_w, row_h)
            stb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            sp = stb.text_frame.paragraphs[0]
            sp.alignment = PP_ALIGN.RIGHT
            if not self._append_math_omml_inline(sp, sym_latex, Pt(sym_pt), self.ACCENT_TEXT):
                sp.text = self._math_text_fallback(sym_latex)
                sp.font.name = self.FONT; sp.font.size = self._fs(Pt(sym_pt - 2))
                sp.font.bold = True; sp.font.italic = True
                sp.font.color.rgb = self.ACCENT_TEXT
            dtb = self._add_textbox(slide, div_x + int(Inches(0.25)), row_top, desc_w, row_h)
            dtb.text_frame.word_wrap = True
            dtb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Legend rows are a fixed grid — a long description gives way on
            # size rather than spilling into the row below it.
            self._set_text_with_inline_math(
                dtb.text_frame.paragraphs[0], desc,
                self._fit_size(desc, Pt(desc_pt), desc_w, row_h), self.FG)

    def build_equations(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.eq_system)
        if n == 0:
            return
        region = self._content_region(has_title=bool(sd.h1))
        _, rtop, _, rheight = region
        label_left = MARGIN_L + Inches(0.3)
        label_w = Inches(1.9)
        eq_left = label_left + label_w + Inches(0.2)
        eq_w = CONTENT_W - (eq_left - MARGIN_L) - Inches(0.3)
        # Reserve the legend at the row height it is actually drawn with. The
        # two used to disagree (0.56in reserved, 0.46in drawn, plus a 0.25/0.32
        # gap mismatch), which walked the last legend row into the footnote.
        legend_row = int(Inches(0.46))
        legend_gap = int(Inches(0.32))
        vars_h = (legend_row * len(sd.eq_vars) + legend_gap) if sd.eq_vars else 0
        row_h = int(min(Inches(1.2), max(Inches(0.8), (rheight - vars_h) // n)))
        pt_size = 30 if n <= 3 else (26 if n <= 4 else 22)
        # Center the whole rows+legend block in the body region.
        block_h = row_h * n + vars_h
        if block_h > rheight:
            # More rows than the region can hold at full rhythm: tighten both
            # rhythms proportionally instead of letting the last legend row
            # walk off the bottom into the footnote.
            k = rheight / block_h
            row_h = max(int(Inches(0.55)), int(row_h * k))
            legend_row = max(int(Inches(0.34)), int(legend_row * k))
            vars_h = (legend_row * len(sd.eq_vars) + legend_gap) if sd.eq_vars else 0
            block_h = row_h * n + vars_h
        top = rtop + max(0, (rheight - block_h) // 2)
        for i, (label, latex) in enumerate(sd.eq_system):
            row_top = top + int(row_h * i)
            if label:
                ltb = self._add_textbox(slide, label_left, row_top, label_w, row_h)
                ltf = ltb.text_frame
                ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
                ltf.word_wrap = True
                lp = ltf.paragraphs[0]
                lp.text = label
                lp.font.name = self.FONT
                lp.font.size = self._fs(Pt(max(14, pt_size - 10)))
                lp.font.color.rgb = self.SECONDARY
                lp.alignment = PP_ALIGN.RIGHT
            el = self._omml_element(latex, display=True, pt_size=pt_size)
            etb = self._add_textbox(slide, eq_left, row_top, eq_w, row_h)
            etf = etb.text_frame
            etf.word_wrap = True
            etf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ep = etf.paragraphs[0]
            ep.alignment = PP_ALIGN.LEFT
            erun = ep.add_run()
            erun.text = ""
            erun.font.name = self.FONT
            erun.font.size = self._fs(Pt(pt_size))
            erun.font.color.rgb = self.FG
            if el is not None:
                ep._p.append(el)
            else:
                etb.element.getparent().remove(etb.element)
                result = self._add_math_image(slide, latex, eq_left, row_top, eq_w, display=True, fontsize=pt_size)
                if not result:
                    tb2 = self._add_textbox(slide, eq_left, row_top, eq_w, row_h)
                    tf2 = tb2.text_frame
                    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p2 = tf2.paragraphs[0]
                    p2.text = latex
                    p2.font.name = self.FONT
                    p2.font.size = self._fs(Pt(pt_size - 4))
                    p2.font.color.rgb = self.FG
        if sd.eq_vars:
            var_top = top + int(row_h * n) + legend_gap
            self._add_var_legend(slide, sd.eq_vars, var_top, sym_pt=18, desc_pt=15,
                                 row_h=legend_row)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_columns(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        region = self._content_region(has_title=bool(sd.h1))
        rleft, rtop, rwidth, rheight = region
        n = len(sd.columns)
        if n == 0:
            return
        cls = sd.slide_class or ""
        gap = int(CARD_GAP)
        if cls == "cols-2-wide-l":
            widths = [int(rwidth * 0.62) - gap // 2, int(rwidth * 0.38) - gap // 2]
        elif cls == "cols-2-wide-r":
            widths = [int(rwidth * 0.38) - gap // 2, int(rwidth * 0.62) - gap // 2]
        else:
            cw = (rwidth - gap * (n - 1)) // n
            widths = [cw] * n
        size = SZ_COL
        # Two passes: estimate every column first (wrap-aware), then place all
        # of them from one shared top — per-column centering left the column
        # heads at different heights whenever the contents were uneven.
        heights = []
        for i, col_lines in enumerate(sd.columns):
            col_w = widths[i] if i < len(widths) else widths[-1]
            ch = self._column_height(
                [l for l in col_lines if not l.strip().startswith("![")], size,
                col_w)
            heights.append(min(ch, rheight))
        top = rtop + max(0, (rheight - max(heights, default=0)) // 2)
        x = rleft
        for i, col_lines in enumerate(sd.columns):
            col_w = widths[i] if i < len(widths) else widths[-1]
            self._add_column_content(slide, col_lines, left=int(x), top=int(top),
                                     width=int(col_w), height=int(heights[i]), size=size)
            x += col_w + gap
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_sandwich(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        n = len(sd.columns)
        lead_h = int(self._estimate_text_height([sd.top_text], SZ_BODY,
                                                width=rwidth)) if sd.top_text else 0
        concl_h = int(Inches(1.0)) if sd.bottom_text else 0
        gaps = (int(BLOCK_GAP) if sd.top_text and n else 0) + (int(BLOCK_GAP) if sd.bottom_text and n else 0)
        # Size the column band to what the columns actually need (a fixed
        # 3.2in band left a hole between short columns and the conclusion);
        # images reserve the 2.5in slot _add_column_content gives them.
        needed = 0
        if n:
            gap_w = int(CARD_GAP)
            est_w = (rwidth - gap_w * (n - 1)) // n
            for col_lines in sd.columns:
                n_imgs = sum(1 for l in col_lines if l.strip().startswith("!["))
                h = self._column_height(
                    [l for l in col_lines if not l.strip().startswith("![")],
                    SZ_COL, est_w) + n_imgs * int(Inches(2.6))
                needed = max(needed, h)
        col_h = int(min(needed, rheight - lead_h - concl_h - gaps)) if n else 0
        block_h = lead_h + (int(BLOCK_GAP) if sd.top_text and n else 0) + \
                  (col_h if n else 0) + (int(BLOCK_GAP) if sd.bottom_text and n else 0) + concl_h
        cur = rtop + max(0, (rheight - block_h) // 2)
        if sd.top_text:
            tb = self._add_textbox(slide, rleft, cur, rwidth, lead_h)
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            self._set_rich_text(p, sd.top_text, SZ_BODY, self.SECONDARY)
            p.alignment = PP_ALIGN.CENTER
            cur += lead_h + int(BLOCK_GAP)
        if n > 0:
            gap = int(CARD_GAP)
            col_w = (rwidth - gap * (n - 1)) // n
            for i, col_lines in enumerate(sd.columns):
                left = rleft + i * (col_w + gap)
                # Through _add_column_content so box divs and images render
                # in sandwich columns exactly as they do in cols-N.
                self._add_column_content(slide, col_lines, left=int(left), top=int(cur),
                                         width=int(col_w), height=int(col_h), size=SZ_COL)
            cur += col_h + int(BLOCK_GAP)
        if sd.bottom_text:
            self._add_conclusion_box(slide, sd.bottom_text, rleft, int(cur), rwidth, concl_h)

    def build_figure(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region_with_lead(slide, sd)
        cap = self._fig_caption(sd.caption, sd.source)
        cap_h = int(self._estimate_text_height([cap], SZ_SMALL,
                                               width=rwidth)) if cap else 0
        desc_h = int(self._estimate_text_height(sd.body_lines, SZ_COL,
                                               width=rwidth)) if sd.body_lines else 0
        img_file = self._image_or_placeholder(sd.image_path) if sd.image_path else None
        img_top = rtop; ph = 0; pw = 0
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(rwidth * 0.9)
            max_h = int(rheight - cap_h - desc_h - Inches(0.3))
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96); ph = int(ih * scale * 914400 / 96)
        block_h = ph + cap_h + desc_h + (int(Inches(0.15)) if cap_h else 0) + (int(Inches(0.2)) if desc_h else 0)
        img_top = rtop + max(0, (rheight - block_h) // 2)
        if img_file:
            slide.shapes.add_picture(img_file, (SW - pw) // 2, int(img_top), pw, ph)
        cap_top = img_top + ph + int(Inches(0.15))
        if cap:
            tb = self._add_textbox(slide, rleft, int(cap_top), rwidth, cap_h)
            p = tb.text_frame.paragraphs[0]
            self._rich_line(p, cap, SZ_SMALL, self.MUTED)
            p.alignment = PP_ALIGN.CENTER
            tb.text_frame.word_wrap = True
        if sd.body_lines:
            self._add_body_text(slide, sd.body_lines, left=rleft,
                                top=int(cap_top + cap_h + Inches(0.2)),
                                width=rwidth, size=SZ_COL)

    def build_table(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.table_rows:
            return
        region = self._content_region(has_title=bool(sd.h1))
        left, _, width, _ = region
        blocks = []
        if sd.h2:
            blocks.append(("sub", int(Inches(0.4))))
        # The table reserves exactly what its measured rows need — no
        # headroom fudge, because the rows no longer grow at render time.
        blocks.append(("table", self._table_height(sd.table_rows, width)))
        if sd.bottom_text:
            blocks.append(("accent", int(Inches(0.9))))
        tops = self._stack_tops([h for _, h in blocks], region, mode="center",
                                gap=int(Inches(0.2)))
        for (kind, h), t in zip(blocks, tops):
            if kind == "sub":
                tb = self._add_textbox(slide, left, t, width, h)
                p = tb.text_frame.paragraphs[0]
                p.text = sd.h2
                p.font.name = self.FONT; p.font.size = self._fs(SZ_H3)
                p.font.bold = True; p.font.color.rgb = self.SECONDARY
            elif kind == "table":
                self._styled_table(slide, sd.table_rows, left, t, width, h)
            elif kind == "accent":
                self._add_accent_box(slide, sd.bottom_text, left, t, width, h)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_references(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.ref_items:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        items = sd.ref_items
        # Two columns when the list is long.
        two_col = len(items) > 6
        gap = int(Inches(0.5))
        col_w = (rwidth - gap) // 2 if two_col else rwidth
        ref_size = Pt(13)

        def fill(col_items, start_idx, left):
            tb = self._add_textbox(slide, int(left), rtop, int(col_w), rheight)
            tf = tb.text_frame; tf.word_wrap = True
            # Short lists sit vertically centered instead of hugging the top.
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            for j, (author, title, venue) in enumerate(col_items):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_before = PARA_GAP
                r = p.add_run(); r.text = f"[{start_idx + j + 1}] "
                r.font.name = self.FONT; r.font.size = self._fs(ref_size)
                r.font.bold = True; r.font.color.rgb = self.ACCENT
                if author:
                    r = p.add_run(); r.text = author + " "
                    r.font.name = self.FONT; r.font.size = self._fs(ref_size)
                    r.font.bold = True; r.font.color.rgb = self.FG
                if title:
                    r = p.add_run(); r.text = title + " "
                    r.font.name = self.FONT; r.font.size = self._fs(ref_size)
                    r.font.color.rgb = self.FG
                if venue:
                    r = p.add_run(); r.text = venue
                    r.font.name = self.FONT; r.font.size = self._fs(ref_size)
                    r.font.italic = True; r.font.color.rgb = self.MUTED

        if two_col:
            half = (len(items) + 1) // 2
            fill(items[:half], 0, rleft)
            fill(items[half:], half, rleft + col_w + gap)
        else:
            fill(items, 0, rleft)

    def build_timeline_h(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.timeline_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        line_y = rtop + rheight // 2 - int(Inches(0.4))
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rleft, line_y, rwidth, int(Pt(3)))
        line.fill.solid(); line.fill.fore_color.rgb = self.HAIRLINE
        line.line.fill.background(); self._no_shadow(line)
        item_w = rwidth // n
        for i, item in enumerate(sd.timeline_items):
            cx = rleft + int(item_w * i) + int(item_w / 2)
            hl = item.get("highlight")
            dot_color = self.ACCENT if hl else self.SECONDARY
            d = int(Pt(22) if hl else Pt(16))
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - d // 2, line_y + int(Pt(1.5)) - d // 2, d, d)
            dot.fill.solid(); dot.fill.fore_color.rgb = dot_color
            dot.line.color.rgb = self.WHITE; dot.line.width = Pt(2)
            self._no_shadow(dot)
            tb_left = rleft + int(item_w * i)
            yr = self._add_textbox(slide, tb_left, line_y - int(Inches(0.55)), int(item_w), int(Inches(0.4)))
            p = yr.text_frame.paragraphs[0]
            p.text = item.get("year", "")
            p.font.name = self.FONT_HEAD; p.font.size = self._fs(SZ_H3)
            p.font.bold = True; p.font.color.rgb = self.ACCENT_TEXT if hl else self.PRIMARY
            p.alignment = PP_ALIGN.CENTER
            txt = self._add_textbox(slide, tb_left + int(Pt(6)), line_y + int(Inches(0.3)),
                                    int(item_w) - int(Pt(12)), int(Inches(0.55)))
            p2 = txt.text_frame.paragraphs[0]
            p2.text = item.get("text", "")
            p2.font.name = self.FONT; p2.font.size = self._fs(SZ_SMALL)
            p2.font.bold = True; p2.font.color.rgb = self.FG
            p2.alignment = PP_ALIGN.CENTER
            txt.text_frame.word_wrap = True
            if item.get("detail"):
                dtl = self._add_textbox(slide, tb_left + int(Pt(6)), line_y + int(Inches(0.85)),
                                        int(item_w) - int(Pt(12)), int(Inches(0.9)))
                p3 = dtl.text_frame.paragraphs[0]
                self._rich_line(p3, item["detail"], SZ_FOOT, self.MUTED)
                p3.alignment = PP_ALIGN.CENTER
                dtl.text_frame.word_wrap = True

    def build_timeline_v(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.timeline_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        row_h = int(min(Inches(1.15), rheight / max(n, 1)))
        block_h = row_h * n
        top = rtop + max(0, (rheight - block_h) // 2)
        line_x = rleft + int(Inches(0.16))
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_x, top + int(Pt(8)),
                                      int(Pt(2.5)), int(block_h - Pt(16)))
        line.fill.solid(); line.fill.fore_color.rgb = self.HAIRLINE
        line.line.fill.background(); self._no_shadow(line)
        content_left = rleft + int(Inches(0.7))
        year_w = int(Inches(1.3))
        for i, item in enumerate(sd.timeline_items):
            ry = top + int(row_h * i)
            hl = item.get("highlight")
            d = int(Pt(16) if hl else Pt(13))
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, line_x + int(Pt(1.25)) - d // 2,
                                         ry + int(Pt(8)), d, d)
            dot.fill.solid(); dot.fill.fore_color.rgb = self.ACCENT if hl else self.SECONDARY
            dot.line.color.rgb = self.WHITE; dot.line.width = Pt(2); self._no_shadow(dot)
            yr = self._add_textbox(slide, content_left, ry, year_w, int(Inches(0.4)))
            p = yr.text_frame.paragraphs[0]
            p.text = item.get("year", "")
            p.font.name = self.FONT_HEAD; p.font.size = self._fs(SZ_H3)
            p.font.bold = True; p.font.color.rgb = self.ACCENT_TEXT if hl else self.PRIMARY
            tx = content_left + year_w + int(Inches(0.25))
            tw = rleft + rwidth - tx
            txt = self._add_textbox(slide, tx, ry, tw, int(Inches(0.4)))
            p2 = txt.text_frame.paragraphs[0]
            self._set_rich_text(p2, item.get("text", ""), SZ_BODY, self.FG)
            txt.text_frame.word_wrap = True
            if item.get("detail"):
                dtl = self._add_textbox(slide, tx, ry + int(Inches(0.38)), tw, int(Inches(0.5)))
                self._rich_line(dtl.text_frame.paragraphs[0], item["detail"],
                                SZ_SMALL, self.MUTED)
                dtl.text_frame.word_wrap = True

    def build_end(self, sd: SlideData):
        slide = self._blank_slide()
        if self.LAYOUT.end_bg == "dark":
            self._set_bg(slide, self.PRIMARY)
        elif self.LAYOUT.end_bg == "light":
            self._set_bg(slide, self.LIGHT)
        is_dark = self.LAYOUT.end_bg == "dark"
        h_color = self.WHITE if is_dark else self.PRIMARY
        sub_color = self._tint(self.PRIMARY, 0.85) if is_dark else self.MUTED
        accent = self._hero_accent(2.5) if is_dark else self.ACCENT
        cx = int(SW // 2)
        # accent hairline above the thank-you line
        self._hairline(slide, int(cx - Inches(0.55)), int(Inches(2.95)),
                       Inches(1.1), thickness=Pt(1.6), color=accent)
        # thank-you
        hero_pt = SZ_DISPLAY.pt * getattr(self.LAYOUT, "title_scale", 1.0)
        tb = self._add_textbox(slide, int(Inches(1)), int(Inches(3.15)),
                               int(SW - Inches(2)), int(self._fs(Pt(hero_pt * 1.3))))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = sd.h1 or "Thank You"
        p.font.name = self.FONT_HEAD; p.font.size = self._fs(Pt(hero_pt))
        p.font.bold = True; p.font.color.rgb = h_color; p.alignment = PP_ALIGN.CENTER
        # sub lines
        remaining = [strip_html(s.strip()) for s in sd.raw.split("\n")
                     if s.strip() and not s.strip().startswith(("#", "<!--", "<div", "</div", ">"))]
        if remaining:
            sb = self._add_textbox(slide, int(Inches(1)), int(Inches(4.55)),
                                   int(SW - Inches(2)), int(Inches(0.9)))
            sb.text_frame.word_wrap = True
            for i, line in enumerate(remaining[:2]):
                sp = sb.text_frame.paragraphs[0] if i == 0 else sb.text_frame.add_paragraph()
                sp.text = line; sp.font.name = self.FONT
                sp.font.size = self._fs(SZ_BODY); sp.font.color.rgb = sub_color
                sp.alignment = PP_ALIGN.CENTER
                if i:
                    sp.space_before = Pt(4)
        if len(remaining) > 2:
            cb = self._add_textbox(slide, int(Inches(1)), int(Inches(5.65)),
                                   int(SW - Inches(2)), int(Inches(0.4)))
            self._kicker_para(cb.text_frame.paragraphs[0], remaining[2],
                              color=sub_color, align=PP_ALIGN.CENTER, tracking=160)

    def build_zone_flow(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.zone_flow_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region_with_lead(slide, sd)
        gap = int(Inches(0.3))
        arrow_w = int(Inches(0.34))
        box_w = (rwidth - arrow_w * (n - 1) - gap * (n - 1)) // n
        box_h = int(min(Inches(3.6), rheight))
        box_top = rtop + max(0, (rheight - box_h) // 2)
        for i, item in enumerate(sd.zone_flow_items):
            x = rleft + i * (box_w + gap + arrow_w)
            self._add_zone_box(slide, int(x), box_top, int(box_w), int(box_h),
                              label=item.get("label", ""), body=item.get("body", ""))
            if i < n - 1:
                ax = int(x + box_w + gap // 2)
                ay = int(box_top + box_h // 2 - int(Pt(11)))
                arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, ax, ay, int(arrow_w), int(Pt(22)))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.ACCENT
                arrow.line.fill.background()
                self._no_shadow(arrow)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_zone_compare(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        mid = int(Inches(1.0))
        box_w = (rwidth - mid) // 2
        box_h = int(min(Inches(4.2), rheight))
        box_top = rtop + max(0, (rheight - box_h) // 2)
        self._add_zone_box(slide, rleft, box_top, int(box_w), int(box_h),
                          label=sd.zone_compare.get("left_label", ""),
                          body=sd.zone_compare.get("left_body", ""))
        # VS badge (circular)
        d = int(Inches(0.72))
        vs_x = int(rleft + box_w + (mid - d) // 2)
        vs_y = int(box_top + box_h // 2 - d // 2)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, vs_x, vs_y, d, d)
        badge.fill.solid(); badge.fill.fore_color.rgb = self.ACCENT
        badge.line.fill.background(); self._no_shadow(badge)
        bp = badge.text_frame.paragraphs[0]
        bp.text = sd.zone_compare.get("vs_text", "VS")
        bp.font.name = self.FONT_HEAD; bp.font.size = self._fs(SZ_H3)
        bp.font.bold = True; bp.font.color.rgb = self.WHITE
        bp.alignment = PP_ALIGN.CENTER
        self._add_zone_box(slide, int(rleft + box_w + mid), box_top, int(box_w), int(box_h),
                          label=sd.zone_compare.get("right_label", ""),
                          body=sd.zone_compare.get("right_body", ""))
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_zone_matrix(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        cells = sd.zone_matrix.get("cells", [{}, {}, {}, {}])
        x_label = sd.zone_matrix.get("x_label", "")
        y_label = sd.zone_matrix.get("y_label", "")
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        # Reserve strips for the axis labels (the defining feature of a 2-axis
        # matrix). Only when present, so label-less decks are unaffected.
        strip = int(Inches(0.42))
        lblgap = int(Inches(0.12))
        gleft, gtop, gwidth, gheight = rleft, rtop, rwidth, rheight
        if y_label:
            gleft += strip + lblgap
            gwidth -= strip + lblgap
        if x_label:
            gheight -= strip + lblgap
        gap = int(CARD_GAP)
        cell_w = (gwidth - gap) // 2
        cell_h = (gheight - gap) // 2
        positions = [
            (gleft, gtop),
            (int(gleft + cell_w + gap), gtop),
            (gleft, int(gtop + cell_h + gap)),
            (int(gleft + cell_w + gap), int(gtop + cell_h + gap)),
        ]
        for i, (x, y) in enumerate(positions):
            if i < len(cells):
                self._add_zone_box(slide, int(x), int(y), int(cell_w), int(cell_h),
                                  label=cells[i].get("label", ""),
                                  body=cells[i].get("body", ""))
        # y-axis label: vertical (rotated 270°) along the left of the grid.
        if y_label:
            yb = self._add_textbox(slide, int(gleft - strip - lblgap + (strip - gheight) // 2),
                                   int(gtop + gheight // 2 - strip // 2),
                                   int(gheight), int(strip))
            yb.rotation = 270
            yp = yb.text_frame.paragraphs[0]
            yp.text = y_label
            yp.font.name = self.FONT_HEAD; yp.font.size = self._fs(SZ_SMALL)
            yp.font.bold = True; yp.font.color.rgb = self.MUTED
            yp.alignment = PP_ALIGN.CENTER
            yb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # x-axis label: horizontal, centered under the grid.
        if x_label:
            xb = self._add_textbox(slide, int(gleft), int(gtop + gheight + lblgap),
                                   int(gwidth), int(strip))
            xp = xb.text_frame.paragraphs[0]
            xp.text = x_label
            xp.font.name = self.FONT_HEAD; xp.font.size = self._fs(SZ_SMALL)
            xp.font.bold = True; xp.font.color.rgb = self.MUTED
            xp.alignment = PP_ALIGN.CENTER
            xb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_zone_process(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.zone_process_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region_with_lead(slide, sd)
        gap = int(Inches(0.28))
        box_w = (rwidth - gap * (n - 1)) // n
        badge_d = int(Inches(0.56))
        block_h = int(min(Inches(4.0) + badge_d, rheight))
        box_h = block_h - badge_d - int(Inches(0.18))
        top = rtop + max(0, (rheight - block_h) // 2)
        for i, item in enumerate(sd.zone_process_items):
            x = rleft + i * (box_w + gap)
            bx = int(x + box_w // 2 - badge_d // 2)
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, int(top), badge_d, badge_d)
            badge.fill.solid(); badge.fill.fore_color.rgb = self.SECONDARY
            badge.line.fill.background(); self._no_shadow(badge)
            bp = badge.text_frame.paragraphs[0]
            bp.text = item.get("step", str(i + 1))
            bp.font.name = self.FONT_HEAD; bp.font.size = self._fs(SZ_H3)
            bp.font.bold = True; bp.font.color.rgb = self.WHITE
            bp.alignment = PP_ALIGN.CENTER
            self._add_zone_box(slide, int(x), int(top + badge_d + Inches(0.18)),
                              int(box_w), int(box_h),
                              label=item.get("title", ""), body=item.get("body", ""))
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_sections(self, sd: SlideData):
        """Hearing-style dense stack: 2-4 topic bands, each a colored lead
        line + indented body, separated by hairlines. The high-information
        workhorse — one slide carries what looser types spread over three."""
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        items = sd.sections_items
        if not items:
            return
        rleft, rtop, rwidth, rheight = self._content_region_with_lead(slide, sd)
        indent = int(Inches(0.28))
        gap = int(Inches(0.18))
        lead_h = int(Pt(SZ_H3.pt * metrics.DEFAULT_LINE_FACTOR))
        heights = []
        for it in items:
            bh = int(self._estimate_text_height(
                it["body"].split("\n"), SZ_ZONE_B,
                width=rwidth - indent, line_spacing=LINE_BODY)) if it["body"] else 0
            heights.append((lead_h if it["title"] else 0)
                           + (int(Pt(3)) if it["title"] and it["body"] else 0) + bh)
        total = sum(heights) + gap * 2 * max(0, len(items) - 1)
        top = rtop + (max(0, (rheight - total) // 2)
                      if self.LAYOUT.vertical_align == "center" else 0)
        cur = top
        for i, (it, h) in enumerate(zip(items, heights)):
            if it["title"]:
                tb = self._add_textbox(slide, rleft, int(cur), rwidth, lead_h)
                p = tb.text_frame.paragraphs[0]
                self._set_rich_text(p, it["title"], SZ_H3, self.SECONDARY)
                for r in p.runs:
                    r.font.bold = True
                cur += lead_h + (int(Pt(3)) if it["body"] else 0)
            if it["body"]:
                bh = h - (lead_h + int(Pt(3)) if it["title"] else 0)
                bt = self._add_textbox(slide, rleft + indent, int(cur),
                                       rwidth - indent, max(bh, int(Pt(14))))
                tf = bt.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.NONE
                for j, line in enumerate(it["body"].split("\n")):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    self._set_rich_text(p, line, SZ_ZONE_B, self.FG)
                    p.line_spacing = LINE_BODY
                cur += bh
            if i < len(items) - 1:
                self._hairline(slide, rleft, int(cur + gap - Pt(1)), rwidth,
                               color=self.HAIRLINE)
                cur += gap * 2
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_agenda(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.agenda_items:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        n = len(sd.agenda_items)
        # 1.15in rows let a 3-4 item agenda occupy the canvas instead of
        # huddling in the middle; longer agendas divide the space as before.
        row_h = int(min(Inches(1.15), rheight / max(n, 1)))
        block_h = row_h * n
        top = rtop + max(0, (rheight - block_h) // 2)
        num_w = int(Inches(0.9))
        for i, item in enumerate(sd.agenda_items):
            y = top + row_h * i
            nb = self._add_textbox(slide, rleft, y, num_w, row_h)
            nb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            np_ = nb.text_frame.paragraphs[0]
            np_.text = f"{i + 1:02d}"
            np_.font.name = self.FONT_HEAD; np_.font.size = self._fs(Pt(26))
            np_.font.bold = True; np_.font.color.rgb = self.ACCENT
            tb = self._add_textbox(slide, rleft + num_w + int(Inches(0.2)), y,
                                   rwidth - num_w - int(Inches(0.2)), row_h)
            tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            tb.text_frame.word_wrap = True
            tp = tb.text_frame.paragraphs[0]
            tp.text = item
            tp.font.name = self.FONT; tp.font.size = self._fs(SZ_H2)
            tp.font.color.rgb = self.FG
            if i < n - 1:
                self._hairline(slide, rleft, y + row_h - int(Pt(1)), rwidth,
                               color=self._tint(self.MUTED, 0.78))

    def build_rq(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        card_w = int(rwidth * 0.82)
        card_x = rleft + (rwidth - card_w) // 2
        # Estimate at the card's INNER text width so a long question that wraps
        # to 2-3 lines sizes the card tall enough — otherwise the text spilled
        # out the bottom and collided with rq_sub.
        main_tw = card_w - int(Inches(0.6))
        # The card's own chrome (accent bar + padding top and bottom) has to be
        # added on top of the text height, or the inner box ends up shorter
        # than the question it holds.
        card_chrome = int(CARD_ACCENT_H) + int(CARD_PAD) * 2
        main_h = int(self._estimate_text_height([sd.rq_main], SZ_H2, width=main_tw)) + card_chrome if sd.rq_main else 0
        sub_h = int(Inches(0.8)) if sd.rq_sub else 0
        gap = int(Inches(0.3)) if sd.rq_sub else 0
        block_h = main_h + gap + sub_h
        top = rtop + max(0, (rheight - block_h) // 2)
        if sd.rq_main:
            self._add_card(slide, (card_x, top, card_w, main_h),
                           body=sd.rq_main, accent_bar=True, anchor="middle",
                           body_size=SZ_H2)
            # center & emphasize the question text
            # (re-style: the card body paragraph)
        if sd.rq_sub:
            tb2 = self._add_textbox(slide, card_x, top + main_h + gap, card_w, sub_h)
            tf2 = tb2.text_frame; tf2.word_wrap = True
            tf2.vertical_anchor = MSO_ANCHOR.TOP
            p2 = tf2.paragraphs[0]
            self._rich_line(p2, sd.rq_sub, SZ_SMALL, self.MUTED)
            p2.alignment = PP_ALIGN.CENTER

    def build_result_dual(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.result_dual_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        gap = int(CARD_GAP)
        col_w = (rwidth - gap * (n - 1)) // n
        cap_h = int(Inches(0.45))
        for i, item in enumerate(sd.result_dual_items):
            x = rleft + i * (col_w + gap)
            img_file = self._image_or_placeholder(item.get("image", ""))
            ph = 0
            img_top = rtop
            if img_file:
                from PIL import Image
                with Image.open(img_file) as im:
                    iw, ih = im.size
                max_w = int(col_w * 0.98)
                max_h = int(rheight - cap_h - Inches(0.15))
                scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
                pw = int(iw * scale * 914400 / 96)
                ph = int(ih * scale * 914400 / 96)
                block_h = ph + (cap_h if item.get("caption") else 0)
                img_top = rtop + max(0, (rheight - block_h) // 2)
                img_left = int(x) + (int(col_w) - pw) // 2
                slide.shapes.add_picture(img_file, img_left, int(img_top), pw, ph)
            cap_text = item.get("caption", "")
            if cap_text:
                ctb = self._add_textbox(slide, int(x), int(img_top + ph + Inches(0.12)),
                                        int(col_w), cap_h)
                cp = ctb.text_frame.paragraphs[0]
                self._rich_line(cp, cap_text, SZ_SMALL, self.MUTED)
                cp.alignment = PP_ALIGN.CENTER

    def build_summary(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.summary_points:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        n = len(sd.summary_points)
        row_h = int(min(Inches(1.0), rheight / max(n, 1)))
        block_h = row_h * n
        top = rtop + max(0, (rheight - block_h) // 2)
        bar_w = int(Pt(5))
        text_x = rleft + bar_w + int(Inches(0.28))
        for i, pt in enumerate(sd.summary_points):
            y = top + row_h * i
            # accent bar synced to the row top (fixes the old desync)
            bh = int(min(row_h - Inches(0.12), Inches(0.7)))
            bdr = self._hairline(slide, rleft, y + (row_h - bh) // 2, bh,
                                 thickness=bar_w, color=self.ACCENT, vertical=True)
            tb = self._add_textbox(slide, text_x, y, rwidth - (text_x - rleft), row_h)
            tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            tb.text_frame.word_wrap = True
            self._set_rich_text(tb.text_frame.paragraphs[0], pt, SZ_H3, self.FG)

    def build_appendix(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            # Leave the corner clear for the APPENDIX label instead of running
            # the title box underneath it.
            self._add_title(slide, sd.h1, color=self.MUTED,
                            width=int(CONTENT_W - Inches(2.6)) if sd.appendix_label else None)
        if sd.appendix_label:
            lbl = self._add_textbox(slide, int(SW - Inches(2.6)), int(MARGIN_T),
                                    int(Inches(2.0)), int(KICKER_H))
            self._kicker_para(lbl.text_frame.paragraphs[0], sd.appendix_label,
                              color=self.MUTED, align=PP_ALIGN.RIGHT)
        if sd.table_rows:
            # Render the table on THIS slide — delegating to build_table()
            # would open a fresh slide and split the appendix in two.
            region = self._content_region(has_title=bool(sd.h1))
            left, _, width, _ = region
            blocks = [("table", self._table_height(sd.table_rows, width))]
            if sd.bottom_text:
                blocks.append(("accent", int(Inches(0.9))))
            tops = self._stack_tops([h for _, h in blocks], region, mode="center",
                                    gap=int(Inches(0.25)))
            for (kind, h), t in zip(blocks, tops):
                if kind == "table":
                    self._styled_table(slide, sd.table_rows, left, t, width, h)
                else:
                    self._add_accent_box(slide, sd.bottom_text, left, t, width, h)
            if sd.footnote:
                self._add_footnote(slide, sd.footnote)
        elif sd.body_lines:
            region = self._content_region(has_title=bool(sd.h1))
            left, _, width, _ = region
            bh = int(self._estimate_text_height(sd.body_lines, SZ_SMALL, width=width))
            top = self._stack_tops([min(bh, region[3])], region, mode="center")[0]
            self._add_body_text(slide, sd.body_lines, left=left, top=top,
                                width=width, height=min(bh, region[3]), size=SZ_SMALL)

    def _build_image_points(self, sd, lead_text, image_path, points):
        """Shared layout for overview/result: optional lead line, then a
        figure on the left and bullet points on the right, vertically centered."""
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        cur_top = rtop
        avail_h = rheight
        if lead_text:
            # Measured, not a fixed 0.55in: the lead is the first thing to
            # clip once font_scale or the keynote density enlarges the type.
            lead_h = int(self._estimate_text_height([lead_text], SZ_H3,
                                                    width=rwidth))
            tb = self._add_textbox(slide, rleft, cur_top, rwidth, lead_h)
            p = tb.text_frame.paragraphs[0]
            self._rich_line(p, lead_text, SZ_H3, self.SECONDARY)
            tb.text_frame.word_wrap = True
            step = lead_h + int(Inches(0.15))
            cur_top += step; avail_h -= step
        left_w = int(rwidth * 0.56)
        right_w = int(rwidth * 0.4)
        right_x = rleft + rwidth - right_w
        # Reserve room for the figure caption — overview/result dropped it
        # entirely (the figure showed, but `Fig. N. …` never rendered).
        cap = self._fig_caption(sd.caption, sd.source)
        cap_h = int(Inches(0.5)) if cap else 0
        img_file = self._image_or_placeholder(image_path) if image_path else None
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(left_w * 0.98); max_h = int((avail_h - cap_h) * 0.96)
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96); ph = int(ih * scale * 914400 / 96)
            img_top = cur_top + max(0, (avail_h - ph - cap_h) // 2)
            slide.shapes.add_picture(img_file, rleft, int(img_top), pw, ph)
            if cap:
                ctb = self._add_textbox(slide, rleft, int(img_top + ph + Inches(0.1)),
                                        left_w, cap_h)
                cp = ctb.text_frame.paragraphs[0]
                self._rich_line(cp, cap, SZ_SMALL, self.MUTED)
                cp.alignment = PP_ALIGN.CENTER
                ctb.text_frame.word_wrap = True
        if points:
            ph_pts = int(self._estimate_text_height([f"\u2022 {p}" for p in points], SZ_COL,
                                                    width=right_w, gap=PARA_GAP))
            pts_top = cur_top + max(0, (avail_h - min(ph_pts, avail_h)) // 2)
            tb = self._add_textbox(slide, right_x, int(pts_top), right_w,
                                   int(min(ph_pts, avail_h)))
            tf = tb.text_frame; tf.word_wrap = True
            for i, pt in enumerate(points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = PARA_GAP
                self._set_rich_text(p, f"\u2022 {pt}", SZ_COL, self.FG)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_overview(self, sd: SlideData):
        self._build_image_points(sd, sd.overview_text, sd.image_path, sd.overview_points)

    def build_result(self, sd: SlideData):
        self._build_image_points(sd, sd.result_text, sd.result_figure, sd.result_analysis)

    def build_steps(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.steps_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        gap = int(Inches(0.28))
        box_w = (rwidth - gap * (n - 1)) // n
        badge_d = int(Inches(0.56))
        block_h = int(min(Inches(3.9) + badge_d, rheight))
        box_h = block_h - badge_d - int(Inches(0.18))
        top = rtop + max(0, (rheight - block_h) // 2)
        for i, item in enumerate(sd.steps_items):
            x = rleft + i * (box_w + gap)
            bx = int(x + box_w // 2 - badge_d // 2)
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, int(top), badge_d, badge_d)
            badge.fill.solid(); badge.fill.fore_color.rgb = self.ACCENT
            badge.line.fill.background(); self._no_shadow(badge)
            bp = badge.text_frame.paragraphs[0]
            bp.text = item.get("num", str(i + 1))
            bp.font.name = self.FONT_HEAD; bp.font.size = self._fs(SZ_H3)
            bp.font.bold = True; bp.font.color.rgb = self.WHITE
            bp.alignment = PP_ALIGN.CENTER
            self._add_zone_box(slide, int(x), int(top + badge_d + Inches(0.18)),
                              int(box_w), int(box_h),
                              label=item.get("title", ""), body=item.get("body", ""))
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_quote(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        qx = rleft + int(Inches(1.0))
        qw = rwidth - int(Inches(1.6))
        # width-aware so a long quote that wraps sizes its block correctly —
        # otherwise the attribution landed on top of the wrapped quote text.
        quote_h = int(self._estimate_text_height([sd.quote_text or ""], SZ_H2, width=qw))
        block_h = int(Inches(0.9)) + quote_h + (int(Inches(0.5)) if sd.quote_source else 0)
        top = rtop + max(0, (rheight - block_h) // 2)
        # opening quotation mark (ghost) \u2014 the box has to be a full line tall
        # for an 80pt glyph, or it reads as a clipped text box to any checker.
        # Its column also stops short of the quote text: a decorative glyph
        # overlapping the text column is what made this read as a collision.
        mark_pt = self._fs(Pt(80))
        mark = self._add_textbox(slide, rleft, top, int(Inches(0.9)),
                                 int(Pt(mark_pt.pt * metrics.DEFAULT_LINE_FACTOR)))
        mp = mark.text_frame.paragraphs[0]
        mp.text = "\u201C"
        mp.font.name = self.FONT_HEAD; mp.font.size = mark_pt
        mp.font.bold = True; mp.font.color.rgb = self._tint(self.ACCENT, 0.55)
        # left accent bar beside the quote
        qtop = top + int(Inches(0.55))
        self._hairline(slide, qx - int(Inches(0.25)), qtop, max(quote_h, int(Inches(0.5))),
                       thickness=Pt(3), color=self.ACCENT, vertical=True)
        if sd.quote_text:
            tb = self._add_textbox(slide, qx, qtop, qw, quote_h + int(Inches(0.2)))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            self._rich_line(p, sd.quote_text, SZ_H2, self.FG)
            p.line_spacing = LINE_BODY
        if sd.quote_source:
            stb = self._add_textbox(slide, qx, qtop + quote_h + int(Inches(0.25)),
                                    qw, int(Inches(0.5)))
            sp = stb.text_frame.paragraphs[0]
            sp.text = f"\u2014 {sd.quote_source}"
            sp.font.name = self.FONT; sp.font.size = self._fs(SZ_SMALL)
            sp.font.color.rgb = self.MUTED

    def build_history(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.history_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        # Same rhythm as agenda: taller rows let a 3-5 item history own the
        # canvas instead of pooling in the middle band.
        row_h = int(min(Inches(1.15), rheight / max(n, 1)))
        block_h = row_h * n
        top = rtop + max(0, (rheight - block_h) // 2)
        year_w = int(Inches(1.5))
        ev_x = rleft + year_w + int(Inches(0.35))
        for i, item in enumerate(sd.history_items):
            y = top + row_h * i
            yr = self._add_textbox(slide, rleft, y, year_w, row_h)
            yr.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = yr.text_frame.paragraphs[0]
            p.text = item.get("year", "")
            p.font.name = self.FONT_HEAD; p.font.size = self._fs(SZ_H3)
            p.font.bold = True; p.font.color.rgb = self.ACCENT_TEXT
            p.alignment = PP_ALIGN.RIGHT
            # vertical connector tick
            self._hairline(slide, ev_x - int(Inches(0.18)), y + int(Inches(0.04)),
                           row_h - int(Inches(0.08)), thickness=Pt(2),
                           color=self.HAIRLINE, vertical=True)
            ev = self._add_textbox(slide, ev_x, y, rwidth - (ev_x - rleft), row_h)
            ev.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            ev.text_frame.word_wrap = True
            p2 = ev.text_frame.paragraphs[0]
            self._set_rich_text(p2, item.get("event", ""), SZ_BODY, self.FG)
            if i < n - 1:
                self._hairline(slide, rleft, y + row_h - int(Pt(1)), rwidth,
                               color=self._tint(self.MUTED, 0.82))

    def build_panorama(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1), full=True)
        cap_h = int(self._estimate_text_height([sd.panorama_text], SZ_SMALL,
                                               width=rwidth)) if sd.panorama_text else 0
        img_file = self._image_or_placeholder(sd.image_path) if sd.image_path else None
        img_top = rtop; ph = 0
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(rwidth); max_h = int(rheight - cap_h - Inches(0.15))
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96); ph = int(ih * scale * 914400 / 96)
            img_top = rtop + max(0, (rheight - ph - cap_h) // 2)
            slide.shapes.add_picture(img_file, (SW - pw) // 2, int(img_top), pw, ph)
        if sd.panorama_text:
            tb = self._add_textbox(slide, rleft, int(img_top + ph + Inches(0.12)), rwidth, cap_h)
            p = tb.text_frame.paragraphs[0]
            self._rich_line(p, sd.panorama_text, SZ_SMALL, self.MUTED)
            p.alignment = PP_ALIGN.CENTER
            tb.text_frame.word_wrap = True

    def build_kpi(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.kpi_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region_with_lead(slide, sd)
        gap = int(CARD_GAP)
        box_w = (rwidth - gap * (n - 1)) // n
        # Height from the content — a fixed 2.7in card left a value and a
        # one-line caption floating in an empty box.
        inner_w = box_w - int(CARD_PAD) * 2
        label_h = 0
        for item in sd.kpi_items:
            label_h = max(label_h, int(self._estimate_text_height(
                [item.get("label", "")], SZ_SMALL, width=inner_w)))
        box_h = int(min(rheight,
                        max(int(Inches(1.7)),
                            int(self._fs(SZ_METRIC) * metrics.DEFAULT_LINE_FACTOR)
                            + label_h + int(CARD_PAD) * 2 + int(CARD_ACCENT_H))))
        kpi_top = rtop + max(0, (rheight - box_h) // 2)
        for i, item in enumerate(sd.kpi_items):
            x = rleft + i * (box_w + gap)
            self._add_card(slide, (int(x), kpi_top, int(box_w), box_h),
                           value=item.get("value", ""), label=item.get("label", ""),
                           accent_bar=True, anchor="middle",
                           value_size=self._fs(SZ_METRIC),
                           label_size=SZ_SMALL, label_color=self.MUTED)

    def build_pros_cons(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        gap = int(CARD_GAP)
        half_w = (rwidth - gap) // 2
        box_h = int(min(Inches(4.4), rheight))
        top = rtop + max(0, (rheight - box_h) // 2)
        head_h = int(Inches(0.62))
        pad = int(Pt(16))
        green = RGBColor(0x3F, 0x7A, 0x52)
        red = RGBColor(0xA8, 0x45, 0x3A)
        for idx, (items, color, label_text, mark) in enumerate([
            (sd.pros_items, green, "Pros", "\u2713"),
            (sd.cons_items, red, "Cons", "\u2717"),
        ]):
            x = rleft + idx * (half_w + gap)
            # card body
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), top, int(half_w), box_h)
            bg.adjustments[0] = CARD_RADIUS
            bg.fill.solid(); bg.fill.fore_color.rgb = self._tint(color, 0.92)
            bg.line.color.rgb = self._tint(color, 0.5); bg.line.width = HAIRLINE_W
            self._no_shadow(bg)
            # header band
            hb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), top, int(half_w), head_h)
            hb.fill.solid(); hb.fill.fore_color.rgb = color
            hb.line.fill.background(); self._no_shadow(hb)
            hp = hb.text_frame.paragraphs[0]
            hb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            hb.text_frame.margin_left = pad
            hp.text = f"{mark}  {label_text}"
            hp.font.name = self.FONT_HEAD; hp.font.size = self._fs(SZ_ZONE_L)
            hp.font.bold = True; hp.font.color.rgb = self.WHITE
            # items
            itb = self._add_textbox(slide, int(x) + pad, top + head_h + int(Pt(12)),
                                    int(half_w) - pad * 2, box_h - head_h - int(Pt(24)))
            tf = itb.text_frame; tf.word_wrap = True
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = PARA_GAP
                self._set_rich_text(p, f"{mark}  {item}", SZ_ZONE_B, self.FG)

    def build_definition(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        bar_x = rleft
        tx = rleft + int(Inches(0.4))
        tw = rwidth - int(Inches(0.4))
        term_h = int(Inches(0.7)) if sd.def_term else 0
        # width-aware: a long definition wraps, and an under-estimate let the
        # body text run under the note below it.
        body_h = int(self._estimate_text_height([sd.def_body], SZ_BODY, width=tw,
                                                 line_spacing=LINE_BODY)) if sd.def_body else 0
        note_h = int(self._estimate_text_height([sd.def_note], SZ_SMALL,
                                                width=tw)) if sd.def_note else 0
        gaps = (int(Inches(0.2)) if sd.def_term and sd.def_body else 0) + \
               (int(Inches(0.3)) if sd.def_note else 0)
        block_h = term_h + body_h + note_h + gaps
        top = rtop + max(0, (rheight - block_h) // 2)
        # left accent bar spanning term+body
        self._hairline(slide, bar_x, top, max(term_h + body_h, int(Inches(0.6))),
                       thickness=Pt(3), color=self.ACCENT, vertical=True)
        cur = top
        if sd.def_term:
            self._eyebrow(slide, "Definition", tx, cur - int(KICKER_H) + int(Inches(0.02)), tw)
            tb = self._add_textbox(slide, tx, cur, tw, term_h)
            p = tb.text_frame.paragraphs[0]
            self._rich_line(p, sd.def_term, Pt(28), self.PRIMARY,
                            font=self.FONT_HEAD, bold=True)
            cur += term_h + int(Inches(0.2))
        if sd.def_body:
            tb2 = self._add_textbox(slide, tx, cur, tw, body_h)
            tf2 = tb2.text_frame; tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            self._set_rich_text(p2, sd.def_body, SZ_BODY, self.FG)
            p2.line_spacing = LINE_BODY
            cur += body_h + int(Inches(0.3))
        if sd.def_note:
            ntb = self._add_textbox(slide, tx, cur, tw, note_h)
            ntb.text_frame.word_wrap = True
            self._rich_line(ntb.text_frame.paragraphs[0], sd.def_note,
                            SZ_SMALL, self.MUTED)

    def build_diagram(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region_with_lead(slide, sd, full=not sd.caption)
        cap = self._fig_caption(sd.caption, sd.source)
        cap_h = (max(int(Inches(0.5)),
                     int(self._estimate_text_height([cap], SZ_SMALL, width=rwidth)))
                 if cap else 0)
        img_file = self._image_or_placeholder(sd.image_path) if sd.image_path else None
        img_top = rtop; ph = 0
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(rwidth * 0.94)
            max_h = int(rheight - cap_h - Inches(0.15))
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96); ph = int(ih * scale * 914400 / 96)
            block_h = ph + cap_h
            img_top = rtop + max(0, (rheight - block_h) // 2)
            left = (SW - pw) // 2
            slide.shapes.add_picture(img_file, left, int(img_top), pw, ph)
        if cap:
            ctb = self._add_textbox(slide, rleft, int(img_top + ph + Inches(0.12)),
                                    rwidth, cap_h)
            ctb.text_frame.word_wrap = True
            p = ctb.text_frame.paragraphs[0]
            self._rich_line(p, cap, SZ_SMALL, self.MUTED)
            p.alignment = PP_ALIGN.CENTER

    def build_gallery_img(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.gallery_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        cols = min(n, 3)
        rows_n = (n + cols - 1) // cols
        gap = int(CARD_GAP)
        cell_w = (rwidth - gap * (cols - 1)) // cols
        cell_h = (rheight - gap * (rows_n - 1)) // rows_n
        for idx, item in enumerate(sd.gallery_items):
            r, c = divmod(idx, cols)
            x = rleft + c * (cell_w + gap)
            y = rtop + r * (cell_h + gap)
            # subtle card behind each image
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y),
                                          int(cell_w), int(cell_h))
            card.adjustments[0] = CARD_RADIUS
            card.fill.solid(); card.fill.fore_color.rgb = self.SURFACE
            card.line.color.rgb = self.HAIRLINE; card.line.width = HAIRLINE_W
            self._no_shadow(card)
            cap_text = item.get("caption", "")
            cap_h = int(Inches(0.34)) if cap_text else 0
            img_file = self._image_or_placeholder(item.get("image", ""))
            if img_file:
                from PIL import Image
                with Image.open(img_file) as im:
                    iw, ih = im.size
                max_w = int(cell_w * 0.86)
                max_h = int((cell_h - cap_h) * 0.82)
                scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
                pw = int(iw * scale * 914400 / 96); ph = int(ih * scale * 914400 / 96)
                ix = int(x) + (int(cell_w) - pw) // 2
                iy = int(y) + (int(cell_h) - cap_h - ph) // 2
                slide.shapes.add_picture(img_file, ix, iy, pw, ph)
            # per-image caption — was dropped entirely (only the title showed)
            if cap_text:
                cb = self._add_textbox(slide, int(x) + int(Inches(0.1)),
                                       int(y) + int(cell_h) - cap_h,
                                       int(cell_w) - int(Inches(0.2)), cap_h)
                cp = cb.text_frame.paragraphs[0]
                self._rich_line(cp, cap_text, SZ_FOOT, self.MUTED)
                cp.alignment = PP_ALIGN.CENTER
                cb.text_frame.word_wrap = True
                cb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def build_highlight(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.highlight_text:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        # width-aware: a long highlight wraps to several lines; an under-estimate
        # made the band too short, so the text spilled out and the accent rule
        # landed on top of the first line.
        text_h = int(self._estimate_text_height([sd.highlight_text], Pt(30),
                                                width=rwidth - int(Inches(1.0))))
        band_h = min(rheight, text_h + int(Inches(1.0)))
        band_top = rtop + max(0, (rheight - band_h) // 2)
        # soft surface band
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rleft, band_top, rwidth, band_h)
        band.fill.solid(); band.fill.fore_color.rgb = self.SURFACE
        band.line.fill.background(); self._no_shadow(band)
        cx = int(SW // 2)
        self._hairline(slide, cx - int(Inches(0.6)), band_top + int(Inches(0.18)),
                       Inches(1.2), thickness=ACCENT_RULE_W, color=self.ACCENT)
        tb = self._add_textbox(slide, rleft + int(Inches(0.5)), band_top,
                               rwidth - int(Inches(1.0)), band_h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        self._set_rich_text(p, sd.highlight_text, Pt(30), self.PRIMARY)
        for r in p.runs:
            r.font.bold = True; r.font.name = self.FONT_HEAD
        p.alignment = PP_ALIGN.CENTER; p.line_spacing = LINE_TITLE

    def build_checklist(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.checklist_items:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        n = len(sd.checklist_items)
        row_h = int(min(Inches(0.7), rheight / max(n, 1)))
        block_h = row_h * n
        top = rtop + max(0, (rheight - block_h) // 2)
        mark_w = int(Inches(0.5))
        for i, item in enumerate(sd.checklist_items):
            y = top + row_h * i
            done = item.get("done")
            d = max(int(Inches(0.26)),
                    int(Pt(self._fs(Pt(13)).pt * metrics.DEFAULT_LINE_FACTOR)))
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rleft,
                                         y + (row_h - d) // 2, d, d)
            box.adjustments[0] = 0.2
            if done:
                box.fill.solid(); box.fill.fore_color.rgb = self.ACCENT
                box.line.fill.background()
                ctf = box.text_frame
                ctf.margin_top = ctf.margin_bottom = 0
                ctf.margin_left = ctf.margin_right = 0
                cp = ctf.paragraphs[0]; cp.text = "\u2713"
                cp.font.size = self._fs(Pt(13)); cp.font.bold = True
                cp.font.color.rgb = self.WHITE; cp.alignment = PP_ALIGN.CENTER
            else:
                box.fill.background()
                box.line.color.rgb = self.MUTED; box.line.width = Pt(1)
            self._no_shadow(box)
            tb = self._add_textbox(slide, rleft + mark_w + int(Inches(0.18)), y,
                                   rwidth - mark_w - int(Inches(0.18)), row_h)
            tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            tb.text_frame.word_wrap = True
            self._set_rich_text(tb.text_frame.paragraphs[0], item.get("text", ""),
                                SZ_H3, self.MUTED if done else self.FG)

    def build_annotation(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        fig_w = int(rwidth * 0.56)
        notes_w = rwidth - fig_w - int(Inches(0.35))
        notes_x = rleft + fig_w + int(Inches(0.35))
        img_file = self._image_or_placeholder(sd.annotation_figure) if sd.annotation_figure else None
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(fig_w * 0.98); max_h = int(rheight * 0.96)
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96); ph = int(ih * scale * 914400 / 96)
            slide.shapes.add_picture(img_file, rleft, int(rtop + max(0, (rheight - ph) // 2)), pw, ph)
        if sd.annotation_notes:
            nh = int(self._estimate_text_height([f"\u2022 {x}" for x in sd.annotation_notes],
                                               SZ_COL, width=notes_w, gap=PARA_GAP))
            ntop = rtop + max(0, (rheight - min(nh, rheight)) // 2)
            tb = self._add_textbox(slide, notes_x, int(ntop), notes_w, int(min(nh, rheight)))
            tf = tb.text_frame; tf.word_wrap = True
            for i, note in enumerate(sd.annotation_notes):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = PARA_GAP
                self._set_rich_text(p, f"\u2022 {note}", SZ_COL, self.FG)

    def build_before_after(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        mid = int(Inches(0.8))
        half_w = (rwidth - mid) // 2
        box_h = int(min(Inches(4.2), rheight))
        top = rtop + max(0, (rheight - box_h) // 2)
        for idx, (data, label_default, fill, bar) in enumerate([
            (sd.ba_before, "Before", self.SURFACE, False),
            (sd.ba_after, "After", self._tint(self.ACCENT, 0.9), True),
        ]):
            x = rleft + idx * (half_w + mid)
            label = data.get("label", label_default) if data else label_default
            body = data.get("body", "") if data else ""
            self._add_zone_box(slide, int(x), top, int(half_w), box_h,
                               label=label, body=body, fill=fill, accent_bar=bar)
        arrow_x = int(rleft + half_w + (mid - Inches(0.55)) // 2)
        arrow_y = int(top + box_h // 2 - Pt(16))
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, arrow_x, arrow_y, int(Inches(0.55)), int(Pt(32)))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = self.ACCENT
        arrow.line.fill.background(); self._no_shadow(arrow)

    def build_funnel(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.funnel_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        max_w = rwidth
        row_h = int(min(Inches(0.95), rheight / max(n, 1)))
        block_h = row_h * n
        top = rtop + max(0, (rheight - block_h) // 2)
        for i, item in enumerate(sd.funnel_items):
            frac = 1 - (i / max(n - 1, 1)) * 0.55
            w = int(max_w * frac)
            x = int(rleft + (max_w - w) / 2)
            y = int(top + row_h * i)
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, int(row_h * 0.85))
            bg.adjustments[0] = 0.05
            bg.fill.solid()
            t = i / max(n - 1, 1)
            r = int(self.PRIMARY.red * (1 - t) + self.ACCENT.red * t) if hasattr(self.PRIMARY, 'red') else 0x16
            g = int(self.PRIMARY.green * (1 - t) + self.ACCENT.green * t) if hasattr(self.PRIMARY, 'green') else 0x21
            b = int(self.PRIMARY.blue * (1 - t) + self.ACCENT.blue * t) if hasattr(self.PRIMARY, 'blue') else 0x3e
            bg.fill.fore_color.rgb = RGBColor(min(r, 255), min(g, 255), min(b, 255))
            bg.line.fill.background()
            tb = self._add_textbox(slide, x + Pt(16), y + Pt(4), w - Pt(32), int(row_h * 0.85) - Pt(8))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_l = tf.paragraphs[0].add_run()
            run_l.text = item.get("label", "")
            run_l.font.name = self.FONT_HEAD
            run_l.font.size = self._fs(SZ_ZONE_B)
            run_l.font.bold = True
            run_l.font.color.rgb = self.WHITE
            if item.get("value"):
                run_v = tf.paragraphs[0].add_run()
                run_v.text = f"  {item['value']}"
                run_v.font.name = self.FONT
                run_v.font.size = self._fs(SZ_SMALL)
                run_v.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)

    def build_stack(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.stack_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        gap = int(Inches(0.12))
        # Rows are as tall as their content needs, within what the region can
        # give — a fixed 1.0in cap clipped every two-line layer description.
        room = int((rheight - gap * (n - 1)) / max(n, 1))
        inner_w = rwidth - int(CARD_PAD) * 2
        need = int(Inches(0.6))
        for item in sd.stack_items:
            h = int(self._estimate_text_height([item.get("desc", "")], SZ_ZONE_B,
                                               width=inner_w)) if item.get("desc") else 0
            if item.get("name"):
                h += int(Pt(self._fs(SZ_ZONE_L).pt * metrics.DEFAULT_LINE_FACTOR)) + int(Pt(10))
            need = max(need, h + int(CARD_PAD) * 2)
        row_h = int(min(room, max(int(Inches(0.6)), need)))
        block_h = row_h * n + gap * (n - 1)
        top = rtop + max(0, (rheight - block_h) // 2)
        for i, item in enumerate(sd.stack_items):
            y = int(top + (row_h + gap) * (n - 1 - i))
            # deeper layers get a slightly stronger surface for depth
            shade = self._tint(self.SECONDARY, 0.93 - 0.04 * i)
            self._add_zone_box(slide, rleft, y, rwidth, int(row_h), fill=shade,
                              label=item.get("name", ""), body=item.get("desc", ""),
                              anchor="middle")

    def build_card_grid(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.card_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        cols = 2 if n <= 4 else 3
        rows_n = (n + cols - 1) // cols
        gap = int(CARD_GAP)
        card_w = (rwidth - gap * (cols - 1)) // cols
        card_h = int(min(Inches(2.4), (rheight - gap * (rows_n - 1)) / max(rows_n, 1)))
        block_h = card_h * rows_n + gap * (rows_n - 1)
        top = rtop + max(0, (rheight - block_h) // 2)
        for idx, item in enumerate(sd.card_items):
            r, c = divmod(idx, cols)
            x = rleft + c * (card_w + gap)
            y = top + r * (card_h + gap)
            self._add_zone_box(slide, int(x), int(y), int(card_w), int(card_h),
                              label=item.get("title", ""), body=item.get("body", ""))

    def build_split_text(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        gap = int(CARD_GAP)
        half_w = (rwidth - gap) // 2
        for idx, data in enumerate([sd.split_left, sd.split_right]):
            x = rleft + idx * (half_w + gap)
            label = data.get("label", "") if data else ""
            body = data.get("body", "") if data else ""
            self._add_zone_box(slide, int(x), rtop, int(half_w), rheight,
                               label=label, body=body, anchor="middle")

    # beamer block conventions: structure color for theorem-likes (follows the
    # theme), green for examples, red for alerts.
    _BLOCK_COLORS = {
        "example": RGBColor(0x14, 0x65, 0x3D),
        "alert":   RGBColor(0x9F, 0x1D, 0x1D),
    }

    def build_blocks(self, sd: SlideData):
        """beamer-style theorem environments: a colored title bar + tinted
        body panel per block (theorem / definition / lemma / example / alert)."""
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.block_items:
            return
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        bar_h = int(Inches(0.34))
        gap = int(Inches(0.22))
        items = []
        for (variant, title, body) in sd.block_items:
            bh = int(self._estimate_text_height([body], SZ_COL,
                                                width=rwidth - int(Inches(0.5))))
            items.append([variant, title, body,
                          max(bh + int(Inches(0.14)), int(Inches(0.5)))])
        total = sum(bar_h + it[3] for it in items) + gap * (len(items) - 1)
        if total > rheight:                      # squeeze instead of overflowing
            over = total - rheight
            shrinkable = sum(it[3] for it in items)
            for it in items:
                it[3] = max(int(Inches(0.4)), int(it[3] - over * it[3] / shrinkable))
            total = sum(bar_h + it[3] for it in items) + gap * (len(items) - 1)
        cur = rtop + max(0, (rheight - total) // 2)
        for (variant, title, body, body_h) in items:
            head_c = self._BLOCK_COLORS.get(variant, self.SECONDARY)
            # sharp rectangles, flush bar+panel — the (non-rounded) beamer block
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         rleft, int(cur), rwidth, bar_h)
            bar.fill.solid(); bar.fill.fore_color.rgb = head_c
            bar.line.fill.background(); self._no_shadow(bar)
            tb = self._add_textbox(slide, rleft + int(Inches(0.18)), int(cur),
                                   rwidth - int(Inches(0.36)), bar_h)
            tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tb.text_frame.paragraphs[0]
            # A theorem title carries math — "定理 1（$O(1/k)$ 収束）" — so it
            # goes through the inline renderer like any other line.
            self._set_rich_text(p, title or variant.capitalize(), Pt(15), self.WHITE)
            for r in p.runs:
                r.font.bold = True
                r.font.name = self.FONT_HEAD
            cur += bar_h
            panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           rleft, int(cur), rwidth, int(body_h))
            panel.fill.solid(); panel.fill.fore_color.rgb = self._tint(head_c, 0.93)
            panel.line.fill.background(); self._no_shadow(panel)
            btb = self._add_textbox(slide, rleft + int(Inches(0.22)),
                                    int(cur + Inches(0.05)),
                                    rwidth - int(Inches(0.44)),
                                    int(body_h - Inches(0.08)))
            btb.text_frame.word_wrap = True
            btb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            self._set_text_with_inline_math(btb.text_frame.paragraphs[0], body,
                                            SZ_COL, self.FG)
            cur += body_h + gap
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    # tmu-cs [!step] actions → (tick bar, line band) on the dark code card.
    # Hues follow the card's Catppuccin-leaning palette so bands read as
    # state, not noise.
    _STEP_STYLE = {
        "highlight": (RGBColor(0xA6, 0xE3, 0xA1), RGBColor(0x2E, 0x3E, 0x30)),
        "focus":     (RGBColor(0x89, 0xB4, 0xFA), RGBColor(0x2B, 0x33, 0x4A)),
        "warning":   (RGBColor(0xF9, 0xE2, 0xAF), RGBColor(0x45, 0x3E, 0x28)),
        "error":     (RGBColor(0xF3, 0x8B, 0xA8), RGBColor(0x46, 0x2A, 0x34)),
        "info":      (RGBColor(0x74, 0xC7, 0xEC), RGBColor(0x28, 0x3C, 0x4A)),
    }

    def build_code(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        desc_h = int(Inches(0.7)) if sd.code_desc else 0
        code_lines = (sd.code_text or "").split("\n")
        n_lines = len(code_lines)
        CODE_PT = 13
        scale = getattr(self.theme, "font_scale", 1.0)
        line_h = int(Pt(CODE_PT * 1.5 * scale))   # fixed line pitch so the
        bar_h = int(Inches(0.34))                 # step bands can be placed
        pad_top = int(Pt(6)); pad_bot = int(Pt(12))
        avail = rheight - desc_h - (int(Inches(0.2)) if desc_h else 0)
        # More lines than the region can hold: tighten the pitch (and the type
        # with it) instead of clipping the tail of the listing.
        room = avail - bar_h - pad_top - pad_bot
        if n_lines and line_h * n_lines > room:
            line_h = max(int(Pt(9)), int(room / n_lines))
            CODE_PT = max(7.0, line_h / 12700.0 / 1.5 / max(scale, 0.01))
        code_h = int(min(avail,
                         max(int(Inches(1.5)), bar_h + pad_top + line_h * n_lines + pad_bot)))
        block_h = code_h + (desc_h + int(Inches(0.2)) if desc_h else 0)
        top = rtop + max(0, (rheight - block_h) // 2)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rleft, top, rwidth, code_h)
        bg.adjustments[0] = 0.02
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
        bg.line.fill.background(); self._no_shadow(bg)
        # window dots
        for j, col in enumerate([RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFE, 0xBC, 0x2E),
                                 RGBColor(0x28, 0xC8, 0x40)]):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         rleft + int(Pt(16)) + j * int(Pt(18)),
                                         top + int(Pt(11)), int(Pt(10)), int(Pt(10)))
            dot.fill.solid(); dot.fill.fore_color.rgb = col
            dot.line.fill.background(); self._no_shadow(dot)
        # Step emphasis: bands go in before the text box so they sit behind it.
        stepping = bool(sd.code_steps) and sd.active_step is not None
        active: dict[int, str] = {}
        if stepping:
            for (idx, step, action, span) in sd.code_steps:
                if step == sd.active_step:
                    for k in range(span):
                        if idx + k < n_lines:
                            active[idx + k] = action
        text_top = top + bar_h + pad_top
        for li, action in sorted(active.items()):
            tick_c, band_c = self._STEP_STYLE.get(action, self._STEP_STYLE["highlight"])
            band = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, rleft + int(Pt(8)),
                int(text_top + li * line_h), rwidth - int(Pt(16)), line_h)
            band.fill.solid(); band.fill.fore_color.rgb = band_c
            band.line.fill.background(); self._no_shadow(band)
            tick = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, rleft + int(Pt(8)),
                int(text_top + li * line_h), int(Pt(3.5)), line_h)
            tick.fill.solid(); tick.fill.fore_color.rgb = tick_c
            tick.line.fill.background(); self._no_shadow(tick)
        if stepping:
            last = max(s for (_, s, _, _) in sd.code_steps)
            sb = self._add_textbox(slide, rleft + rwidth - int(Inches(1.7)),
                                   top + int(Pt(9)), int(Inches(1.5)), int(Pt(15)))
            sp = sb.text_frame.paragraphs[0]
            sp.text = f"STEP {sd.active_step} / {last}"
            sp.alignment = PP_ALIGN.RIGHT
            sp.font.name = self.FONT_HEAD; sp.font.size = self._fs(Pt(10))
            sp.font.bold = True
            sp.font.color.rgb = RGBColor(0x9C, 0xA0, 0xB0)
        tb = self._add_textbox(slide, rleft + int(Pt(20)), int(text_top),
                               rwidth - int(Pt(40)), code_h - bar_h - pad_top - pad_bot)
        tf = tb.text_frame
        tf.word_wrap = False   # line index must map 1:1 to band position
        # Dimmed-but-legible: 4.9:1 on the panel (was 3.4:1 — doctor's
        # contrast check is the deck's own bar, de-emphasis included).
        dim = RGBColor(0x88, 0x8C, 0xA0)
        normal = RGBColor(0xCD, 0xD6, 0xF4)
        for i, ln in enumerate(code_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ln if ln.strip() else " "   # keep blank lines tall
            p.font.name = self.FONT_MONO; p.font.size = self._fs(Pt(CODE_PT))
            p.font.color.rgb = normal if (not stepping or i in active) else dim
            p.line_spacing = Pt(CODE_PT * 1.5 * scale)
        if sd.code_desc:
            dtb = self._add_textbox(slide, rleft, int(top + code_h + Inches(0.2)),
                                    rwidth, desc_h)
            dtb.text_frame.word_wrap = True
            self._rich_line(dtb.text_frame.paragraphs[0], sd.code_desc,
                            SZ_SMALL, self.MUTED)

    def build_multi_result(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.multi_result_items)
        if n == 0:
            return
        rleft, rtop, rwidth, rheight = self._content_region_with_lead(slide, sd)
        gap = int(CARD_GAP)
        box_w = (rwidth - gap * (n - 1)) // n
        box_h = int(min(Inches(3.4), rheight))
        top = rtop + max(0, (rheight - box_h) // 2)
        for i, item in enumerate(sd.multi_result_items):
            x = rleft + i * (box_w + gap)
            bg, tb = self._add_card(slide, (int(x), top, int(box_w), box_h),
                                    accent_bar=True, anchor="middle")
            tf = tb.text_frame
            tf.clear()
            pm = tf.paragraphs[0]
            self._kicker_para(pm, item.get("metric", ""), size=SZ_SMALL,
                              color=self.MUTED, align=PP_ALIGN.CENTER, tracking=120)
            pv = tf.add_paragraph(); pv.space_before = Pt(6)
            rv = pv.add_run(); rv.text = item.get("value", "")
            rv.font.name = self.FONT_HEAD; rv.font.size = self._fs(Pt(38))
            rv.font.bold = True; rv.font.color.rgb = self.ACCENT
            pv.alignment = PP_ALIGN.CENTER
            if item.get("desc"):
                pd = tf.add_paragraph(); pd.space_before = Pt(8)
                self._set_rich_text(pd, item.get("desc", ""), SZ_ZONE_B, self.FG)
                pd.alignment = PP_ALIGN.CENTER

    def build_takeaway(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        main_w = rwidth - int(Inches(1.0))
        pts_w = rwidth - int(Inches(2.8))
        main_h = int(self._estimate_text_height([sd.takeaway_main], Pt(28), width=main_w)) if sd.takeaway_main else 0
        pts_h = int(self._estimate_text_height([f"\u2022 {p}" for p in sd.takeaway_points], SZ_H3, width=pts_w)) if sd.takeaway_points else 0
        gap = int(Inches(0.4)) if (main_h and pts_h) else 0
        block_h = main_h + gap + pts_h
        top = rtop + max(0, (rheight - block_h) // 2)
        cx = int(SW // 2)
        if sd.takeaway_main:
            self._hairline(slide, cx - int(Inches(0.6)), top - int(Inches(0.22)),
                           Inches(1.2), thickness=ACCENT_RULE_W, color=self.ACCENT)
            tb = self._add_textbox(slide, rleft + int(Inches(0.5)), top,
                                   main_w, main_h)
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            self._set_rich_text(p, sd.takeaway_main, Pt(28), self.PRIMARY)
            for r in p.runs:
                r.font.bold = True; r.font.name = self.FONT_HEAD
            p.alignment = PP_ALIGN.CENTER; p.line_spacing = LINE_TITLE
        if sd.takeaway_points:
            ptb = self._add_textbox(slide, rleft + int(Inches(1.4)), top + main_h + gap,
                                    pts_w, pts_h)
            tf = ptb.text_frame; tf.word_wrap = True
            for i, pt in enumerate(sd.takeaway_points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = PARA_GAP
                self._set_rich_text(p, f"\u2022  {pt}", SZ_H3, self.FG)

    def build_profile(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        rleft, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1))
        left_w = int(Inches(3.2))
        right_x = rleft + left_w + int(Inches(0.5))
        right_w = rwidth - left_w - int(Inches(0.5))
        img_file = self._image_or_placeholder(sd.image_path) if sd.image_path else None
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(left_w * 0.95); max_h = int(rheight * 0.92)
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96); ph = int(ih * scale * 914400 / 96)
            iy = rtop + max(0, (rheight - ph) // 2)
            slide.shapes.add_picture(img_file, rleft, int(iy), pw, ph)
        # Right column: name / affiliation / bio, vertically centered as a block
        name_h = int(Pt(self._fs(Pt(24)).pt * metrics.DEFAULT_LINE_FACTOR)
                     + Pt(6)) if sd.profile_name else 0
        affil_h = int(Inches(0.4)) if sd.profile_affiliation else 0
        # width-aware: long bio bullets wrap in the right column; a 1-line-each
        # estimate under-sized the block, mis-centering it (pushed low / could
        # run off the bottom).
        bio_h = int(self._estimate_text_height(["\u2022 " + b for b in sd.profile_bio],
                                               SZ_ZONE_B, width=right_w,
                                               gap=PARA_GAP)) if sd.profile_bio else 0
        block_h = name_h + affil_h + (int(Inches(0.2)) if bio_h else 0) + bio_h
        cur_y = rtop + max(0, (rheight - block_h) // 2)
        if sd.profile_name:
            ntb = self._add_textbox(slide, int(right_x), int(cur_y), int(right_w), name_h)
            p = ntb.text_frame.paragraphs[0]; p.text = sd.profile_name
            p.font.name = self.FONT_HEAD; p.font.size = self._fs(Pt(24))
            p.font.bold = True; p.font.color.rgb = self.PRIMARY
            cur_y += name_h
        if sd.profile_affiliation:
            atb = self._add_textbox(slide, int(right_x), int(cur_y), int(right_w), affil_h)
            p = atb.text_frame.paragraphs[0]
            self._kicker_para(p, sd.profile_affiliation, size=SZ_SMALL,
                              color=self.ACCENT_TEXT, tracking=120)
            cur_y += affil_h + int(Inches(0.2))
        if sd.profile_bio:
            btb = self._add_textbox(slide, int(right_x), int(cur_y), int(right_w), bio_h)
            btf = btb.text_frame; btf.word_wrap = True
            for i, item in enumerate(sd.profile_bio):
                p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                p.space_before = PARA_GAP
                self._set_rich_text(p, "\u2022  " + item, SZ_ZONE_B, self.FG)

    # ── Okabe-Ito color-vision-safe categorical palette ──
    OKABE_ITO = [
        RGBColor(0xE6, 0x9F, 0x00), RGBColor(0x56, 0xB4, 0xE9),
        RGBColor(0x00, 0x9E, 0x73), RGBColor(0xF0, 0xE4, 0x42),
        RGBColor(0x00, 0x72, 0xB2), RGBColor(0xD5, 0x5E, 0x00),
        RGBColor(0xCC, 0x79, 0xA7), RGBColor(0x00, 0x00, 0x00),
    ]

    def _chart_colors(self, n: int) -> list[RGBColor]:
        """Series colors drawn from the theme, accent first (dominance).

        A foreign categorical set (Okabe-Ito) reads as a different design
        pasted into the deck; up to four series the theme itself can supply
        distinguishable colors. The luminance-ratio gate is measured, not
        assumed — themes whose ramp collapses (accent ≈ secondary) fall back
        to Okabe-Ito, which also stays the palette beyond four series.
        """
        if n <= 1:
            return [self.ACCENT]
        # A palette that declares accent2..accent6 has hand-picked its
        # categorical ramp (e.g. research's MATLAB-style 6-color system) —
        # use it verbatim, in order.
        explicit = [self.ACCENT] + [
            c for c in (getattr(self.theme, f"accent{i}", None) for i in range(2, 7))
            if c is not None]
        if len(explicit) >= 2 and n <= len(explicit):
            return explicit[:n]
        from marp_pptx.audit import contrast_ratio
        ramp = [self.ACCENT, self.PRIMARY,
                self._tint(self.ACCENT, 0.55), self.SECONDARY]
        if n <= len(ramp):
            cols = ramp[:n]
            distinct = all(
                contrast_ratio(tuple(a), tuple(b)) >= 1.3
                for i, a in enumerate(cols) for b in cols[i + 1:])
            if distinct:
                return cols
        return [self.OKABE_ITO[i % len(self.OKABE_ITO)] for i in range(n)]

    def build_statement(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.dark:
            self._set_bg(slide, self.PRIMARY)
        fg = self.WHITE if sd.dark else self.PRIMARY
        accent = self.WHITE if sd.dark else self.ACCENT
        text = sd.statement_text or sd.h1 or ""
        cx = int(SW // 2)
        tw = int(SW - Inches(2.4))
        # Size from the *rendered* width, not the character count: 40 narrow
        # Latin characters and 40 full-width ones differ by more than 2x, and
        # the old count-based rule put both at 44pt — where the second wrapped
        # to two lines in a box reserved for one.
        pt = 44
        for cand in (44, 36, 28):
            pt = cand
            if metrics.line_count(self._plain(text), self.FONT_HEAD, cand,
                                  tw / 12700.0, ea_font=self.FONT_EA,
                                  bold=True) <= 2:
                break
        th = int(self._estimate_text_height([text], Pt(pt), width=tw))
        top = int(MARGIN_T) + max(0, (int(SH - 2 * MARGIN_T) - th - int(Inches(0.3))) // 2)
        self._hairline(slide, cx - int(Inches(0.5)), top - int(Inches(0.28)),
                       Inches(1.0), thickness=ACCENT_RULE_W, color=accent)
        tb = self._add_textbox(slide, int(Inches(1.2)), top, tw,
                               th + int(Inches(0.4)))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        self._set_rich_text(p, text, Pt(pt), fg)
        for r in p.runs:
            r.font.bold = True; r.font.name = self.FONT_HEAD
        p.alignment = PP_ALIGN.CENTER; p.line_spacing = LINE_TITLE

    def build_big_number(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.dark:
            self._set_bg(slide, self.PRIMARY)
        if sd.h1:
            self._add_title(slide, sd.h1)
        _, rtop, rwidth, rheight = self._content_region(has_title=bool(sd.h1),
                                                        full=not sd.footnote)
        cx = int(SW // 2)
        big = sd.bignum_value or "—"
        num_pt = 110 if len(big) <= 4 else (88 if len(big) <= 7 else 64)
        num_h = int(self._fs(Pt(num_pt)) * metrics.DEFAULT_LINE_FACTOR)
        lbl_h = int(Inches(0.5)) if sd.bignum_label else 0
        cap_h = int(Inches(0.45)) if sd.bignum_caption else 0
        block = num_h + lbl_h + cap_h
        top = rtop + max(0, (rheight - block) // 2)
        tb = self._add_textbox(slide, int(MARGIN_L), top, int(CONTENT_W), num_h)
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tb.text_frame.paragraphs[0]; p.text = big
        p.font.name = self.FONT_HEAD; p.font.size = self._fs(Pt(num_pt))
        p.font.bold = True; p.font.color.rgb = self.ACCENT
        p.alignment = PP_ALIGN.CENTER
        self._hairline(slide, cx - int(Inches(0.5)), top + num_h + int(Inches(0.02)),
                       Inches(1.0), thickness=ACCENT_RULE_W, color=self.ACCENT)
        cur = top + num_h + int(Inches(0.16))
        if sd.bignum_label:
            lb = self._add_textbox(slide, int(MARGIN_L), cur, int(CONTENT_W), lbl_h)
            self._kicker_para(lb.text_frame.paragraphs[0], sd.bignum_label,
                              size=SZ_H3, color=(self.WHITE if sd.dark else self.SECONDARY),
                              align=PP_ALIGN.CENTER, tracking=140)
            cur += lbl_h
        if sd.bignum_caption:
            cb = self._add_textbox(slide, int(MARGIN_L), cur, int(CONTENT_W), cap_h)
            cp = cb.text_frame.paragraphs[0]
            self._rich_line(cp, sd.bignum_caption, SZ_SMALL,
                            RGBColor(0xD8, 0xD8, 0xDE) if sd.dark else self.MUTED)
            cp.alignment = PP_ALIGN.CENTER
            cb.text_frame.word_wrap = True
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_chart(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        left, rtop, width, rheight = self._content_region_with_lead(slide, sd)
        if not sd.chart_series or not sd.chart_categories:
            self._add_body_text(slide, ["(chart: データ未指定 — Markdown表で系列を渡してください)"],
                                left=left, top=rtop, width=width)
            return
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        kind = {"line": XL_CHART_TYPE.LINE_MARKERS,
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED}.get(
                    sd.chart_kind, XL_CHART_TYPE.COLUMN_CLUSTERED)
        data = CategoryChartData()
        data.categories = sd.chart_categories
        for name, vals in sd.chart_series:
            data.add_series(name, vals)
        cap_text = self._fig_caption(sd.chart_caption, sd.source) or sd.footnote
        cap_h = (max(int(Inches(0.4)),
                     int(self._estimate_text_height([cap_text], SZ_SMALL, width=width)))
                 if cap_text else 0)
        ch_h = rheight - cap_h
        gf = slide.shapes.add_chart(kind, left, rtop, width, ch_h, data)
        chart = gf.chart
        chart.has_title = False
        n_series = len(sd.chart_series)
        chart.has_legend = n_series > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = self._fs(SZ_SMALL)
        series_cols = self._chart_colors(n_series)
        for i, plot_series in enumerate(chart.series):
            col = series_cols[i]
            try:
                fmt = plot_series.format
                fmt.fill.solid(); fmt.fill.fore_color.rgb = col
                fmt.line.color.rgb = col
            except Exception:
                pass
        try:
            for ax in (chart.category_axis, chart.value_axis):
                ax.tick_labels.font.size = self._fs(SZ_SMALL)
                ax.format.line.color.rgb = self.HAIRLINE
        except Exception:
            pass
        if cap_text:
            cb = self._add_textbox(slide, left, int(rtop + ch_h + Inches(0.05)),
                                   width, cap_h)
            cb.text_frame.word_wrap = True
            cp = cb.text_frame.paragraphs[0]
            self._rich_line(cp, cap_text, SZ_SMALL, self.MUTED)
            cp.alignment = PP_ALIGN.CENTER

    # ══════════════════════════════════════════════
    # Build all slides
    # ══════════════════════════════════════════════
    BUILDERS = {
        "title": "build_title",
        "divider": "build_divider",
        "cols-2": "build_columns",
        "cols-2-wide-l": "build_columns",
        "cols-2-wide-r": "build_columns",
        "cols-3": "build_columns",
        "sandwich": "build_sandwich",
        "equation": "build_equation",
        "equations": "build_equations",
        "figure": "build_figure",
        "table-slide": "build_table",
        "references": "build_references",
        "timeline-h": "build_timeline_h",
        "timeline": "build_timeline_v",
        "end": "build_end",
        "zone-flow": "build_zone_flow",
        "zone-compare": "build_zone_compare",
        "zone-matrix": "build_zone_matrix",
        "zone-process": "build_zone_process",
        "agenda": "build_agenda",
        "sections": "build_sections",
        "rq": "build_rq",
        "result-dual": "build_result_dual",
        "summary": "build_summary",
        "appendix": "build_appendix",
        "overview": "build_overview",
        "result": "build_result",
        "steps": "build_steps",
        "quote": "build_quote",
        "history": "build_history",
        "panorama": "build_panorama",
        "kpi": "build_kpi",
        "pros-cons": "build_pros_cons",
        "definition": "build_definition",
        "diagram": "build_diagram",
        "gallery-img": "build_gallery_img",
        "highlight": "build_highlight",
        "checklist": "build_checklist",
        "annotation": "build_annotation",
        "before-after": "build_before_after",
        "funnel": "build_funnel",
        "stack": "build_stack",
        "card-grid": "build_card_grid",
        "split-text": "build_split_text",
        "code": "build_code",
        "multi-result": "build_multi_result",
        "takeaway": "build_takeaway",
        "profile": "build_profile",
        "statement": "build_statement",
        "big-statement": "build_statement",
        "dark": "build_statement",
        "big-number": "build_big_number",
        "big-number-dark": "build_big_number",
        "chart": "build_chart",
        "blocks": "build_blocks",
    }

    # Generic noun-label titles that should be assertions instead.
    _LABEL_TITLES = {
        "結果", "考察", "背景", "方法", "手法", "まとめ", "結論", "議論", "目的",
        "序論", "概要", "課題", "提案手法", "実験", "実験結果", "今後の課題",
        "results", "result", "method", "methods", "methodology", "conclusion",
        "conclusions", "background", "introduction", "discussion", "summary",
        "objective", "overview", "agenda", "outline", "motivation", "approach",
        "experiments", "evaluation", "future work", "related work",
    }
    # density -> (max_elements, max_body_lines, max_line_chars, max_words)
    _LINT_LIMITS = {
        "academic": (6, 6, 42, 75),
        "keynote":  (4, 3, 30, 40),
        "dense":    (9, 10, 56, 130),   # hearing-deck density: organized, not sparse
    }

    def _lint_slide(self, sd, idx: int):
        """Emit stderr warnings for over-stuffed slides / weak titles. Advisory
        only (never blocks). Grounded in PLOS '≤6 elements', 6x6, action-titles."""
        import sys as _sys
        density = getattr(self.theme, "density", "academic")
        max_el, max_lines, max_chars, max_words = self._LINT_LIMITS.get(
            density, self._LINT_LIMITS["academic"])
        warn = lambda m: print(f"[lint] slide {idx}: {m}", file=_sys.stderr)

        def cjk_len(s):  # full-width chars count as 2
            return sum(2 if ord(c) > 0x2E7F else 1 for c in s)

        body = [l.strip() for l in (sd.body_lines or []) if l.strip()
                and not l.strip().startswith(("#", "<"))]
        if len(body) > max_lines:
            warn(f"{len(body)} body lines (>{max_lines}) — split or trim")
        for l in body:
            if cjk_len(l) > max_chars * 2:
                warn(f"a line is long (~{cjk_len(l)//2} chars >{max_chars}) — shorten")
                break
        words = sum(len(l.split()) for l in body)
        if words > max_words:
            warn(f"~{words} words (>{max_words}) — this reads as a document, not a slide")

        # element count (max items across the type's collection + body + media)
        item_fields = ("columns", "kpi_items", "zone_flow_items", "zone_process_items",
                       "steps_items", "card_items", "timeline_items", "agenda_items",
                       "summary_points", "checklist_items", "funnel_items", "stack_items",
                       "multi_result_items", "gallery_items", "history_items",
                       "eq_system", "ref_items", "annotation_notes")
        items = max([len(getattr(sd, f, []) or []) for f in item_fields] + [0])
        elements = max(items, len(body)) + (1 if sd.image_path else 0) + (1 if sd.table_rows else 0)
        if elements > max_el:
            warn(f"{elements} elements (>{max_el}) — one slide, one message")

        # assertion (action-title) check
        h = (sd.h2 or sd.h1 or "").strip()
        h = re.sub(r"^\[[^\]]*\]\s*", "", h)   # drop a "[KPI]" style prefix
        if h and h.strip().rstrip("。.:：").lower() in self._LABEL_TITLES:
            warn(f'title "{h}" is a label — state the takeaway (e.g. add the result/number)')

    def build_all(self, slides: list[SlideData]):
        self._divider_no = 0
        self._deck_title = ""
        self._fig_no = 0
        self._slide_sections: list[str] = []   # per-built-slide section name
        self._slide_classes: list[str] = []    # per-built-slide semantic type
        current_section = ""
        for n, sd in enumerate(slides, 1):
            try:
                self._lint_slide(sd, n)
            except Exception:
                pass
            if sd.slide_class == "divider":
                self._divider_no += 1
                current_section = re.sub(r"^\s*\d+[\.．]?\s*", "", sd.h1 or "")
            if sd.slide_class == "title" and not self._deck_title:
                self._deck_title = sd.h1 or ""
            # A specified-but-unknown _class silently fell back to a plain
            # bullet slide before — warn so typos in the type name surface.
            if sd.slide_class and sd.slide_class not in self.BUILDERS:
                self._warn(f"slide {n}: unknown type '{sd.slide_class}', "
                           "rendered as a plain slide (check the _class name)")
            method_name = self.BUILDERS.get(sd.slide_class, "build_default")
            before_n = len(self.prs.slides)
            getattr(self, method_name)(sd)
            # Embed the slide class (+ any speaker note) into the notes of the
            # slide(s) just created so pptx2md can recover the semantic type.
            cls = sd.slide_class or "default"
            user_note = getattr(sd, "speaker_notes", "") or ""
            for i in range(before_n, len(self.prs.slides)):
                self._write_notes(self.prs.slides[i], cls,
                                  user_note if i == before_n else "")
                self._slide_sections.append(current_section)
                self._slide_classes.append(cls)
        self._add_global_footer()
        self._warn_text_collisions()

    def _warn_text_collisions(self):
        """Post-build geometric self-check: flag any text that overflows its
        box onto the element below it (the wrap-overlap bug class). Surfaces
        through self.warnings → CLI summary / MCP authoring_warnings."""
        try:
            from marp_pptx.visuallint import detect_text_collisions
            for c in detect_text_collisions(self.prs):
                where = (c["onto"] if c["onto"].startswith("(")
                         else f'onto "{c["onto"][:24]}"')
                self._warn(f"slide {c['slide']}: text overflows {where} "
                           f"— \"{c['over'][:24]}\" (box too short for its content)")
        except Exception:
            pass  # a self-check must never break a build

    def _write_notes(self, slide, cls: str, user_note: str = ""):
        """Write speaker notes: user note first (what the presenter reads),
        then the machine-readable `_class:` tag last (pptx2md round-trip)."""
        try:
            tf = slide.notes_slide.notes_text_frame
            existing = tf.text or ""
            tag = f"_class: {cls}"
            body = "\n".join(l for l in existing.splitlines()
                             if not l.strip().startswith("_class:")).strip()
            parts = []
            if body:
                parts.append(body)
            if user_note and user_note not in body:
                parts.append(user_note)
            note_text = "\n\n".join(p for p in parts if p.strip())
            tf.text = (note_text + "\n\n" + tag) if note_text else tag
        except Exception:
            pass

    # Backward-compatible alias
    def _write_class_note(self, slide, cls: str):
        self._write_notes(slide, cls, "")

    def _fig_caption(self, caption: str, source: str = "") -> str:
        """Hearing-deck caption convention (opt-in via figure_numbers):
        "図 N｜caption　出典｜source". Call ONCE per figure — the counter
        advances on every numbered caption."""
        text = (caption or "").strip()
        if text and self.LAYOUT.figure_numbers:
            self._fig_no = getattr(self, "_fig_no", 0) + 1
            text = f"図 {self._fig_no}｜{text}"
        if source:
            # A bare source line stays an attribution, not a numbered figure.
            text = (text + ("　" if text else "") + f"出典｜{source}").strip()
        return text

    def _add_header_crumb(self):
        """Hearing-style top chrome: section breadcrumb left, n／m right.

        Replaces the bottom-right page number (both at once is noise). Hero
        slides (title/divider/end/statement) keep their clean canvas.
        """
        n = len(self.prs.slides)
        sections = getattr(self, "_slide_sections", [])
        classes = getattr(self, "_slide_classes", [])
        skip = {"title", "divider", "end", "statement", "dark", "big-statement"}
        for i, slide in enumerate(self.prs.slides):
            cls = classes[i] if i < len(classes) else ""
            if i == 0 or cls in skip:
                continue
            crumb = sections[i] if i < len(sections) else ""
            if crumb:
                tb = self._add_textbox(slide, int(MARGIN_L), int(Inches(0.10)),
                                       int(CONTENT_W * 0.7), int(Inches(0.22)))
                p = tb.text_frame.paragraphs[0]
                p.text = crumb
                p.font.name = self.FONT
                p.font.size = self._fs(SZ_FOOT)
                p.font.color.rgb = self.MUTED
            tb = self._add_textbox(slide, int(SW - MARGIN_R - Inches(1.6)),
                                   int(Inches(0.10)), int(Inches(1.6)),
                                   int(Inches(0.22)))
            p = tb.text_frame.paragraphs[0]
            p.text = f"{i + 1}／{n}"
            p.font.name = self.FONT
            p.font.size = self._fs(SZ_FOOT)
            p.font.color.rgb = self.MUTED
            p.alignment = PP_ALIGN.RIGHT

    def _add_global_footer(self):
        if self.LAYOUT.footer_bar:
            return self._add_beamer_footer()
        if self.LAYOUT.header_crumb:
            return self._add_header_crumb()
        n = len(self.prs.slides)
        for i, slide in enumerate(self.prs.slides):
            if i == 0 or i == n - 1:
                continue
            tb = self._add_textbox(slide, int(MARGIN_L), int(SH - Inches(0.42)),
                                   int(CONTENT_W), int(Inches(0.25)))
            p = tb.text_frame.paragraphs[0]
            p.text = f"{i + 1} / {n}"
            p.font.name = self.FONT
            p.font.size = self._fs(SZ_FOOT)
            p.font.color.rgb = self.MUTED
            p.alignment = PP_ALIGN.RIGHT

    def _add_beamer_footer(self):
        """beamer-style footer bar: deck title | current section | page.

        The middle cell is a slightly lighter tint of the structure color, the
        classic Madrid three-cell look. The title slide keeps its hero layout."""
        n = len(self.prs.slides)
        bar_h = int(Inches(0.30))
        sections = getattr(self, "_slide_sections", [])
        deck_title = getattr(self, "_deck_title", "")
        cell_pts = (0.0, 0.40, 0.74, 1.0)      # cell boundaries as SW fractions
        for i, slide in enumerate(self.prs.slides):
            if i == 0:
                continue
            mid = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(SW * cell_pts[1]), int(SH - bar_h),
                int(SW * (cell_pts[2] - cell_pts[1])), bar_h)
            mid.fill.solid(); mid.fill.fore_color.rgb = self._tint(self.PRIMARY, 0.18)
            mid.line.fill.background(); self._no_shadow(mid)
            for cell, text, align in (
                (0, deck_title, PP_ALIGN.LEFT),
                (1, sections[i] if i < len(sections) else "", PP_ALIGN.CENTER),
                (2, f"{i + 1} / {n}", PP_ALIGN.RIGHT),
            ):
                if cell != 1:   # three tones, darkest at the left (Madrid)
                    bgc = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, int(SW * cell_pts[cell]), int(SH - bar_h),
                        int(SW * (cell_pts[cell + 1] - cell_pts[cell])), bar_h)
                    bgc.fill.solid()
                    bgc.fill.fore_color.rgb = (self.PRIMARY if cell == 0
                                               else self._tint(self.PRIMARY, 0.35))
                    bgc.line.fill.background(); self._no_shadow(bgc)
                tb = self._add_textbox(
                    slide, int(SW * cell_pts[cell] + Inches(0.18)), int(SH - bar_h),
                    int(SW * (cell_pts[cell + 1] - cell_pts[cell]) - Inches(0.36)), bar_h)
                tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tb.text_frame.paragraphs[0]
                p.text = text
                p.alignment = align
                p.font.name = self.FONT
                p.font.size = self._fs(Pt(9.5))
                p.font.color.rgb = self.WHITE
