"""Deterministic deck audit — read a .pptx, report what a reviewer would.

`visuallint` looks at pixels, so it needs LibreOffice and can only judge what
the renderer chose to draw.  This module works on the geometry instead: every
shape's box, every run's font and size, measured against real font metrics
(`marp_pptx.metrics`).  That makes it exact where pixels are approximate —
"the caption overflows its box by 11.4 pt" rather than "content reaches the
slide edge" — and it runs anywhere, in milliseconds, with no renderer at all.

    findings = audit_pptx("deck.pptx")          # list[Finding], worst first
    report   = format_findings(findings)        # human-readable text

Checks, in the order they bite in practice:

    overflow    wrapped text is taller (or wider) than the box holding it
    offslide    a shape sits outside the slide, or inside the safe margin
    overlap     two text shapes cover each other
    contrast    a run fails WCAG AA against the background actually behind it
    font        a typeface that is missing locally or renders untrue in QA
    deck        deck-level: layout monotony, text-only slides, size contrast
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from . import metrics as M

# "7 / 17" in the corner is chrome the builder stamps, not author content:
# measuring it reports the same two findings on every slide of every deck.
_PAGE_NO = re.compile(r"\d+\s*[/／]\s*\d+")

EMU_PT = 12700.0
EMU_IN = 914400.0

# Severity ordering for sorting / exit codes.
_RANK = {"error": 0, "warn": 1, "info": 2}

DEFAULT_SAFE_MARGIN_IN = 0.35   # content closer than this to an edge is cramped
MIN_GAP_IN = 0.12               # text blocks closer than this read as touching
DEFAULT_SIZE_PT = 18.0          # when a run inherits its size from the layout


@dataclass
class Finding:
    kind: str
    severity: str
    slide: int
    message: str
    fix: str = ""
    shape: str = ""
    detail: dict = field(default_factory=dict)

    def __str__(self) -> str:
        where = f"slide {self.slide}" + (f" · {self.shape}" if self.shape else "")
        line = f"[{self.severity}] {where}: {self.message}"
        return line + (f"\n        → {self.fix}" if self.fix else "")


# ── geometry helpers ────────────────────────────────────────────────────────
def _rect(shape):
    """The shape's on-screen bounding box in EMU, rotation included.

    A rotated axis label is stored as a wide flat box that PowerPoint spins
    about its centre; measuring the stored box puts it 1.5" off the slide.
    """
    try:
        x0, y0, w, h = shape.left, shape.top, shape.width, shape.height
    except TypeError:                                   # inherited placeholder
        return None
    if x0 is None or y0 is None or w is None or h is None:
        return None
    rot = getattr(shape, "rotation", 0) or 0
    if rot % 180:
        import math
        a = math.radians(rot)
        cx, cy = x0 + w / 2, y0 + h / 2
        bw = abs(w * math.cos(a)) + abs(h * math.sin(a))
        bh = abs(w * math.sin(a)) + abs(h * math.cos(a))
        return (int(cx - bw / 2), int(cy - bh / 2),
                int(cx + bw / 2), int(cy + bh / 2))
    return (x0, y0, x0 + w, y0 + h)


def _intersection(a, b) -> float:
    """Overlapping area of two EMU rects, in square inches."""
    x0 = max(a[0], b[0])
    y0 = max(a[1], b[1])
    x1 = min(a[2], b[2])
    y1 = min(a[3], b[3])
    if x1 <= x0 or y1 <= y0:
        return 0.0
    return ((x1 - x0) / EMU_IN) * ((y1 - y0) / EMU_IN)


def _label(shape) -> str:
    txt = ""
    if shape.has_text_frame:
        txt = shape.text_frame.text.strip().replace("\n", " ⏎ ")
    elif getattr(shape, "has_table", False) and shape.has_table:
        head = [c.text.strip() for c in shape.table.rows[0].cells]
        txt = "表: " + " / ".join(h for h in head if h)
    if len(txt) > 34:
        txt = txt[:33] + "…"
    return txt or (shape.name or "shape")


# ── text model ──────────────────────────────────────────────────────────────
@dataclass
class _Para:
    text: str
    size: float
    bold: bool
    font: str
    line_spacing: float
    space_before: float
    space_after: float
    ea: str | None = None


def _emu_or_pt(v, default=0.0) -> float:
    if v is None:
        return default
    return v / EMU_PT if hasattr(v, "emu") else float(v)


_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _ea_typeface(run) -> str | None:
    """The `<a:ea>` typeface on a run — what PowerPoint uses for CJK text.

    python-pptx's `font.name` only exposes `<a:latin>`, so a Japanese deck
    would otherwise be measured with a Latin face and read ~40% narrow.
    """
    try:
        rpr = run._r.find(f"{{{_A_NS}}}rPr")
        if rpr is None:
            return None
        ea = rpr.find(f"{{{_A_NS}}}ea")
        return ea.get("typeface") if ea is not None else None
    except Exception:
        return None


def _paragraphs(shape, latin: str, ea: str) -> list[_Para]:
    out: list[_Para] = []
    tf = shape.text_frame
    for p in tf.paragraphs:
        runs = list(p.runs)
        text = "".join(r.text for r in runs) or p.text
        if not text.strip():
            # An empty paragraph still occupies a line.
            size = _emu_or_pt(p.font.size, DEFAULT_SIZE_PT)
            out.append(_Para("", size, False, latin, _spacing(p), 0.0, 0.0))
            continue
        size = None
        bold = False
        font = None
        ea_face = None
        for r in runs:
            if size is None and r.font.size is not None:
                size = r.font.size.pt
            bold = bold or bool(r.font.bold)
            if font is None and r.font.name:
                font = r.font.name
            if ea_face is None:
                ea_face = _ea_typeface(r)
        if size is None:
            size = _emu_or_pt(p.font.size, DEFAULT_SIZE_PT)
        out.append(_Para(text, float(size), bold, font or latin, _spacing(p),
                         _emu_or_pt(p.space_before), _emu_or_pt(p.space_after),
                         ea_face or ea))
    return out


def _spacing(p) -> float:
    """Line spacing as a multiple of the modelled single-line pitch.

    PowerPoint takes `lnSpc` either as a percentage or as an exact height in
    points. An exact height *replaces* the pitch, so it has to be divided by
    the pitch — not by the font size — or a code block set to a tight 1.5x
    exact spacing measures ~25% taller than it renders.
    """
    ls = p.line_spacing
    if ls is None:
        return 1.0
    if hasattr(ls, "pt"):
        size = _emu_or_pt(p.font.size, DEFAULT_SIZE_PT)
        return max(0.5, ls.pt / max(1.0, size * M.DEFAULT_LINE_FACTOR))
    return float(ls)


def _text_height_pt(paras: list[_Para], width_pt: float, ea: str,
                    wrap: bool = True) -> tuple[float, int, float]:
    """(height needed, line count, widest line) for a block of paragraphs.

    With `wrap=False` — a box whose `word_wrap` is off, like a code panel where
    line N must stay on row N — each paragraph stays one line and the widest
    line is what tells you whether it fits.
    """
    total = 0.0
    lines = 0
    widest = 0.0
    for i, para in enumerate(paras):
        face = para.ea or ea
        if not para.text:
            total += M.DEFAULT_LINE_FACTOR * para.size * para.line_spacing
            lines += 1
            continue
        if wrap:
            wrapped = M.wrap_text(para.text, para.font, para.size, width_pt,
                                  ea_font=face, bold=para.bold)
        else:
            wrapped = [para.text]
        for line in wrapped:
            widest = max(widest, M.measure_pt(line, para.font, para.size,
                                              ea_font=face, bold=para.bold))
        h = M.block_height_pt(wrapped, para.size, line_spacing=para.line_spacing)
        total += h + (para.space_before if i else 0.0) + para.space_after
        lines += len(wrapped)
    return total, lines, widest


# ── colour / contrast ───────────────────────────────────────────────────────
def _rgb(color) -> tuple[int, int, int] | None:
    try:
        if color is None or color.type is None:
            return None
        rgb = color.rgb
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def _lum(c: tuple[int, int, int]) -> float:
    def ch(v):
        v /= 255.0
        return v / 12.92 if v <= 0.04045 else ((v + 0.055) / 1.055) ** 2.4
    r, g, b = (ch(x) for x in c)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def contrast_ratio(fg: tuple[int, int, int], bg: tuple[int, int, int]) -> float:
    a, b = _lum(fg), _lum(bg)
    hi, lo = max(a, b), min(a, b)
    return (hi + 0.05) / (lo + 0.05)


def _slide_bg(slide, prs) -> tuple[int, int, int]:
    for src in (slide.background, ):
        try:
            c = _rgb(src.fill.fore_color)
            if c:
                return c
        except Exception:
            pass
    return (255, 255, 255)


def _fill_rgb(shape) -> tuple[int, int, int] | None:
    try:
        if shape.fill.type is None or shape.fill.type == 5:   # inherit/none
            return None
        return _rgb(shape.fill.fore_color)
    except Exception:
        return None


# ── main entry ──────────────────────────────────────────────────────────────
def audit_pptx(path: str | Path, *, ea_font: str | None = None,
               latin_font: str | None = None,
               safe_margin_in: float = DEFAULT_SAFE_MARGIN_IN,
               ignore_top_in: float = 0.0,
               ignore_bottom_in: float = 0.0) -> list[Finding]:
    """Audit a built deck.  `ignore_*_in` skip theme chrome bands (frametitle
    band, footer bar) that are meant to touch the slide edge."""
    from pptx import Presentation

    prs = Presentation(str(path))
    sw, sh = prs.slide_width, prs.slide_height
    latin = latin_font or "Helvetica Neue"
    ea = ea_font or "Hiragino Sans"
    out: list[Finding] = []
    layouts: list[str] = []
    fonts_seen: set[str] = set()

    for idx, slide in enumerate(prs.slides, start=1):
        bg = _slide_bg(slide, prs)
        # A full-width band at an edge is theme chrome (frametitle bar, footer
        # bar) — it is *supposed* to touch the edge, and so is the text on it.
        top_chrome, bottom_chrome = _chrome_bands(slide, sw, sh)
        top_ig = max(ignore_top_in, top_chrome)
        bot_ig = max(ignore_bottom_in, bottom_chrome)
        text_boxes: list[tuple] = []          # (rect, shape, needed_h_pt)
        filled: list[tuple] = []              # (rect, rgb) — card backgrounds
        pictures = 0
        tables = 0
        charts = 0

        for shape in slide.shapes:
            rect = _rect(shape)
            if rect is None:
                continue
            if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                pictures += 1
            if getattr(shape, "has_table", False) and shape.has_table:
                tables += 1
            if getattr(shape, "has_chart", False) and shape.has_chart:
                charts += 1
            rgb = _fill_rgb(shape)
            if rgb and not (shape.has_text_frame and shape.text_frame.text.strip()):
                filled.append((rect, rgb))

            if (shape.has_text_frame
                    and _PAGE_NO.fullmatch(shape.text_frame.text.strip())):
                continue

            # Header-crumb chrome (research/hearing style): a small text box
            # living entirely above the top margin is builder chrome — author
            # content starts at MARGIN_T (0.5in), so nothing legitimate fits
            # wholly inside the top 0.38in.
            if (shape.has_text_frame and rect is not None
                    and rect[3] <= int(0.38 * 914400)):
                continue

            if getattr(shape, "has_table", False) and shape.has_table:
                out += _check_table(idx, shape, rect, sh, latin, ea)
                text_boxes.append((rect, shape,
                                   _table_needed_pt(shape, latin, ea), False))
                continue

            # ── off-slide / cramped margin ──
            _check_bounds(out, idx, shape, rect, sw, sh, safe_margin_in,
                          top_ig, bot_ig)

            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            tf = shape.text_frame
            paras = _paragraphs(shape, latin, ea)
            for p in paras:
                fonts_seen.add(p.font)
                if p.ea:
                    fonts_seen.add(p.ea)
            inner_w = (shape.width - (tf.margin_left or 0) - (tf.margin_right or 0)) / EMU_PT
            inner_h = (shape.height - (tf.margin_top or 0) - (tf.margin_bottom or 0)) / EMU_PT
            if inner_w <= 1:
                continue
            wraps = tf.word_wrap is not False
            needed, nlines, widest = _text_height_pt(paras, inner_w, ea, wraps)
            text_boxes.append((rect, shape, needed, _is_ghost(shape, bg)))

            # ── overflow ──
            if not wraps and widest > inner_w + 2.0:
                out.append(Finding(
                    "overflow", "error", idx,
                    f"line is {widest - inner_w:.0f}pt wider than its box, and the "
                    f"box does not wrap — it runs out the side",
                    fix="shorten the line, or widen the block",
                    shape=_label(shape),
                    detail={"needed_pt": widest, "box_pt": inner_w}))
            grows = str(getattr(tf, "auto_size", "")).startswith("SHAPE_TO_FIT")
            slack = max(2.0, inner_h * 0.04)     # PowerPoint's own rounding
            if needed > inner_h + slack:
                over = needed - inner_h
                sev = "error" if not grows else "warn"
                how = ("the box is set to grow, so it will push into whatever "
                       "sits below it") if grows else "text will be clipped"
                out.append(Finding(
                    "overflow", sev, idx,
                    f"text needs {needed:.0f}pt in a {inner_h:.0f}pt box "
                    f"(over by {over:.0f}pt, {nlines} lines) — {how}",
                    fix="shorten the text, split the slide, or lower the font size "
                        f"to ~{_suggest_size(paras, inner_w, inner_h, ea):.0f}pt",
                    shape=_label(shape),
                    detail={"needed_pt": needed, "box_pt": inner_h, "lines": nlines}))

            # ── contrast ──
            _check_contrast(out, idx, shape, rect, paras, filled, bg)

        _check_overlaps(out, idx, text_boxes)
        layouts.append(_layout_signature(slide))
        if not (pictures or tables or charts) and len(text_boxes) > 1:
            out.append(Finding("deck", "info", idx,
                               "text-only slide — no figure, table, chart or diagram",
                               fix="add a visual, or merge it with a neighbouring slide",
                               shape=""))

    out.extend(_deck_level(layouts, fonts_seen))
    out.sort(key=lambda f: (_RANK.get(f.severity, 3), f.slide))
    return out


def _table_cell_rows(shape, latin: str, ea: str) -> list[float]:
    """Height (pt) each row of a table actually needs for its cells."""
    table = shape.table
    widths = [c.width for c in table.columns]
    needed = []
    for ri, row in enumerate(table.rows):
        tallest = 0.0
        for ci, cell in enumerate(row.cells):
            if not cell.text.strip():
                continue
            inner = (widths[ci] - (cell.margin_left or 0)
                     - (cell.margin_right or 0)) / EMU_PT
            if inner <= 1:
                continue
            h, _, _ = _text_height_pt(_paragraphs(cell, latin, ea), inner, ea)
            h += ((cell.margin_top or 0) + (cell.margin_bottom or 0)) / EMU_PT
            tallest = max(tallest, h)
        needed.append(max(tallest, row.height / EMU_PT))
    return needed


def _table_needed_pt(shape, latin: str, ea: str) -> float:
    try:
        return sum(_table_cell_rows(shape, latin, ea))
    except Exception:
        return shape.height / EMU_PT


def _check_table(idx: int, shape, rect, sh: int, latin: str, ea: str) -> list[Finding]:
    """Tables grow rather than clip: a cell that needs more room pushes every
    row below it down, and the table off the bottom of the slide. python-pptx
    exposes them as a graphic frame with no text frame, so the ordinary text
    checks never saw one."""
    out: list[Finding] = []
    try:
        rows = _table_cell_rows(shape, latin, ea)
    except Exception:
        return out
    stored = [r.height / EMU_PT for r in shape.table.rows]
    grew = [(i, n, s) for i, (n, s) in enumerate(zip(rows, stored))
            if n > s + max(2.0, s * 0.04)]
    needed = sum(rows)
    frame = shape.height / EMU_PT
    bottom = rect[1] / EMU_PT + needed
    limit = sh / EMU_PT - 0.3 * 72
    if bottom > limit:
        out.append(Finding(
            "overflow", "error", idx,
            f"table needs {needed:.0f}pt from y={rect[1] / EMU_IN:.1f}\" and runs "
            f"{(bottom - limit) / 72:.2f}\" past the bottom of the slide",
            fix="shorten the cells, drop a column, or split the table over two slides",
            shape=_label(shape),
            detail={"needed_pt": needed, "box_pt": frame,
                    "rows": [round(r, 1) for r in rows]}))
    elif grew:
        first = grew[0]
        out.append(Finding(
            "overflow", "warn", idx,
            f"row {first[0] + 1} of the table needs {first[1]:.0f}pt in a "
            f"{first[2]:.0f}pt row — the table will grow {needed - frame:.0f}pt "
            f"and push what is under it",
            fix="shorten the longest cell in that row, or widen its column",
            shape=_label(shape),
            detail={"needed_pt": needed, "box_pt": frame}))
    return out


def _chrome_bands(slide, sw: int, sh: int) -> tuple[float, float]:
    """(top, bottom) bands in inches that belong to the theme, not the author.

    A frametitle band or a footer bar spans the full width and sits flush to
    its edge; flagging its text for "too close to the edge" would be noise on
    every slide of a deck that uses one.
    """
    top_parts: list[tuple[int, int, int]] = []
    bottom_parts: list[tuple[int, int, int]] = []
    for shape in slide.shapes:
        rect = _rect(shape)
        if rect is None or _fill_rgb(shape) is None:
            continue
        if rect[1] <= sh * 0.02:
            top_parts.append((rect[0], rect[2], rect[3]))
        if rect[3] >= sh * 0.98:
            bottom_parts.append((rect[0], rect[2], sh - rect[1]))
    return _band_depth(top_parts, sw), _band_depth(bottom_parts, sw)


def _band_depth(parts, sw: int) -> float:
    """Depth of an edge band, in inches, or 0 if the parts don't span the slide.

    A bar is often drawn as several segments — beamer's footer is three tones —
    so coverage is measured on the union, not on any single shape.
    """
    if not parts:
        return 0.0
    covered = 0
    end = -1
    for x0, x1, _ in sorted(parts):
        if x1 <= end:
            continue
        covered += x1 - max(x0, end)
        end = x1
    if covered < sw * 0.9:
        return 0.0
    return max(p[2] for p in parts) / EMU_IN


def _suggest_size(paras: list[_Para], w: float, h: float, ea: str) -> float:
    """Largest uniform scale of the current sizes that fits."""
    scale = 1.0
    while scale > 0.55:
        total = 0.0
        for i, p in enumerate(paras):
            lines = M.wrap_text(p.text, p.font, p.size * scale, w,
                                ea_font=p.ea or ea, bold=p.bold) if p.text else [""]
            total += M.block_height_pt(lines, p.size * scale,
                                       line_spacing=p.line_spacing)
            total += (p.space_before if i else 0.0) + p.space_after
        if total <= h:
            break
        scale -= 0.05
    return max(p.size for p in paras) * scale


def needed_height_emu(shape, *, latin: str = "Helvetica Neue",
                      ea: str = "Hiragino Sans") -> int:
    """Height (EMU) the shape's wrapped text needs at the shape's own width.

    The one place height is estimated: `builder` reserves boxes with it,
    `visuallint`'s collision detector checks the reservation held, and the
    audit reports what didn't.
    """
    tf = shape.text_frame
    inner_w = (shape.width - (tf.margin_left or 0) - (tf.margin_right or 0)) / EMU_PT
    if inner_w <= 1:
        return 0
    needed, _, _ = _text_height_pt(_paragraphs(shape, latin, ea), inner_w, ea,
                                   tf.word_wrap is not False)
    return int(needed * EMU_PT)


def _check_bounds(out, idx, shape, rect, sw, sh, margin_in, top_ig, bot_ig):
    # deco:* shapes (hero corner motifs) bleed off the edge by design.
    if (getattr(shape, "name", "") or "").startswith("deco:"):
        return
    lo, t, r, b = rect
    if r <= 0 or b <= 0 or lo >= sw or t >= sh:
        out.append(Finding("offslide", "error", idx,
                           "shape is entirely outside the slide",
                           fix="check the coordinates that placed it",
                           shape=_label(shape)))
        return
    if lo < 0 or t < 0 or r > sw or b > sh:
        over = max(-lo, -t, r - sw, b - sh) / EMU_IN
        out.append(Finding("offslide", "error", idx,
                           f"shape extends {over:.2f}\" past the slide edge",
                           fix="shrink the box or move it inside the content area",
                           shape=_label(shape)))
        return
    if not shape.has_text_frame or not shape.text_frame.text.strip():
        return
    # Text that lives inside a theme band — the frametitle bar, the footer bar —
    # is chrome. The band is meant to run edge to edge, and so is its label.
    if t < top_ig * EMU_IN or b > sh - bot_ig * EMU_IN:
        return
    m = margin_in * EMU_IN
    # Fine print — a source line or a page number — belongs in the band below
    # the content area. Only body-sized text down there is a crowding problem.
    fine = _max_pt(shape) <= 12.0
    near = []
    if lo < m:
        near.append("left")
    if r > sw - m:
        near.append("right")
    if t < m:
        near.append("top")
    if b > sh - m and not fine:
        near.append("bottom")
    if near:
        out.append(Finding("offslide", "warn", idx,
                           f"text sits within {margin_in}\" of the "
                           f"{'/'.join(near)} edge",
                           fix="pull it into the safe area (≥0.5\" is comfortable)",
                           shape=_label(shape)))


def _text_colors(shape) -> list[tuple[int, int, int]]:
    """Every colour the shape's text is written in.

    Run-level first, then paragraph-level: the builder sets `p.font.color` for
    single-run blocks, and reading only runs would miss those entirely.
    """
    out = []
    for p in shape.text_frame.paragraphs:
        run_colors = [_rgb(r.font.color) for r in p.runs]
        run_colors = [c for c in run_colors if c]
        if run_colors:
            out += run_colors
        else:
            c = _rgb(p.font.color)
            if c:
                out.append(c)
    return out


def _is_ghost(shape, bg) -> bool:
    """True for a watermark — type tinted so far toward the background that it
    reads as decoration (the giant section numeral behind a divider title).
    Content is meant to sit on top of it, so it is not an overlap."""
    colors = _text_colors(shape)
    if not colors:
        return False
    return max(contrast_ratio(c, bg) for c in colors) < 2.0


def _max_pt(shape) -> float:
    sizes = [r.font.size.pt for p in shape.text_frame.paragraphs for r in p.runs
             if r.font.size is not None]
    sizes += [p.font.size.pt for p in shape.text_frame.paragraphs
              if p.font.size is not None]
    return max(sizes) if sizes else DEFAULT_SIZE_PT


# A text box is a rectangle; the ink inside it is not. Two boxes have to
# interpenetrate by this much before their glyphs can actually touch — the
# figure `visuallint` was tuned to, which keeps a small-caps eyebrow sitting
# directly on its heading from reading as a collision.
_OVERLAP_DEPTH = 0.12 * EMU_IN
_OVERLAP_SPAN = 0.20 * EMU_IN


def _check_overlaps(out, idx, boxes):
    for i in range(len(boxes)):
        ra, sa, needa, ghosta = boxes[i]
        # Use the height the text actually needs — a box set to grow will
        # occupy that much on screen even though its stored height is smaller.
        ra = (ra[0], ra[1], ra[2], max(ra[3], ra[1] + int(needa * EMU_PT)))
        for j in range(i + 1, len(boxes)):
            rb, sb, needb, ghostb = boxes[j]
            if ghosta or ghostb:
                continue        # a watermark is meant to sit under the content
            rb = (rb[0], rb[1], rb[2], max(rb[3], rb[1] + int(needb * EMU_PT)))
            span = min(ra[2], rb[2]) - max(ra[0], rb[0])
            depth = min(ra[3], rb[3]) - max(ra[1], rb[1])
            if span > _OVERLAP_SPAN and depth > _OVERLAP_DEPTH:
                area = _intersection(ra, rb)
                out.append(Finding(
                    "overlap", "error", idx,
                    f"“{_label(sa)}” and “{_label(sb)}” overlap by {area:.2f} in² "
                    f"({depth / EMU_IN:.2f}\" deep)",
                    fix="move one block down, or shorten the text above it",
                    shape=_label(sa)))


def _check_contrast(out, idx, shape, rect, paras, filled, slide_bg):
    # A shape that carries its own fill *is* its text's background — a number
    # badge in an accent circle, a dark callout. Checking its white label
    # against the slide behind it reports 1.1:1 for something perfectly legible.
    own = _fill_rgb(shape)
    bg = own or slide_bg
    best = None
    for frect, frgb in filled if own is None else ():
        if (frect[0] <= rect[0] and frect[1] <= rect[1]
                and frect[2] >= rect[2] and frect[3] >= rect[3]):
            area = (frect[2] - frect[0]) * (frect[3] - frect[1])
            if best is None or area < best[0]:
                best = (area, frgb)
    if best:
        bg = best[1]
    seen = set()
    for para, xml_p in zip(paras, shape.text_frame.paragraphs):
        if not para.text.strip():
            continue
        color = next((_rgb(r.font.color) for r in xml_p.runs
                      if _rgb(r.font.color)), None) or _rgb(xml_p.font.color)
        if color is None or (color, para.size) in seen:
            continue
        seen.add((color, para.size))
        ratio = contrast_ratio(color, bg)
        large = para.size >= 18 or (para.size >= 14 and para.bold)
        need = 3.0 if large else 4.5
        if ratio + 0.05 < need:      # a pair that lands exactly on the bar passes
            out.append(Finding(
                "contrast", "warn", idx,
                f"{para.size:.0f}pt text at {ratio:.1f}:1 against its background "
                f"(WCAG AA needs {need}:1)",
                fix="darken the text colour, or drop it onto a lighter panel",
                shape=_label(shape), detail={"ratio": ratio}))
            break


def _layout_signature(slide) -> str:
    """A coarse fingerprint of where things sit, for monotony detection."""
    cells = set()
    for shape in slide.shapes:
        rect = _rect(shape)
        if rect is None or not shape.has_text_frame:
            continue
        if not shape.text_frame.text.strip():
            continue
        cx = int((rect[0] + rect[2]) / 2 / EMU_IN * 2)
        cy = int((rect[1] + rect[3]) / 2 / EMU_IN * 2)
        cells.add((cx, cy))
    return ",".join(f"{x}:{y}" for x, y in sorted(cells))


def _deck_level(layouts: list[str], fonts: set[str]) -> list[Finding]:
    out: list[Finding] = []
    run_start = 0
    for i in range(1, len(layouts) + 1):
        if i == len(layouts) or layouts[i] != layouts[run_start]:
            span = i - run_start
            if span >= 4 and layouts[run_start]:
                out.append(Finding(
                    "deck", "info", run_start + 1,
                    f"slides {run_start + 1}–{i} all use the same layout "
                    f"({span} in a row)",
                    fix="vary the type: a comparison, a figure, or a KPI row "
                        "breaks the rhythm"))
            run_start = i
    for name in sorted(fonts):
        status = M.font_status(name)
        if status == "missing":
            out.append(Finding(
                "font", "warn", 1,
                f"“{name}” is not installed here — fit checks for it are "
                f"approximate (measured with a metric-compatible stand-in)",
                fix="install the font, or pick one from the safe list "
                    "(Arial, Calibri, Cambria, Times New Roman)"))
        elif status == "unreliable":
            out.append(Finding(
                "font", "info", 1,
                f"“{name}” renders with different widths in LibreOffice preview "
                f"than in PowerPoint",
                fix="leave ~10% slack in boxes using it, or use a safe-list font"))
    return out


def format_findings(findings: list[Finding], *, limit: int | None = None) -> str:
    if not findings:
        return "clean — no layout, contrast or font problems found"
    counts: dict[str, int] = {}
    for f in findings:
        counts[f.severity] = counts.get(f.severity, 0) + 1
    head = "  ".join(f"{k}: {v}" for k, v in
                     sorted(counts.items(), key=lambda kv: _RANK.get(kv[0], 3)))
    body = findings if limit is None else findings[:limit]
    lines = [str(f) for f in body]
    if limit is not None and len(findings) > limit:
        lines.append(f"… and {len(findings) - limit} more")
    return head + "\n" + "\n".join(lines)
